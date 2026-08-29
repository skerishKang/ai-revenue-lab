from __future__ import annotations

from io import BytesIO

from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation
from pptx.util import Inches


DOCX_TEXT = "Padiem DOCX worker compatibility"
PPTX_TEXT = "Padiem PPTX worker compatibility"
XLSX_VALUE = "Padiem XLSX worker compatibility"


def test_python_docx_in_memory_round_trip_extracts_known_text() -> None:
    output = BytesIO()
    document = Document()
    document.add_paragraph(DOCX_TEXT)
    document.save(output)

    output.seek(0)
    parsed = Document(output)
    extracted = "\n".join(paragraph.text for paragraph in parsed.paragraphs)

    assert DOCX_TEXT in extracted


def test_python_pptx_in_memory_round_trip_extracts_known_text() -> None:
    output = BytesIO()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    textbox.text = PPTX_TEXT
    presentation.save(output)

    output.seek(0)
    parsed = Presentation(output)
    extracted = [
        shape.text
        for parsed_slide in parsed.slides
        for shape in parsed_slide.shapes
        if hasattr(shape, "text")
    ]

    assert PPTX_TEXT in extracted


def test_openpyxl_in_memory_round_trip_extracts_known_value() -> None:
    output = BytesIO()
    workbook = Workbook()
    sheet = workbook.active
    sheet["A1"] = XLSX_VALUE
    workbook.save(output)

    output.seek(0)
    parsed = load_workbook(output, read_only=True, data_only=True)
    try:
        assert parsed.active["A1"].value == XLSX_VALUE
    finally:
        parsed.close()
