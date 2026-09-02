#!/usr/bin/env python3
"""Generate the Business 35 customer-facing one-page offer source PPTX.

Commercial truth is consumed from the exact accepted Lane A revision via
``accepted_source`` (01-one-page-offer.md). Product identity is the V3.1
six-stage primary journey; no seven-step education identity.

Single slide, 16:9, editable text, speaker notes. Exported to PDF separately.
Layout: cover band 0..1.35; value 1.5; two boxes 2.15..5.15; next steps 5.35;
footer 6.55. No overlapping layers.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from pathlib import Path
import datetime
import sys
sys.path.insert(0, str(Path(__file__).resolve().parent))
from accepted_source import require_accepted_source, user_journey_six  # noqa: E402
FIXED_DT = datetime.datetime(2026, 9, 3, 0, 0, 0)
OUT = Path(__file__).resolve().parent.parent / "Business35_OnePage_Offer_Source.pptx"

NAVY = RGBColor(0x1F, 0x3A, 0x5F)
BLUE = RGBColor(0x2E, 0x5E, 0x8C)
GRAY = RGBColor(0x55, 0x5A, 0x60)
LIGHT = RGBColor(0xF2, 0xF4, 0xF7)
ACCENT = RGBColor(0xC2, 0x7B, 0x2D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SW, SH = Inches(13.333), Inches(7.5)


def para(tf, text, size=14, bold=False, color=GRAY, first=False, align=None, space_before=None):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    if align:
        p.alignment = align
    if space_before is not None:
        p.space_before = Pt(space_before)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Malgun Gothic"
    return p


def build():
    snapshot = require_accepted_source()
    journey = user_journey_six(snapshot)
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)

    bg = slide.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT
    bg.line.fill.background()
    bg.shadow.inherit = False

    # Cover band
    band = slide.shapes.add_shape(1, 0, 0, SW, Inches(1.35))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    band.shadow.inherit = False
    tf = band.text_frame
    tf.margin_left = Inches(0.55)
    tf.margin_top = Inches(0.22)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Business 35 · AI Media Education & DX"
    r.font.size = Pt(28)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "Malgun Gothic"
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = "AI 업무전환 프로그램 — 1페이지 소개"
    r2.font.size = Pt(16)
    r2.font.color.rgb = RGBColor(0xC9, 0xD6, 0xE3)
    r2.font.name = "Malgun Gothic"
    p3 = tf.add_paragraph()
    r3 = p3.add_run()
    r3.text = "제공: 파디엠"
    r3.font.size = Pt(13)
    r3.font.bold = True
    r3.font.color.rgb = RGBColor(0xC9, 0xD6, 0xE3)
    r3.font.name = "Malgun Gothic"

    # Value proposition (accepted 01 한 줄 가치 제안)
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(1.5), Inches(12.2), Inches(0.6))
    vtf = tb.text_frame
    vtf.word_wrap = True
    para(vtf, "조직의 실제 미디어 업무 한 흐름을 입력하면, AI 적용 후보·사람 검토 지점·추천 파일럿·운영 산출물을 한 번에 구성한다.",
         size=14, bold=True, color=NAVY, first=True)

    # Left column: V3.1 six-stage user journey (box bottom at 5.15)
    left = slide.shapes.add_shape(1, Inches(0.55), Inches(2.15), Inches(6.0), Inches(3.0))
    left.fill.solid()
    left.fill.fore_color.rgb = WHITE
    left.line.color.rgb = BLUE
    left.line.width = Pt(1)
    left.shadow.inherit = False
    ltf = left.text_frame
    ltf.word_wrap = True
    ltf.margin_left = Inches(0.2)
    ltf.margin_right = Inches(0.2)
    ltf.margin_top = Inches(0.12)
    para(ltf, "사용자가 실제로 하는 일 (V3.1)", size=12, bold=True, color=BLUE, first=True)
    for step in journey:
        short = step.split(". ", 1)[-1] if ". " in step else step
        para(ltf, "• " + short, size=10, color=NAVY)

    # Right column: products A / B / C (box bottom at 5.15)
    right = slide.shapes.add_shape(1, Inches(6.75), Inches(2.15), Inches(6.0), Inches(3.0))
    right.fill.solid()
    right.fill.fore_color.rgb = WHITE
    right.line.color.rgb = BLUE
    right.line.width = Pt(1)
    right.shadow.inherit = False
    rtf = right.text_frame
    rtf.word_wrap = True
    rtf.margin_left = Inches(0.2)
    rtf.margin_right = Inches(0.2)
    rtf.margin_top = Inches(0.12)
    para(rtf, "상품 A · 진단 워크숍", size=13, bold=True, color=BLUE, first=True)
    para(rtf, "• 1~2일 · 초기형 300만–500만원 · 확장형 500만–800만원", size=11, color=GRAY)
    para(rtf, "• 현재 흐름 진단 · 후보 선정 · 위험 분리", size=11, color=GRAY)
    para(rtf, "상품 B1 · 디자인 파트너 파일럿", size=13, bold=True, color=BLUE, space_before=6)
    para(rtf, "• 6주 · 1팀 · 1업무 · 1,000만–1,500만원", size=11, color=GRAY)
    para(rtf, "상품 B2 · 표준 6주 파일럿", size=13, bold=True, color=BLUE, space_before=4)
    para(rtf, "• 6주 · 1,500만–2,500만원", size=11, color=GRAY)
    para(rtf, "• 기준선·교육·재설계·파일럿·성과·플레이북", size=11, color=GRAY)
    para(rtf, "상품 C · 운영 자문 · 월 300만–600만원", size=13, bold=True, color=BLUE, space_before=6)

    # Next steps (accepted 01 다음 행동)
    tb = slide.shapes.add_textbox(Inches(0.55), Inches(5.35), Inches(12.2), Inches(0.9))
    ntf = tb.text_frame
    ntf.word_wrap = True
    para(ntf, "다음 행동: 진단 워크숍 또는 6주 파일럿 범위 확인", size=14, bold=True, color=BLUE, first=True)
    para(ntf, "바꾸고 싶은 미디어 업무 1건 선정 → 현재 흐름·병목·사람 검토 지점 확인 → 고객별 범위와 가격 가설을 다시 승인한 뒤 견적 제안",
         size=12, color=GRAY, space_before=4)

    # Footer (customer review copy only, no internal English markers)
    foot = slide.shapes.add_textbox(Inches(0.55), Inches(6.45), Inches(12.2), Inches(0.85))
    ftf = foot.text_frame
    ftf.word_wrap = True
    para(ftf, "파디엠 · DRAFT", size=9, color=GRAY, first=True)
    para(ftf, "가격은 시장 검증 전 가설이며 범위·인원·기간에 따라 달라질 수 있습니다.",
         size=9, color=GRAY)

    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "핵심 말할 내용: 1페이지 안에 문제·방식·상품·가격·다음 단계를 설명한다.\n\n"
        "고객에게 물어볼 질문: 어떤 업무가 가장 시간이 많이 걸리나요?\n\n"
        "과장해서는 안 되는 부분: 가격을 확정 가격처럼 말하지 않고, 성과·정부지원금을 보장하지 않는다."
    )

    Path(OUT).parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.core_properties.created = FIXED_DT
        prs.core_properties.modified = FIXED_DT
        prs.core_properties.revision = 1
    except Exception:
        pass
    prs.save(str(OUT))
    from normalize_ooxml import normalize_ooxml
    normalize_ooxml(OUT)
    print(f"saved {OUT}")


if __name__ == "__main__":
    build()
