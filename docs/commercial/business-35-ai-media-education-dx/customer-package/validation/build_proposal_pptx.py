#!/usr/bin/env python3
"""Generate the Business 35 customer-facing master proposal PPTX (V3.1 deterministic).

Commercial truth is consumed from the exact accepted Lane A revision via
``accepted_source`` (SOURCE_REVISION=63adbefcf24a91a5a064c6b8e13779e151ba7de7),
aligned to accepted 02-ten-page-proposal.md (10 pages, 16:9).

Product identity is the V3.1 six-stage primary journey. Week 0-6 / step
sequences appear only as downstream delivery detail, never as the product
identity.

Layout rule: title band 0..1.15; headline 1.3..1.92; body starts at 2.1;
footer at 7.05. No overlapping text layers. All text/shapes editable; each slide
carries a page number and speaker notes. No external images/fonts/stock assets.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from pathlib import Path
import datetime
import sys
sys.path.insert(0, str(Path(__file__).resolve().parent))
from accepted_source import require_accepted_source, six_stage_journey  # noqa: E402
FIXED_DT = datetime.datetime(2026, 9, 3, 0, 0, 0)
OUT = Path(__file__).resolve().parent.parent / "Business35_Master_Proposal_10p.pptx"

NAVY = RGBColor(0x1F, 0x3A, 0x5F)
BLUE = RGBColor(0x2E, 0x5E, 0x8C)
GRAY = RGBColor(0x55, 0x5A, 0x60)
LIGHT = RGBColor(0xF2, 0xF4, 0xF7)
ACCENT = RGBColor(0xC2, 0x7B, 0x2D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SOFT = RGBColor(0xD5, 0xDE, 0xE8)

SW, SH = Inches(13.333), Inches(7.5)

FOOTER_COVER = "파디엠 · DRAFT"
FOOTER_INNER = "파디엠 · DRAFT"
FOOTER_LAST = "파디엠 · DRAFT"
PROVIDER = "파디엠"

TOTAL_SLIDES = 10

HEADLINE_Y = Inches(1.3)
HEADLINE_H = Inches(0.62)
BODY_Y = Inches(2.1)


def new_deck():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    return prs


def add_slide(prs, title, footer_mode="inner", slide_no=None, headline=None):
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    bg = slide.shapes.add_shape(1, 0, 0, SW, SH)
    bg.fill.solid()
    bg.fill.fore_color.rgb = LIGHT
    bg.line.fill.background()
    bg.shadow.inherit = False
    band = slide.shapes.add_shape(1, 0, 0, SW, Inches(1.15))
    band.fill.solid()
    band.fill.fore_color.rgb = NAVY
    band.line.fill.background()
    band.shadow.inherit = False
    tf = band.text_frame
    tf.margin_left = Inches(0.55)
    tf.margin_right = Inches(0.4)
    tf.margin_top = Inches(0.18)
    tf.margin_bottom = Inches(0.12)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "Malgun Gothic"
    if slide_no is not None:
        # Badge sits in the headline zone's right side, NOT overlapping the title band
        badge = slide.shapes.add_shape(5, Inches(12.05), Inches(1.38), Inches(0.95), Inches(0.6))
        badge.fill.solid()
        badge.fill.fore_color.rgb = ACCENT
        badge.line.fill.background()
        badge.shadow.inherit = False
        btf = badge.text_frame
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run()
        br.text = f"{slide_no:02d}"
        br.font.size = Pt(18)
        br.font.bold = True
        br.font.color.rgb = WHITE
        br.font.name = "Malgun Gothic"
    # Single headline layer (one message, no duplicate body text near it)
    if headline:
        hb = slide.shapes.add_shape(1, Inches(0.7), HEADLINE_Y, Inches(11.15), HEADLINE_H)
        hb.fill.solid()
        hb.fill.fore_color.rgb = RGBColor(0xE3, 0xE9, 0xF0)
        hb.line.color.rgb = BLUE
        hb.line.width = Pt(0.75)
        hb.shadow.inherit = False
        htf = hb.text_frame
        htf.word_wrap = True
        htf.vertical_anchor = MSO_ANCHOR.MIDDLE
        htf.margin_left = Inches(0.2)
        hp = htf.paragraphs[0]
        hr = hp.add_run()
        hr.text = headline
        hr.font.size = Pt(15)
        hr.font.bold = True
        hr.font.color.rgb = NAVY
        hr.font.name = "Malgun Gothic"
    # Footer
    footer = slide.shapes.add_textbox(Inches(0.55), Inches(7.05), Inches(11.7), Inches(0.35))
    ftf = footer.text_frame
    ftf.word_wrap = True
    fp = ftf.paragraphs[0]
    fr = fp.add_run()
    if footer_mode == "cover":
        fr.text = FOOTER_COVER
    elif footer_mode == "last":
        fr.text = FOOTER_LAST
    else:
        fr.text = FOOTER_INNER
    fr.font.size = Pt(11)
    fr.font.color.rgb = GRAY
    fr.font.name = "Malgun Gothic"
    return slide


def add_body_box(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)
    return tf


def add_page_number(prs, slide, n):
    pn = slide.shapes.add_textbox(Inches(12.35), Inches(7.05), Inches(0.8), Inches(0.35))
    tf = pn.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = f"{n} / {TOTAL_SLIDES}"
    r.font.size = Pt(10)
    r.font.color.rgb = GRAY
    r.font.name = "Malgun Gothic"


def para(tf, text, size=16, bold=False, color=GRAY, first=False, space_before=None, space_after=8, align=None):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    if space_before is not None:
        p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if align:
        p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Malgun Gothic"
    return p


def add_step_box(slide, x, y, w, h, num, label):
    box = slide.shapes.add_shape(5, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = BLUE
    box.line.color.rgb = NAVY
    box.line.width = Pt(1)
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.04)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = num
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(13)
    r2.font.color.rgb = WHITE
    r2.font.name = "Malgun Gothic"


def add_notes(slide, core, questions, caution):
    notes = slide.notes_slide.notes_text_frame
    notes.text = (
        "핵심 말할 내용: " + core + "\n\n"
        "고객에게 물어볼 질문: " + questions + "\n\n"
        "과장해서는 안 되는 부분: " + caution
    )


def add_card(slide, x, y, w, h, title, body_lines, title_color=NAVY, body_size=14):
    card = slide.shapes.add_shape(1, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = BLUE
    card.line.width = Pt(1)
    card.shadow.inherit = False
    ctf = card.text_frame
    ctf.word_wrap = True
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ctf.margin_left = Inches(0.2)
    ctf.margin_right = Inches(0.2)
    ctf.margin_top = Inches(0.15)
    cp = ctf.paragraphs[0]
    cp.alignment = PP_ALIGN.CENTER
    cr = cp.add_run()
    cr.text = title
    cr.font.size = Pt(18)
    cr.font.bold = True
    cr.font.color.rgb = title_color
    cr.font.name = "Malgun Gothic"
    for line in body_lines:
        cp2 = ctf.add_paragraph()
        cp2.space_before = Pt(6)
        cp2.alignment = PP_ALIGN.CENTER
        cr2 = cp2.add_run()
        cr2.text = line
        cr2.font.size = Pt(body_size)
        cr2.font.color.rgb = GRAY
        cr2.font.name = "Malgun Gothic"
    return card


def build():
    snapshot = require_accepted_source()
    journey = six_stage_journey(snapshot)
    journey_line = " → ".join(journey)
    prs = new_deck()

    # ---- Slide 1 · 제품과 결과 (accepted 02 Page 1) ----
    s = add_slide(prs, "1. 제품과 결과 — 파디엠 AI 미디어 업무전환 스튜디오",
                  footer_mode="cover", slide_no=1,
                  headline="팀의 실제 미디어 업무 한 흐름을 사람이 승인하는 운영체계로 바꾼다")
    tf = add_body_box(s, Inches(0.7), BODY_Y, Inches(11.9), Inches(4.2))
    para(tf, "파디엠 AI 미디어 업무전환 스튜디오 — 서비스 주도형 업무전환 스튜디오",
         size=18, bold=True, color=NAVY, first=True)
    para(tf, "입력 한 번으로 진단·새 업무 흐름·사람 검토 지점·추천 파일럿·운영 산출물을 구성한다",
         size=15, color=GRAY, space_before=10)
    para(tf, "대상: 지역 문화기관 · 교육기관 · 협회·단체 · 미디어·콘텐츠 기관 · 기업 홍보·콘텐츠팀",
         size=14, color=GRAY, space_before=10)
    para(tf, "제품 흐름: 고객 입력 → 진단 → 새 workflow → 추천 pilot",
         size=14, bold=True, color=BLUE, space_before=10)
    para(tf, "V3.1 여정: " + journey_line,
         size=11, color=GRAY, space_before=10)
    para(tf, "실제 계약·매출 발생 주장이 아닙니다. 제공: 파디엠",
         size=12, color=GRAY, space_before=8)
    add_notes(
        s,
        "첫 30초 안에 누구를 위한 제품인지 / 무엇을 입력하는지 / 무엇이 나오는지를 설명한다. 일반 AI 강의가 아니라 고객의 현재 업무가 출발점이다.",
        "조직에서 가장 시간이 오래 걸리는 콘텐츠 업무는 무엇인가요?",
        "가격을 확정 가격처럼 말하지 않는다. 실제 계약·매출 발생 주장이 아님을 명시한다."
    )
    add_page_number(prs, s, 1)

    # ---- Slide 2 · 지금 바꿀 업무를 고른다 (accepted 02 Page 2) ----
    s = add_slide(prs, "2. 지금 바꿀 업무를 고른다", slide_no=2,
                  headline="어떤 미디어 업무의 어디가 막혀 있는지를 먼저 정한다")
    cards = [
        ("입력", "다섯 가지 입력", ["조직·결과물·병목·팀 규모·", "AI 사용 상태를 입력합니다."]),
        ("후보", "결과물·병목 후보", ["홍보물·교육자료·영상·이미지·", "캠페인 중 병목: 기획·초안·제작·검토·승인·배포"]),
        ("원칙", "합성·비식별 입력", ["비공개·민감 원문 없이도", "제품 구조와 파일럿 후보를 설명합니다."]),
    ]
    for i, (badge, title, lines) in enumerate(cards):
        cx = Inches(0.7) + Inches(4.1) * i
        card = s.shapes.add_shape(5, cx, BODY_Y, Inches(3.85), Inches(3.2))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = BLUE
        card.line.width = Pt(1)
        card.shadow.inherit = False
        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.margin_left = Inches(0.2)
        ctf.margin_right = Inches(0.2)
        ctf.margin_top = Inches(0.25)
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = badge
        cr.font.size = Pt(26)
        cr.font.bold = True
        cr.font.color.rgb = ACCENT
        cr.font.name = "Malgun Gothic"
        cp2 = ctf.add_paragraph()
        cp2.alignment = PP_ALIGN.CENTER
        cr2 = cp2.add_run()
        cr2.text = title
        cr2.font.size = Pt(19)
        cr2.font.bold = True
        cr2.font.color.rgb = NAVY
        cr2.font.name = "Malgun Gothic"
        for line in lines:
            cp3 = ctf.add_paragraph()
            cp3.space_before = Pt(8)
            cp3.alignment = PP_ALIGN.CENTER
            cr3 = cp3.add_run()
            cr3.text = line
            cr3.font.size = Pt(14)
            cr3.font.color.rgb = GRAY
            cr3.font.name = "Malgun Gothic"
    add_notes(
        s,
        "AI 도구 목록보다 먼저 어떤 미디어 업무의 어디가 막혀 있는지를 정한다고 설명한다. 팀 규모와 AI 사용 상태를 함께 확인한다.",
        "지금 어떤 업무를 바꾸고 싶습니까? 병목은 어디인가요?",
        "비공개·민감 원문을 초기 상담 단계에서 받지 않는다."
    )
    add_page_number(prs, s, 2)

    # ---- Slide 3 · 조직별 진단이 나온다 (accepted 02 Page 3) ----
    s = add_slide(prs, "3. 조직별 진단이 나온다", slide_no=3,
                  headline="조직마다 적용 후보·제외 영역·사람 검토 지점이 달라야 한다")
    card1 = add_card(s, Inches(0.7), BODY_Y, Inches(5.9), Inches(2.3),
                     "적용 후보 진단",
                     ["현재 흐름과 병목을 기준으로", "적용 후보를 좁힙니다."])
    card2 = add_card(s, Inches(6.9), BODY_Y, Inches(5.9), Inches(2.3),
                     "사람 검토 지점",
                     ["자동화하면 안 되는 지점을 별도 표시하고,", "승인 도구와 금지·주의 자료의 경계를 기록합니다."])
    tf = add_body_box(s, Inches(0.7), Inches(4.8), Inches(11.9), Inches(1.6))
    para(tf, "AI 활용 확산과 함께 개인정보·저작권·투명성·안전성·사람 검토 등 조직 차원의 사용정책과 거버넌스 요구가 강화되고 있다.",
         size=14, color=GRAY, first=True)
    para(tf, "진단은 법률 판단을 대신하지 않는다. 개인정보·저작권·조달은 고객별 확인이 필요하다.",
         size=13, color=GRAY, space_before=8)
    add_notes(
        s,
        "같은 AI 도구라도 조직마다 적용 후보·제외 영역·사람 검토 지점이 달라야 한다고 설명한다.",
        "검토와 승인 없이 자동화하면 안 되는 지점이 어디인가요?",
        "진단을 법률 판단처럼 말하지 않는다."
    )
    add_page_number(prs, s, 3)

    # ---- Slide 4 · 새 업무 흐름을 설계한다 (accepted 02 Page 4) ----
    # V3.1: Before/After workflow with human approval gates. This slide must
    # NOT present a seven-step education sequence as the product identity.
    s = add_slide(prs, "4. 새 업무 흐름을 설계한다", slide_no=4,
                  headline="AI는 workflow 안에 들어가고, 사람 승인 gate는 사라지지 않는다")
    flow_now = s.shapes.add_shape(1, Inches(0.7), BODY_Y, Inches(5.9), Inches(2.0))
    flow_now.fill.solid()
    flow_now.fill.fore_color.rgb = WHITE
    flow_now.line.color.rgb = BLUE
    flow_now.line.width = Pt(1)
    flow_now.shadow.inherit = False
    ntf = flow_now.text_frame
    ntf.word_wrap = True
    ntf.margin_left = Inches(0.2)
    ntf.margin_top = Inches(0.15)
    np_ = ntf.paragraphs[0]
    nr = np_.add_run()
    nr.text = "현재 흐름"
    nr.font.size = Pt(17)
    nr.font.bold = True
    nr.font.color.rgb = NAVY
    nr.font.name = "Malgun Gothic"
    for line in ["기획 → 초안 → 제작", "→ 검토 → 승인 → 게시"]:
        np2 = ntf.add_paragraph()
        np2.space_before = Pt(8)
        nr2 = np2.add_run()
        nr2.text = line
        nr2.font.size = Pt(15)
        nr2.font.color.rgb = GRAY
        nr2.font.name = "Malgun Gothic"
    flow_new = s.shapes.add_shape(1, Inches(6.9), BODY_Y, Inches(5.9), Inches(2.0))
    flow_new.fill.solid()
    flow_new.fill.fore_color.rgb = WHITE
    flow_new.line.color.rgb = BLUE
    flow_new.line.width = Pt(1.25)
    flow_new.shadow.inherit = False
    wtf = flow_new.text_frame
    wtf.word_wrap = True
    wtf.margin_left = Inches(0.2)
    wtf.margin_top = Inches(0.15)
    wp = wtf.paragraphs[0]
    wr = wp.add_run()
    wr.text = "전환 흐름 예시"
    wr.font.size = Pt(17)
    wr.font.bold = True
    wr.font.color.rgb = NAVY
    wr.font.name = "Malgun Gothic"
    for line in ["요청 정의 → AI 보조 초안 → 제작", "→ 사람 검토 → 수정 → 승인 → 게시"]:
        wp2 = wtf.add_paragraph()
        wp2.space_before = Pt(8)
        wr2 = wp2.add_run()
        wr2.text = line
        wr2.font.size = Pt(15)
        wr2.font.bold = True
        wr2.font.color.rgb = BLUE
        wr2.font.name = "Malgun Gothic"
    tf = add_body_box(s, Inches(0.7), Inches(4.5), Inches(11.9), Inches(1.9))
    para(tf, "자동화 대상과 사람 담당자의 책임을 구분하고, 예외 처리와 중단 조건을 설계한다. 결과물 품질보다 속도만 높이는 구조를 제안하지 않는다.",
         size=15, color=GRAY, first=True)
    add_notes(
        s,
        "사람이 빠지는 자동화가 아니라 사람이 더 명확하게 승인하는 workflow임을 강조한다.",
        "현재 가장 바꾸고 싶은 업무가 무엇인가요?",
        "전체 업무 전환을 보장하지 않는다."
    )
    add_page_number(prs, s, 4)

    # ---- Slide 5 · 운영 산출물을 확인한다 (accepted 02 Page 5) ----
    s = add_slide(prs, "5. 운영 산출물을 확인한다", slide_no=5,
                  headline="발표 자료가 아니라 팀이 실제로 쓰는 운영 산출물로 남는다")
    items = [
        "업무 요청서",
        "AI 사용정책 초안",
        "금지·주의 자료 기준",
        "사람 검토 지도",
        "프롬프트·작업 템플릿",
        "KPI 기준선·측정표",
        "파일럿 운영 체크리스트",
        "전환 요약",
    ]
    col1, col2 = Inches(0.7), Inches(6.5)
    row_y = BODY_Y
    for i, item in enumerate(items):
        x = col1 if i < 4 else col2
        yy = row_y + Inches(0.7) * (i % 4)
        box = s.shapes.add_shape(5, x, yy, Inches(5.3), Inches(0.58))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = BLUE
        box.line.width = Pt(1)
        box.shadow.inherit = False
        btf = box.text_frame
        btf.word_wrap = True
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]
        br = bp.add_run()
        br.text = "•  " + item
        br.font.size = Pt(15)
        br.font.color.rgb = NAVY
        br.font.name = "Malgun Gothic"
    tf = add_body_box(s, Inches(0.7), Inches(5.1), Inches(11.9), Inches(1.2))
    para(tf, "각 산출물은 고객 조직의 범위에 맞춰 선택한다. 법률·계약 문서는 전문 검토가 필요한 영역을 별도로 표시한다.",
         size=13, color=GRAY, first=True)
    para(tf, "교육자료 몇 장이 아니라 조직 운영에 남는 결과를 보여준다. 제공: 파디엠",
         size=13, color=GRAY, space_before=6)
    add_notes(
        s,
        "전환안은 발표 자료로 끝나지 않고 팀이 실제로 쓰는 운영 산출물로 남는다고 설명한다.",
        "이 중 조직에 가장 필요한 산출물은 무엇인가요?",
        "법률·계약 문서를 확정 문서처럼 표현하지 않는다."
    )
    add_page_number(prs, s, 5)

    # ---- Slide 6 · 상품 A (accepted 02 Page 6) ----
    s = add_slide(prs, "6. 상품 A · 진단 워크숍", slide_no=6,
                  headline="바꿀 업무와 파일럿 범위를 먼저 확정한다")
    dur = s.shapes.add_shape(5, Inches(0.7), BODY_Y, Inches(2.9), Inches(1.3))
    dur.fill.solid()
    dur.fill.fore_color.rgb = BLUE
    dur.line.fill.background()
    dur.shadow.inherit = False
    dtf = dur.text_frame
    dtf.vertical_anchor = MSO_ANCHOR.MIDDLE
    dp = dtf.paragraphs[0]
    dp.alignment = PP_ALIGN.CENTER
    dr = dp.add_run()
    dr.text = "1~2일 워크숍"
    dr.font.size = Pt(18)
    dr.font.bold = True
    dr.font.color.rgb = WHITE
    dr.font.name = "Malgun Gothic"
    prc = s.shapes.add_shape(5, Inches(0.7), Inches(3.6), Inches(2.9), Inches(1.6))
    prc.fill.solid()
    prc.fill.fore_color.rgb = ACCENT
    prc.line.fill.background()
    prc.shadow.inherit = False
    ptf2 = prc.text_frame
    ptf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    ptf2.word_wrap = True
    pp = ptf2.paragraphs[0]
    pp.alignment = PP_ALIGN.CENTER
    prr = pp.add_run()
    prr.text = "초기형 300만–500만원"
    prr.font.size = Pt(15)
    prr.font.bold = True
    prr.font.color.rgb = WHITE
    prr.font.name = "Malgun Gothic"
    pp2 = ptf2.add_paragraph()
    pp2.alignment = PP_ALIGN.CENTER
    prr2 = pp2.add_run()
    prr2.text = "확장형 500만–800만원"
    prr2.font.size = Pt(13)
    prr2.font.color.rgb = RGBColor(0xFF, 0xE8, 0xC8)
    prr2.font.name = "Malgun Gothic"
    out = s.shapes.add_shape(1, Inches(3.95), BODY_Y, Inches(8.65), Inches(2.8))
    out.fill.solid()
    out.fill.fore_color.rgb = WHITE
    out.line.color.rgb = BLUE
    out.line.width = Pt(1)
    out.shadow.inherit = False
    otf = out.text_frame
    otf.word_wrap = True
    otf.margin_left = Inches(0.25)
    otf.margin_top = Inches(0.2)
    op = otf.paragraphs[0]
    orr = op.add_run()
    orr.text = "주요 산출물"
    orr.font.size = Pt(17)
    orr.font.bold = True
    orr.font.color.rgb = NAVY
    orr.font.name = "Malgun Gothic"
    outputs = [
        "진단 요약",
        "추천 workflow 초안",
        "파일럿 후보 및 범위 제안",
        "위험·확인 항목 목록",
        "사람 검토 gate 초안",
    ]
    for o in outputs:
        op2 = otf.add_paragraph()
        orr2 = op2.add_run()
        orr2.text = "•  " + o
        orr2.font.size = Pt(15)
        orr2.font.color.rgb = GRAY
        orr2.font.name = "Malgun Gothic"
    para_tf = add_body_box(s, Inches(0.7), Inches(5.2), Inches(11.9), Inches(1.4))
    para(para_tf, "가격은 시장 검증 전 자사 가격 가설입니다. 범위와 인원·기간에 따라 최종 견적이 달라집니다.",
         size=14, color=GRAY, first=True)
    add_notes(
        s,
        "진단 워크숍은 변화의 출발점이며, 고객 조직이 무엇을 바꿀지 함께 찾는 자리라고 설명한다.",
        "워크숍에 참여할 수 있는 팀과 일정이 있나요?",
        "가격을 확정 가격처럼 말하지 않는다 — 가설임을 명시한다."
    )
    add_page_number(prs, s, 6)

    # ---- Slide 7 · 상품 B1/B2 · 6주 파일럿 (accepted 02 Page 7) ----
    # Week 0-6 is downstream delivery detail, not the product identity.
    s = add_slide(prs, "7. 상품 B1/B2 · 6주 파일럿", slide_no=7,
                  headline="실제 workflow와 사람 검토 체계를 6주간 시험한다")
    sc = s.shapes.add_shape(5, Inches(0.7), BODY_Y, Inches(2.9), Inches(1.3))
    sc.fill.solid()
    sc.fill.fore_color.rgb = BLUE
    sc.line.fill.background()
    sc.shadow.inherit = False
    stf = sc.text_frame
    stf.vertical_anchor = MSO_ANCHOR.MIDDLE
    sp = stf.paragraphs[0]
    sp.alignment = PP_ALIGN.CENTER
    srr = sp.add_run()
    srr.text = "6주 · 1팀 · 1핵심 업무"
    srr.font.size = Pt(17)
    srr.font.bold = True
    srr.font.color.rgb = WHITE
    srr.font.name = "Malgun Gothic"
    prc6 = s.shapes.add_shape(5, Inches(0.7), Inches(3.6), Inches(2.9), Inches(1.5))
    prc6.fill.solid()
    prc6.fill.fore_color.rgb = ACCENT
    prc6.line.fill.background()
    prc6.shadow.inherit = False
    ptf6 = prc6.text_frame
    ptf6.vertical_anchor = MSO_ANCHOR.MIDDLE
    ptf6.word_wrap = True
    pp6 = ptf6.paragraphs[0]
    pp6.alignment = PP_ALIGN.CENTER
    prr6 = pp6.add_run()
    prr6.text = "1,000만–1,500만원"
    prr6.font.size = Pt(17)
    prr6.font.bold = True
    prr6.font.color.rgb = WHITE
    prr6.font.name = "Malgun Gothic"
    pp6b = ptf6.add_paragraph()
    pp6b.alignment = PP_ALIGN.CENTER
    prr6b = pp6b.add_run()
    prr6b.text = "상품 B1 · 디자인 파트너"
    prr6b.font.size = Pt(11)
    prr6b.font.color.rgb = RGBColor(0xFF, 0xE8, 0xC8)
    prr6b.font.name = "Malgun Gothic"
    pp6c = ptf6.add_paragraph()
    pp6c.alignment = PP_ALIGN.CENTER
    prr6c = pp6c.add_run()
    prr6c.text = "B2 표준 1,500만–2,500만원"
    prr6c.font.size = Pt(11)
    prr6c.font.color.rgb = RGBColor(0xFF, 0xE8, 0xC8)
    prr6c.font.name = "Malgun Gothic"
    out6 = s.shapes.add_shape(1, Inches(3.95), BODY_Y, Inches(8.65), Inches(2.8))
    out6.fill.solid()
    out6.fill.fore_color.rgb = WHITE
    out6.line.color.rgb = BLUE
    out6.line.width = Pt(1)
    out6.shadow.inherit = False
    otf6 = out6.text_frame
    otf6.word_wrap = True
    otf6.margin_left = Inches(0.25)
    otf6.margin_top = Inches(0.2)
    op6 = otf6.paragraphs[0]
    orr6b = op6.add_run()
    orr6b.text = "파일럿 진행 흐름"
    orr6b.font.size = Pt(17)
    orr6b.font.bold = True
    orr6b.font.color.rgb = NAVY
    orr6b.font.name = "Malgun Gothic"
    flow6 = [
        "범위·책임·중단 조건 확정",
        "현재 흐름·기준선 진단",
        "승인 도구·금지 규칙 교육",
        "새 workflow·운영 템플릿 설계",
        "제한 파일럿 운영",
        "KPI·위험·산출물 검토",
    ]
    for o in flow6:
        op6b = otf6.add_paragraph()
        orr6c = op6b.add_run()
        orr6c.text = "•  " + o
        orr6c.font.size = Pt(15)
        orr6c.font.color.rgb = GRAY
        orr6c.font.name = "Malgun Gothic"
    # Week 0-6 timeline row (accepted 04 delivery detail)
    rows = [
        ("W0", "범위·책임", "범위 확정"),
        ("W1", "기준선 진단", "기준선 측정"),
        ("W2", "직무별 교육", "규칙 확인"),
        ("W3", "실제 실습", "실습 결과"),
        ("W4", "재설계", "검토 gate"),
        ("W5", "제한 파일럿", "중간 보고"),
        ("W6", "측정·플레이북", "플레이북 승인"),
    ]
    tl_y = Inches(5.1)
    tl_h = Inches(0.95)
    tl_w = Inches(1.62)
    tl_gap = Inches(0.12)
    for i, (label, desc, result) in enumerate(rows):
        tx = Inches(0.7) + (tl_w + tl_gap) * i
        box = s.shapes.add_shape(5, tx, tl_y, tl_w, tl_h)
        box.fill.solid()
        box.fill.fore_color.rgb = BLUE
        box.line.fill.background()
        box.shadow.inherit = False
        btf = box.text_frame
        btf.word_wrap = True
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        btf.margin_left = Inches(0.04)
        btf.margin_right = Inches(0.04)
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run()
        br.text = label + " " + desc
        br.font.size = Pt(11)
        br.font.bold = True
        br.font.color.rgb = WHITE
        br.font.name = "Malgun Gothic"
        bp2 = btf.add_paragraph()
        bp2.alignment = PP_ALIGN.CENTER
        br2 = bp2.add_run()
        br2.text = "▸ " + result
        br2.font.size = Pt(10)
        br2.font.color.rgb = RGBColor(0xE8, 0xEE, 0xF5)
        br2.font.name = "Malgun Gothic"
        # connecting timeline line between weeks
        if i < len(rows) - 1:
            line_x = tx + tl_w + Inches(0.01)
            line = s.shapes.add_shape(1, line_x, tl_y + Inches(0.42), Inches(0.1), Inches(0.1))
            line.fill.solid()
            line.fill.fore_color.rgb = ACCENT
            line.line.fill.background()
            line.shadow.inherit = False
    tf6 = add_body_box(s, Inches(0.7), Inches(6.25), Inches(11.9), Inches(0.7))
    para(tf6, "이 주차 구조는 제품을 정의하는 7단계 정체성이 아니라, 선택된 파일럿을 수행하는 delivery detail이다. 가격은 시장 검증 전 가설이다.",
         size=12, color=GRAY, first=True)
    add_notes(
        s,
        "파일럿은 1팀·1핵심 업무로 제한되어 위험을 낮추고, 측정 가능한 결과를 만든다고 설명한다.",
        "파일럿 후보 업무가 있나요? 담당자를 지정할 수 있나요?",
        "성과를 보장하지 않는다. 사람 검토 gate가 필수임을 명시한다."
    )
    add_page_number(prs, s, 7)

    # ---- Slide 8 · KPI와 위험을 함께 본다 (accepted 02 Page 8) ----
    s = add_slide(prs, "8. KPI와 위험을 함께 본다", slide_no=8,
                  headline="품질·사람 검토·정책 준수·팀 사용성을 함께 측정한다")
    tf = add_body_box(s, Inches(0.7), BODY_Y, Inches(11.9), Inches(0.8))
    para(tf, "KPI 후보", size=20, bold=True, color=NAVY, first=True)
    kpis = ["기준 생산시간 vs 파일럿 생산시간", "재작업률", "사람 검토 통과율",
            "정책 위반 또는 중단 건수", "팀 참여·실제 사용률", "결과물 품질 평가",
            "미해결 위험·운영 산출물 승인"]
    for i, k in enumerate(kpis[:6]):
        col = Inches(0.7) if i < 3 else Inches(6.5)
        yy = Inches(2.9) + Inches(0.7) * (i % 3)
        box = s.shapes.add_shape(5, col, yy, Inches(5.3), Inches(0.55))
        box.fill.solid()
        box.fill.fore_color.rgb = WHITE
        box.line.color.rgb = BLUE
        box.line.width = Pt(0.75)
        box.shadow.inherit = False
        btf = box.text_frame
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]
        br = bp.add_run()
        br.text = k
        br.font.size = Pt(13)
        br.font.color.rgb = NAVY
        br.font.name = "Malgun Gothic"
    tf2 = add_body_box(s, Inches(0.7), Inches(5.0), Inches(11.9), Inches(1.5))
    para(tf2, "KPI 후보: 미해결 위험과 운영 산출물 승인 여부 포함. 위험 경계: 민감정보 무제한 외부 입력 금지 · 사람 검토 없는 자동 게시 금지 · 고위험 판단 자동화 제외 · 사고·위반 시 중단 절차 사전 합의",
         size=14, color=GRAY, first=True)
    para(tf2, "성과를 보장하지 않는다. 기준선과 KPI로 변화 여부를 측정한다.", size=14, bold=True, color=ACCENT, space_before=10)
    add_notes(
        s,
        "KPI는 속도·품질·채택·거버넌스 준수를 분리해 측정한다고 설명한다. 성과 보장이 아님을 강조한다.",
        "현재 어떤 기준으로 성과를 판단하고 있나요?",
        "KPI 목표 가설을 성과 보장으로 오인하지 않게 한다."
    )
    add_page_number(prs, s, 8)

    # ---- Slide 9 · 상품 C와 가격 가설 (accepted 02 Page 9) ----
    s = add_slide(prs, "9. 상품 C와 가격 가설", slide_no=9,
                  headline="범위 확인 후 최종 견적 — 모두 시장 검증 전 가설이다")
    price_cards = [
        ("상품 A", "진단 워크숍", "초기형 300만–500만원", "확장형 500만–800만원"),
        ("상품 B1", "디자인 파트너", "1,000만–1,500만원", "6주 파일럿"),
        ("상품 B2", "표준 파일럿", "1,500만–2,500만원", "6주 파일럿"),
        ("상품 C", "운영 자문", "월 300만–600만원", "월 단위"),
    ]
    for i, (prod, name, price, note) in enumerate(price_cards):
        cx = Inches(0.7) + Inches(3.1) * i
        card = s.shapes.add_shape(5, cx, BODY_Y, Inches(2.85), Inches(2.6))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = BLUE
        card.line.width = Pt(1.25)
        card.shadow.inherit = False
        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ctf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = prod
        cr.font.size = Pt(13)
        cr.font.bold = True
        cr.font.color.rgb = BLUE
        cr.font.name = "Malgun Gothic"
        cp2 = ctf.add_paragraph()
        cp2.alignment = PP_ALIGN.CENTER
        cr2 = cp2.add_run()
        cr2.text = name
        cr2.font.size = Pt(17)
        cr2.font.bold = True
        cr2.font.color.rgb = NAVY
        cr2.font.name = "Malgun Gothic"
        cp3 = ctf.add_paragraph()
        cp3.alignment = PP_ALIGN.CENTER
        cr3 = cp3.add_run()
        cr3.text = price
        cr3.font.size = Pt(17)
        cr3.font.bold = True
        cr3.font.color.rgb = ACCENT
        cr3.font.name = "Malgun Gothic"
        cp4 = ctf.add_paragraph()
        cp4.alignment = PP_ALIGN.CENTER
        cr4 = cp4.add_run()
        cr4.text = note
        cr4.font.size = Pt(11)
        cr4.font.color.rgb = GRAY
        cr4.font.name = "Malgun Gothic"
    tf = add_body_box(s, Inches(0.7), Inches(4.9), Inches(11.9), Inches(1.6))
    para(tf, "시장 검증 전 자사 가격 가설", size=15, bold=True, color=ACCENT, first=True)
    para(tf, "범위·인원·기간에 따라 최종 견적 · VAT 조건은 최종 견적서에서 확정",
         size=15, bold=True, color=ACCENT, space_before=6)
    add_notes(
        s,
        "가격은 가설이며 범위·인원·기간에 따라 달라진다고 명확히 한다. VAT는 최종 견적서에서 확정된다고 말한다.",
        "어느 상품 범위가 관심 있나요?",
        "경쟁사 평균 가격 주장, 정부지원금 수령 가능 주장을 하지 않는다."
    )
    add_page_number(prs, s, 9)

    # ---- Slide 10 · 다음 행동과 계약 경계 (accepted 02 Page 10) ----
    s = add_slide(prs, "10. 다음 행동과 계약 경계", footer_mode="last", slide_no=10,
                  headline="진단 워크숍 또는 6주 파일럿의 범위를 먼저 확인한다")
    steps10 = [
        ("1", "바꿀 업무 1건 선정", "현재 흐름·병목 확인"),
        ("2", "검토 지점 확인", "사람 검토 지점 포함"),
        ("3", "A·B1/B2 적합 판단", "진단 또는 파일럿"),
        ("4", "범위·중단 조건 합의", "일정·참여시간 포함"),
    ]
    for i, (num, label, desc) in enumerate(steps10):
        box = s.shapes.add_shape(5, Inches(0.7) + Inches(3.1) * i, BODY_Y, Inches(2.7), Inches(2.4))
        box.fill.solid()
        box.fill.fore_color.rgb = BLUE
        box.line.color.rgb = NAVY
        box.line.width = Pt(1)
        box.shadow.inherit = False
        btf = box.text_frame
        btf.word_wrap = True
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run()
        br.text = num
        br.font.size = Pt(28)
        br.font.bold = True
        br.font.color.rgb = WHITE
        p2 = btf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = label
        r2.font.size = Pt(15)
        r2.font.bold = True
        r2.font.color.rgb = WHITE
        r2.font.name = "Malgun Gothic"
        p3 = btf.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        r3 = p3.add_run()
        r3.text = desc
        r3.font.size = Pt(11)
        r3.font.color.rgb = SOFT
        r3.font.name = "Malgun Gothic"
    tf = add_body_box(s, Inches(0.7), Inches(4.9), Inches(11.9), Inches(1.8))
    para(tf, "다음 행동: 범위·일정·고객 참여시간·중단 조건 합의 → 고객별 가격 가설 재승인 → 필요한 법률·계약 검토 후 최종 제안.",
         size=14, color=GRAY, first=True)
    para(tf, "개인정보·저작권·조달은 고객별 확인 필요. SOW와 위험·데이터 부속서는 전문 법률·계약 검토 필요. 공공기관 계약방식은 기관 계약담당자 확인 필요.",
         size=13, color=GRAY, space_before=8)
    para(tf, "제공 및 계약 주체: 파디엠", size=14, bold=True, color=NAVY, space_before=10)
    para(tf, "가격은 시장 검증 전 가설이며 범위·인원·기간에 따라 달라질 수 있다. 현재 문서는 내부 상업 초안이며 실제 계약·매출 발생 주장이 아니다.",
         size=12, color=GRAY, space_before=6)
    add_notes(
        s,
        "제품 데모의 마지막 행동과 제안서의 마지막 행동을 동일하게 유지한다: 진단 또는 파일럿 범위 확인.",
        "진단 워크숍 또는 6주 파일럿 중 어디부터 범위를 확인할까요?",
        "정부지원금 수령 가능, 1억원 이하 자동 수의계약을 말하지 않는다."
    )
    add_page_number(prs, s, 10)

    Path(OUT).parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.core_properties.created = FIXED_DT
        prs.core_properties.modified = FIXED_DT
        prs.core_properties.revision = 1
    except Exception:
        pass
    prs.save(str(OUT))
    from normalize_ooxml import normalize_ooxml
    normalize_ooxml(OUT)
    print(f"saved {OUT} with {len(prs.slides._sldIdLst)} slides")


if __name__ == "__main__":
    build()
