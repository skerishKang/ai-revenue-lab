#!/usr/bin/env python3
"""Generate Business 35 V3.1 customer-package review artifacts.

Outputs are intentionally written under customer-package/v3-regenerated/ so the
pre-V3.1 binaries remain historical evidence until Web CTO visual QA accepts the
new package.
"""

from __future__ import annotations

from pathlib import Path
from typing import Iterable

from docx import Document
from docx.enum.section import WD_SECTION
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_BREAK
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches as DInches, Pt as DPt, RGBColor as DRGB
from openpyxl import Workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "v3-regenerated"
OUT.mkdir(parents=True, exist_ok=True)

FONT = "Noto Sans CJK KR"
IVORY = "F5F0E6"
PAPER = "FFFDFC"
FOREST = "21372E"
FOREST_2 = "315044"
COBALT = "2757C8"
COBALT_SOFT = "E8EEFC"
MOSS = "91A493"
INK_MUTED = "657168"
LINE = "D8D8CF"
WARM = "C17D58"
WHITE = "FFFFFF"
BLACK = "111713"

SW, SH = Inches(13.333), Inches(7.5)


def rgb(hexstr: str) -> RGBColor:
    return RGBColor.from_string(hexstr)


def add_rect(slide, x, y, w, h, fill, line=None, radius=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid(); shape.fill.fore_color.rgb = rgb(fill)
    if line:
        shape.line.color.rgb = rgb(line); shape.line.width = Pt(0.8)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, text, x, y, w, h, size=18, color=FOREST, bold=False,
             align=PP_ALIGN.LEFT, font=FONT, margin=0.02, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear(); tf.word_wrap = True; tf.vertical_anchor = valign
    tf.margin_left = tf.margin_right = Inches(margin)
    tf.margin_top = tf.margin_bottom = Inches(margin)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = font; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = rgb(color)
    return box


def add_rich_lines(slide, lines: Iterable[tuple[str, int, str, bool]], x, y, w, h, gap=6):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    for i, (text, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        r = p.add_run(); r.text = text
        r.font.name = FONT; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = rgb(color)
    return box


def footer(slide, page: int, label="파디엠 · DRAFT"):
    add_text(slide, label, .55, 7.05, 3.5, .22, 8, INK_MUTED, True)
    add_text(slide, f"{page:02d} / 10", 11.9, 7.05, .85, .22, 8, INK_MUTED, True, PP_ALIGN.RIGHT)


def slide_base(prs, page, section, dark=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 13.333, 7.5, FOREST if dark else IVORY)
    add_text(slide, "PADIEM · BUSINESS 35", .55, .3, 2.8, .3, 9, COBALT if not dark else "C9D8D0", True)
    add_text(slide, section, 9.1, .3, 3.65, .3, 9, INK_MUTED if not dark else "C9D8D0", True, PP_ALIGN.RIGHT)
    footer(slide, page, "파디엠 · DRAFT" if not dark else "파디엠 · DRAFT · 내부 검토용")
    return slide


def card(slide, title, body, x, y, w, h, accent=COBALT, number=None, fill=PAPER):
    add_rect(slide, x, y, w, h, fill, LINE, radius=True)
    if number:
        add_text(slide, number, x+.22, y+.2, .6, .28, 9, accent, True)
        title_y = y+.58
    else:
        title_y = y+.28
    add_text(slide, title, x+.22, title_y, w-.44, .45, 15, FOREST, True)
    add_text(slide, body, x+.22, title_y+.55, w-.44, h-(title_y-y)-.72, 10.5, INK_MUTED, False)


def add_flow(slide, items, x, y, w, h):
    gap = .12
    box_w = (w - gap*(len(items)-1))/len(items)
    for i, (title, sub) in enumerate(items):
        bx = x + i*(box_w+gap)
        fill = COBALT if i in (1, 2) else PAPER
        title_color = WHITE if fill == COBALT else FOREST
        sub_color = "DCE5FF" if fill == COBALT else INK_MUTED
        add_rect(slide, bx, y, box_w, h, fill, COBALT if fill != COBALT else COBALT, radius=True)
        add_text(slide, f"{i+1:02d}", bx+.18, y+.16, .48, .22, 8, sub_color, True)
        add_text(slide, title, bx+.18, y+.48, box_w-.36, .44, 13, title_color, True)
        add_text(slide, sub, bx+.18, y+.98, box_w-.36, h-1.12, 9, sub_color)


def build_proposal():
    prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH

    # 1
    s = slide_base(prs, 1, "제품과 결과")
    add_text(s, "파디엠 AI 미디어\n업무전환 스튜디오", .65, 1.05, 7.1, 1.35, 34, FOREST, True)
    add_text(s, "AI 교육을 듣는 데서 끝내지 않고, 팀의 실제 미디어 업무 한 흐름을 사람이 승인하는 운영체계로 바꿉니다.", .68, 2.55, 6.7, 1.0, 15, INK_MUTED)
    add_rect(s, 8.05, 1.05, 4.55, 4.95, PAPER, LINE, radius=True)
    add_text(s, "TRANSFORMATION BRIEF", 8.4, 1.45, 3.8, .25, 9, COBALT, True)
    add_text(s, "조직 입력", 8.4, 1.95, 1.4, .3, 12, FOREST, True)
    add_text(s, "결과물 · 병목 · 팀 규모 · AI 사용 상태", 8.4, 2.3, 3.6, .45, 10, INK_MUTED)
    add_rect(s, 8.4, 2.95, 3.78, .95, COBALT_SOFT, None, radius=True)
    add_text(s, "진단 + 새 업무 흐름", 8.65, 3.18, 3.25, .3, 15, COBALT, True)
    add_rect(s, 8.4, 4.05, 3.78, .95, FOREST, None, radius=True)
    add_text(s, "사람 승인 gate + 추천 파일럿", 8.65, 4.28, 3.2, .42, 13, WHITE, True)
    add_text(s, "상담 준비용 전환 요약", 8.4, 5.25, 3.4, .35, 11, FOREST, True)
    add_text(s, "대상 · 문화/교육/미디어기관 · 기업 홍보/콘텐츠팀", .68, 6.28, 8.3, .35, 10, COBALT, True)

    # 2
    s = slide_base(prs, 2, "현재 업무 병목")
    add_text(s, "도구보다 먼저,\n바꿀 업무를 고릅니다.", .65, 1.02, 6.2, 1.28, 31, FOREST, True)
    add_text(s, "초기 상담은 비공개 원문 없이도 가능합니다. 현재 결과물과 흐름을 기준으로 병목을 좁힙니다.", .68, 2.46, 6.2, .8, 14, INK_MUTED)
    add_flow(s, [
        ("결과물", "홍보물 · 교육자료 · 영상 · 캠페인"),
        ("병목", "기획 · 초안 · 제작 · 검토 · 승인"),
        ("현재 AI", "개별 사용 · 승인 도구 · 금지 자료"),
        ("팀", "담당 인력 · 승인 책임 · 참여 가능 시간"),
    ], .68, 3.55, 11.95, 1.65)
    add_rect(s, .68, 5.55, 11.95, .9, PAPER, LINE, radius=True)
    add_text(s, "질문", .92, 5.82, .7, .22, 9, COBALT, True)
    add_text(s, "지금 한 가지 업무를 바꾼다면 무엇입니까?", 1.72, 5.72, 5.4, .38, 16, FOREST, True)
    add_text(s, "→ 이 답이 진단과 파일럿의 출발점입니다.", 7.35, 5.78, 4.65, .3, 11, INK_MUTED, False, PP_ALIGN.RIGHT)

    # 3
    s = slide_base(prs, 3, "조직별 진단")
    add_text(s, "같은 AI라도, 조직마다\n적용 경계가 달라야 합니다.", .65, 1.02, 7.2, 1.25, 30, FOREST, True)
    card(s, "적용 후보", "반복 초안 · 변형안 · 분류 · 검색 보조처럼 사람이 검토할 수 있는 업무부터 좁힙니다.", .7, 2.65, 3.7, 2.55, COBALT, "01")
    card(s, "사람 검토", "외부 공개·브랜드·정책·민감 판단은 자동 게시하지 않고 승인 지점을 명시합니다.", 4.62, 2.65, 3.7, 2.55, FOREST_2, "02", COBALT_SOFT)
    card(s, "주의 경계", "개인정보·저작권·계약·조달은 고객별 확인이 필요하며 전문 검토가 필요한 영역을 분리합니다.", 8.54, 2.65, 3.7, 2.55, WARM, "03")
    add_text(s, "제품이 하는 일: 법률 판단을 대신하는 것이 아니라, 적용 후보와 검토 책임을 조직이 논의할 수 있는 구조로 만듭니다.", .72, 5.65, 11.5, .7, 12, INK_MUTED)

    # 4
    s = slide_base(prs, 4, "새 업무 흐름")
    add_text(s, "AI는 workflow 안에 들어가고,\n사람 승인 gate는 남습니다.", .65, 1.02, 7.1, 1.25, 30, FOREST, True)
    add_text(s, "BEFORE", .72, 2.62, 1.1, .22, 9, WARM, True)
    add_flow(s, [("기획","요청 불명확"),("초안","개인별 방식"),("제작","반복 수작업"),("검토","기준 분산"),("승인","책임 불명확")], .72, 3.0, 11.85, 1.18)
    add_text(s, "AFTER", .72, 4.48, 1.1, .22, 9, COBALT, True)
    add_flow(s, [("요청 정의","목적·금지범위"),("AI 보조","허용 범위만"),("제작","팀 템플릿"),("사람 검토","품질·정책"),("승인·게시","책임 명시")], .72, 4.86, 11.85, 1.18)
    add_text(s, "중단 조건과 예외 처리도 workflow의 일부로 문서화합니다.", .72, 6.33, 8.2, .3, 11, INK_MUTED)

    # 5
    s = slide_base(prs, 5, "운영 산출물")
    add_text(s, "교육자료가 아니라,\n팀이 계속 쓰는 운영물로 남깁니다.", .65, 1.02, 7.2, 1.25, 30, FOREST, True)
    items = [
        ("업무 요청서","목적·입력·금지범위"),("AI 사용정책","승인 도구·금지 자료"),
        ("사람 검토 지도","검토·승인 책임"),("작업 템플릿","prompt·checklist"),
        ("KPI 기준선","시간·재작업·품질"),("전환 요약","다음 결정과 책임"),
    ]
    for i,(a,b) in enumerate(items):
        x=.72+(i%3)*4.0; y=2.72+(i//3)*1.55
        card(s,a,b,x,y,3.72,1.25,COBALT if i in (0,5) else FOREST_2,f"{i+1:02d}")
    add_text(s, "법률·계약 문서는 전문 검토가 필요한 영역을 별도로 표시합니다.", .72, 6.12, 11.3, .3, 10.5, INK_MUTED)

    # 6
    s = slide_base(prs, 6, "상품 A · 진단 워크숍")
    add_text(s, "작게 시작하려면,\n바꿀 업무와 파일럿 범위를 먼저 확정합니다.", .65, 1.02, 7.35, 1.3, 29, FOREST, True)
    add_rect(s, .72, 2.75, 4.15, 3.15, COBALT, None, radius=True)
    add_text(s, "A", 1.05, 3.05, .6, .8, 34, WHITE, True)
    add_text(s, "진단 워크숍", 1.05, 3.9, 3.0, .42, 19, WHITE, True)
    add_text(s, "초기형 300만–500만원\n확장형 500만–800만원", 1.05, 4.5, 3.2, .75, 14, "E7EDFF", True)
    add_text(s, "가격은 시장 검증 전 가설", 1.05, 5.45, 3.1, .25, 9, "CFD8F8", False)
    add_rich_lines(s, [
        ("포함 범위", 11, COBALT, True),
        ("현재 workflow 인터뷰", 14, FOREST, True),
        ("적용 후보 / 주의 업무 분류", 14, FOREST, True),
        ("사람 검토 gate 초안", 14, FOREST, True),
        ("파일럿 후보 1건 + 다음 단계", 14, FOREST, True),
    ], 5.35, 2.72, 6.6, 3.25, 9)

    # 7
    s = slide_base(prs, 7, "상품 B1/B2 · 6주 파일럿")
    add_text(s, "1팀 · 1핵심업무로\n실제 운영 방식을 6주간 시험합니다.", .65, .95, 7.45, 1.3, 29, FOREST, True)
    add_rect(s, 8.65, 1.0, 3.95, 1.25, PAPER, LINE, radius=True)
    add_text(s, "B1 디자인 파트너", 8.92, 1.27, 3.35, .3, 13, FOREST, True)
    add_text(s, "1,000만–1,500만원", 8.92, 1.65, 3.35, .3, 12, COBALT, True)
    weeks=[("W0","범위·책임"),("W1","기준선"),("W2","적용·검토"),("W3","workflow"),("W4","제한 실습"),("W5","파일럿"),("W6","KPI·결정")]
    for i,(wk,txt) in enumerate(weeks):
        x=.72+i*1.73
        add_rect(s,x,3.05,1.52,1.75,COBALT if i in (2,3,5) else PAPER,COBALT if i in (2,3,5) else LINE,radius=True)
        col=WHITE if i in (2,3,5) else FOREST
        sub="DDE5FF" if i in (2,3,5) else INK_MUTED
        add_text(s,wk,x+.16,3.25,1.15,.28,10,sub,True)
        add_text(s,txt,x+.16,3.7,1.18,.55,13,col,True)
    add_rect(s,.72,5.25,11.85,.85,COBALT_SOFT,None,radius=True)
    add_text(s,"B2 표준 6주 파일럿 · 1,500만–2,500만원",1.0,5.52,6.1,.3,14,COBALT,True)
    add_text(s,"주차 구조는 delivery detail이며 제품 정체성 자체가 아닙니다.",7.3,5.52,4.8,.3,10,INK_MUTED,False,PP_ALIGN.RIGHT)

    # 8
    s = slide_base(prs, 8, "KPI와 위험")
    add_text(s, "속도만 보지 않고,\n품질·검토·정책·사용성을 함께 봅니다.", .65, 1.02, 7.35, 1.3, 29, FOREST, True)
    kpis=[("생산시간","기준선 대비"),("재작업률","반복 수정"),("검토 통과율","사람 승인"),("정책 중단","위험 관리"),("팀 사용성","실제 사용"),("품질 평가","고객 기준")]
    for i,(a,b) in enumerate(kpis):
        x=.72+(i%3)*2.55; y=2.65+(i//3)*1.4
        add_rect(s,x,y,2.3,1.15,PAPER,LINE,radius=True)
        add_text(s,a,x+.18,y+.2,1.95,.3,13,FOREST,True)
        add_text(s,b,x+.18,y+.65,1.95,.2,9,INK_MUTED)
    add_rect(s,8.65,2.65,3.92,2.55,FOREST,None,radius=True)
    add_text(s,"중단 조건",8.95,2.95,3.2,.3,14,WHITE,True)
    add_text(s,"• 민감정보 무제한 외부 입력\n• 사람 검토 없는 자동 게시\n• 고위험 판단 자동화\n• 정책 위반 또는 사고",8.95,3.45,3.1,1.3,11,"DDE8E2",False)
    add_text(s,"성과를 보장하지 않습니다. 합의한 기준선으로 변화 여부를 측정합니다.",.72,5.65,11.6,.4,11,INK_MUTED)

    # 9
    s = slide_base(prs, 9, "가격 가설과 운영 자문")
    add_text(s, "도입 단계에 따라\n진단 · 파일럿 · 운영 자문으로 확장합니다.", .65, 1.02, 7.25, 1.3, 29, FOREST, True)
    prices=[
        ("A","진단 워크숍","300만–800만원","초기/확장"),
        ("B1","디자인 파트너","1,000만–1,500만원","6주"),
        ("B2","표준 파일럿","1,500만–2,500만원","6주"),
        ("C","운영 자문","월 300만–600만원","월 단위"),
    ]
    for i,(code,name,price,term) in enumerate(prices):
        x=.72+i*3.0
        fill=COBALT if code in ("B1","B2") else PAPER
        col=WHITE if fill==COBALT else FOREST
        sub="DDE5FF" if fill==COBALT else INK_MUTED
        add_rect(s,x,2.78,2.75,2.45,fill,COBALT if fill==COBALT else LINE,radius=True)
        add_text(s,code,x+.2,3.0,.6,.35,18,sub,True)
        add_text(s,name,x+.2,3.48,2.3,.38,14,col,True)
        add_text(s,price,x+.2,4.05,2.3,.55,13,col,True)
        add_text(s,term,x+.2,4.72,2.3,.24,9,sub)
    add_text(s,"모든 가격은 시장 검증 전 가설 · 실제 견적은 범위 확인 후 별도 승인 · 실제 계약/매출 주장 아님",.72,5.68,11.65,.55,10.5,INK_MUTED)

    # 10
    s = slide_base(prs, 10, "다음 행동", dark=True)
    add_text(s,"다음 결정은 하나입니다.",.68,1.0,6.8,.5,13,"C9D8D0",True)
    add_text(s,"진단 워크숍 또는\n6주 파일럿의 범위를 확인합니다.",.68,1.65,8.4,1.55,31,WHITE,True)
    steps=["바꿀 업무 1건","현재 흐름·병목","A / B1 / B2 적합성","범위·일정·중단조건","가격 가설 재승인","법률·계약 확인"]
    for i,t in enumerate(steps):
        x=.72+(i%3)*4.0; y=3.75+(i//3)*1.15
        add_rect(s,x,y,3.72,.88,FOREST_2,"536F62",radius=True)
        add_text(s,f"{i+1:02d}",x+.18,y+.2,.45,.2,8,"AFC5BA",True)
        add_text(s,t,x+.68,y+.19,2.75,.32,12,WHITE,True)
    add_text(s,"개인정보·저작권·조달은 고객별 확인 · 전문 법률/계약 검토가 필요한 문서는 별도 gate",.72,6.25,11.7,.36,10,"C9D8D0")

    out=OUT/"Business35_V3_1_Master_Proposal_10p.pptx"; prs.save(out); return out


def build_one_page():
    prs=Presentation(); prs.slide_width=SW; prs.slide_height=SH
    s=prs.slides.add_slide(prs.slide_layouts[6]); add_rect(s,0,0,13.333,7.5,IVORY)
    add_text(s,"PADIEM · BUSINESS 35 · V3.1",.55,.32,3.0,.25,9,COBALT,True)
    add_text(s,"파디엠 AI 미디어\n업무전환 스튜디오",.6,.9,6.4,1.25,30,FOREST,True)
    add_text(s,"조직의 실제 미디어 업무 한 흐름을 입력하면 진단·새 workflow·사람 검토 지점·추천 파일럿·운영 산출물을 한 번에 구성합니다.",.62,2.25,6.1,.85,13.5,INK_MUTED)
    add_flow(s,[("입력","조직·결과물·병목"),("진단","적용 후보·주의"),("workflow","사람 승인 gate"),("pilot","A 또는 B1/B2")],.62,3.35,6.2,1.45)
    add_rect(s,7.25,.85,5.45,5.65,PAPER,LINE,radius=True)
    add_text(s,"도입 옵션",7.6,1.2,2.4,.3,15,COBALT,True)
    options=[("A","진단 워크숍","300만–800만원"),("B1","디자인 파트너 6주","1,000만–1,500만원"),("B2","표준 6주 파일럿","1,500만–2,500만원"),("C","운영 자문","월 300만–600만원")]
    for i,(c,n,p) in enumerate(options):
        y=1.75+i*1.0
        add_text(s,c,7.6,y,.52,.25,10,COBALT,True)
        add_text(s,n,8.25,y,2.4,.25,12,FOREST,True)
        add_text(s,p,10.55,y,1.65,.25,10,INK_MUTED,True,PP_ALIGN.RIGHT)
        if i<3: add_rect(s,7.6,y+.48,4.55,.01,LINE)
    add_rect(s,7.6,5.75,4.55,.45,COBALT_SOFT,None,radius=True)
    add_text(s,"다음: 진단 또는 6주 파일럿 범위 확인",7.82,5.87,4.1,.2,10,COBALT,True)
    add_text(s,"대상 · 문화/교육/협회/미디어기관 · 기업 홍보·콘텐츠팀",.62,5.35,6.1,.34,10,FOREST,True)
    add_text(s,"가격은 시장 검증 전 가설 · 개인정보/저작권/조달은 고객별 확인 · 전문 법률/계약 검토 필요",.62,6.1,11.7,.45,9,INK_MUTED)
    add_text(s,"파디엠 · DRAFT · 제공/계약 주체 세부정보 확인 필요",.62,6.82,7.0,.25,8,INK_MUTED,True)
    out=OUT/"Business35_V3_1_OnePage_Offer_Source.pptx"; prs.save(out); return out


def set_cell_shading(cell, fill):
    tcPr = cell._tc.get_or_add_tcPr(); shd=OxmlElement('w:shd'); shd.set(qn('w:fill'),fill); tcPr.append(shd)


def build_questionnaire():
    doc=Document()
    for sec in doc.sections:
        sec.top_margin=DInches(.55); sec.bottom_margin=DInches(.55); sec.left_margin=DInches(.65); sec.right_margin=DInches(.65)
    normal=doc.styles['Normal']; normal.font.name=FONT; normal.font.size=DPt(10); normal.font.color.rgb=DRGB.from_string(FOREST)
    title=doc.add_paragraph(); title.alignment=WD_ALIGN_PARAGRAPH.LEFT
    r=title.add_run("파디엠 AI 미디어 업무전환 스튜디오"); r.font.name=FONT; r.font.size=DPt(20); r.font.bold=True; r.font.color.rgb=DRGB.from_string(FOREST)
    p=doc.add_paragraph(); rr=p.add_run("고객 진단 질문지 · V3.1 DRAFT"); rr.font.name=FONT; rr.font.size=DPt(11); rr.font.bold=True; rr.font.color.rgb=DRGB.from_string(COBALT)
    p=doc.add_paragraph("실제 개인정보·비공개 원문은 적지 말고 유형과 포함 여부만 표시해 주세요. 이 문서는 견적·일정 확정 전 진단용입니다.")
    p.style=doc.styles['Normal']

    groups=[
        ("1. 조직과 결과물",[("조직/팀 이름과 참여 인원은?","서술"),("주요 결과물은?","홍보물 / 교육자료 / 영상·이미지 / 캠페인 / 기타"),("월평균 제작량은?","짧게")]),
        ("2. 현재 workflow와 병목",[("요청→기획→초안→제작→검토→승인→게시 흐름을 적어 주세요.","서술"),("가장 오래 걸리거나 재작업이 많은 단계는?","서술"),("최종 승인 책임자는 누구입니까?","직책만")]),
        ("3. 현재 AI 사용 상태",[("현재 사용하는 AI 도구와 용도는?","서술"),("승인된 도구/금지된 도구 기준이 있습니까?","예 / 아니오 / 부분"),("사람 검토 없이 자동화하면 안 되는 업무는?","서술")]),
        ("4. 데이터·권리 경계",[("개인정보가 포함됩니까?","예 / 아니오 / 일부"),("타인 저작물 또는 외부 라이선스 자료가 포함됩니까?","예 / 아니오 / 일부"),("외부 AI에 입력하면 안 되는 자료 유형은?","유형만")]),
        ("5. 파일럿 후보",[("6주 파일럿 후보 업무 1개는?","1건"),("기준선으로 측정 가능한 시간/재작업/품질 데이터는?","서술"),("성공 기준은?","서술"),("중단 조건은?","서술")]),
        ("6. 실행 조건",[("운영 책임자 1인을 지정할 수 있습니까?","예 / 아니오"),("참여자가 주당 실습/검토 시간을 확보할 수 있습니까?","예 / 아니오"),("A 진단 / B1 디자인 파트너 / B2 표준 중 현재 선호는?","선택")]),
    ]
    qnum=1
    for gi,(gtitle,qs) in enumerate(groups):
        h=doc.add_paragraph(); hr=h.add_run(gtitle); hr.font.name=FONT; hr.font.size=DPt(13); hr.font.bold=True; hr.font.color.rgb=DRGB.from_string(COBALT)
        for question,guide in qs:
            tbl=doc.add_table(rows=1,cols=1); tbl.autofit=False; cell=tbl.cell(0,0); cell.vertical_alignment=WD_CELL_VERTICAL_ALIGNMENT.TOP
            set_cell_shading(cell, PAPER)
            p=cell.paragraphs[0]; r=p.add_run(f"{qnum:02d}. {question}"); r.font.name=FONT; r.font.size=DPt(10.5); r.font.bold=True; r.font.color.rgb=DRGB.from_string(FOREST)
            p2=cell.add_paragraph(f"응답 가이드 · {guide}"); p2.runs[0].font.name=FONT; p2.runs[0].font.size=DPt(8.5); p2.runs[0].font.color.rgb=DRGB.from_string(INK_MUTED)
            for _ in range(2 if guide=='서술' else 1):
                line=cell.add_paragraph("____________________________________________________________")
                line.runs[0].font.name=FONT; line.runs[0].font.size=DPt(9); line.runs[0].font.color.rgb=DRGB.from_string(LINE)
            doc.add_paragraph().paragraph_format.space_after=DPt(1)
            qnum+=1
        if gi in (1,3):
            doc.add_paragraph().add_run().add_break(WD_BREAK.PAGE)
    note=doc.add_paragraph(); nr=note.add_run("DRAFT · 개인정보·저작권·조달은 고객별 확인이 필요하며 전문 법률·계약 검토가 필요한 영역을 별도로 구분합니다. 가격은 시장 검증 전 가설입니다."); nr.font.name=FONT; nr.font.size=DPt(8.5); nr.font.color.rgb=DRGB.from_string(INK_MUTED)
    out=OUT/"Business35_V3_1_Diagnostic_Questionnaire.docx"; doc.save(out); return out


def style_ws(ws, title):
    ws.sheet_view.showGridLines=False
    ws.freeze_panes="A5"
    ws.merge_cells("A1:F1"); ws["A1"]=title
    ws["A1"].font=Font(name=FONT,size=18,bold=True,color=FOREST); ws["A1"].fill=PatternFill('solid',fgColor=IVORY); ws["A1"].alignment=Alignment(vertical='center')
    ws.row_dimensions[1].height=32
    for c,w in zip(range(1,7),[18,30,24,18,18,36]): ws.column_dimensions[get_column_letter(c)].width=w
    thin=Side(style='thin',color=LINE)
    for row in ws.iter_rows(min_row=4,max_row=ws.max_row,min_col=1,max_col=6):
        for cell in row:
            cell.font=Font(name=FONT,size=10,color=FOREST)
            cell.alignment=Alignment(vertical='top',wrap_text=True)
            cell.border=Border(bottom=thin)


def build_quote():
    wb=Workbook(); ws=wb.active; ws.title="Quote"
    ws.append([]); ws.append(["파디엠 AI 미디어 업무전환 스튜디오 · 견적 검토표"]); ws.append(["DRAFT · 고객 범위 확인 전 · 실제 계약/매출 주장 아님"])
    headers=["코드","상품","가격 가설","기간","선택","범위/비고"]
    ws.append(headers)
    rows=[
        ["A","진단 워크숍","초기 300만–500만원 / 확장 500만–800만원","1–2일","","현재 workflow 진단 · 적용 후보 · 사람 검토 gate · 파일럿 후보"],
        ["B1","디자인 파트너 6주 파일럿","1,000만–1,500만원","6주","","1팀 · 1업무 · 학습형 범위"],
        ["B2","표준 6주 파일럿","1,500만–2,500만원","6주","","기준선 · workflow · 제한 파일럿 · KPI · 운영 산출물"],
        ["C","조직 운영 자문","월 300만–600만원","월","","정책/workflow/KPI 정기 리뷰"],
    ]
    for row in rows: ws.append(row)
    ws.append([]); ws.append(["고객 범위 확인"])
    fields=[("조직/팀",""),("바꿀 업무 1건",""),("참여 인원",""),("운영 책임자",""),("개인정보 포함 여부",""),("저작권/라이선스 확인",""),("선호 옵션",""),("고객별 승인 가격","미정"),("사업자 정보","발송 전 공식 정보 입력 필요"),("법률/계약 검토","필요")]
    for k,v in fields: ws.append([k,v])
    style_ws(ws,"파디엠 AI 미디어 업무전환 스튜디오 · 견적 검토표")
    for cell in ws[4]: cell.font=Font(name=FONT,size=10,bold=True,color=WHITE); cell.fill=PatternFill('solid',fgColor=COBALT); cell.alignment=Alignment(horizontal='center',vertical='center',wrap_text=True)
    ws.auto_filter.ref=f"A4:F{4+len(rows)}"

    terms=wb.create_sheet("Terms"); terms.sheet_view.showGridLines=False
    terms["A1"]="견적·계약 경계"; terms["A1"].font=Font(name=FONT,size=18,bold=True,color=FOREST)
    notes=["가격은 시장 검증 전 가설이며 고객 범위 확인 후 별도 승인합니다.","VAT 별도 여부는 최종 견적에서 명시합니다.","개인정보·저작권·조달은 고객별 확인이 필요합니다.","SOW/위험·데이터 부속서는 전문 법률·계약 검토가 필요합니다.","파디엠 공식 사업자 정보 입력 전 고객 발송 금지."]
    for i,n in enumerate(notes,3): terms[f"A{i}"]=f"• {n}"; terms[f"A{i}"].font=Font(name=FONT,size=11,color=FOREST); terms[f"A{i}"].alignment=Alignment(wrap_text=True,vertical='top')
    terms.column_dimensions['A'].width=90

    out=OUT/"Business35_V3_1_Pilot_Quote_Template.xlsx"; wb.save(out); return out


def main():
    built=[build_proposal(),build_one_page(),build_questionnaire(),build_quote()]
    print("GENERATED_V3_1_PACKAGE")
    for p in built: print(p)


if __name__ == "__main__":
    main()
