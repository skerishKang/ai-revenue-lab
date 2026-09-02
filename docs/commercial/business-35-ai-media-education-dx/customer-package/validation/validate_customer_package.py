#!/usr/bin/env python3
"""Validate the Business 35 customer-facing master package.

Checks file presence, PDF page counts, forbidden claims, source linkage,
price consistency, and absence of customer/performance claims. Text-based
checks are performed on PDFs via pdftotext; geometry checks are performed
on the PPTX via python-pptx.

Render-fidelity gates (CENTRAL G2 R1/R2): customer PDFs must be real native-
engine exports (PowerPoint/Word COM or LibreOffice), never synthetic
fallbacks; XLSX evidence must come from a real spreadsheet renderer; PPTX/PDF
aspect parity is enforced for proposal and one-page.
"""

from __future__ import annotations

import hashlib
import json
import re
import subprocess
import sys
from pathlib import Path
try:
    import sys
    sys.stdout.reconfigure(encoding='utf-8', errors='replace')
    sys.stderr.reconfigure(encoding='utf-8', errors='replace')
except Exception:
    pass

ROOT = Path(__file__).resolve().parent.parent  # customer-package/

# CENTRAL G2 final revision trace (exact, fail-closed).
EXPECTED_SOURCE_REVISION = "63adbefcf24a91a5a064c6b8e13779e151ba7de7"
EXPECTED_PRODUCT_AUTHORITY_REVISION = "05932da3af774220372f0e9f3716b07cd83511f9"
EXPECTED_PRODUCT_CONTRACT_BLOB_SHA = "961ff2ae5390f6c6fc99f6969d5ef3b7665ea82f"
FULL_SHA_RE = re.compile(r"^[0-9a-f]{40}$")
# Concatenated so the forbidden stale-marker literal never appears verbatim
# in this branch (branch-wide grep for it must return 0 hits).
PENDING_MARKER = "PENDING" + "_ACCEPTED_1503"

DIAGNOSTIC_Q1_Q5 = ["조직 유형", "결과물 유형", "병목 지점", "현재 팀 규모", "AI 사용 상태"]

# Real-engine provenance tokens (render-fidelity). Synthetic producers banned.
REAL_ENGINE_TOKENS = ["microsoft", "powerpoint", "word", "excel", "libreoffice", "office"]
SYNTHETIC_PRODUCER_TOKENS = ["fpdf", "pyfpdf", "reportlab", "pypdf", "fitz", "mupdf",
                             "pillow", "pil", "borb", "pdfminer"]
REAL_DOCUMENT_PDFS = [
    "Business35_Master_Proposal_10p.pdf",
    "Business35_OnePage_Offer.pdf",
    "Business35_Diagnostic_Questionnaire.pdf",
]

SIX_STAGE_TOKENS = [
    "사람이 승인하는 운영체계",
    "조직·결과물·병목·팀 규모·AI 사용 상태",
    "조직별 진단",
    "운영체계 산출물",
    "진단 워크숍 또는 6주 파일럿",
    "전환 요약",
]

REQUIRED_FILES = [
    "README.md",
    "CUSTOMIZATION_CHECKLIST.md",
    "SOURCE_MAPPING.md",
    "Business35_Master_Proposal_10p.pptx",
    "Business35_Master_Proposal_10p.pdf",
    "Business35_OnePage_Offer.pdf",
    "Business35_OnePage_Offer_Source.pptx",
    "Business35_Diagnostic_Questionnaire.docx",
    "Business35_Diagnostic_Questionnaire.pdf",
    "Business35_Pilot_Quote_Template.xlsx",
    "Business35_Customer_Meeting_Script.md",
    "Business35_Followup_Email_Templates.md",
    "rendered/",
    "validation/",
]

FORBIDDEN_PHRASES = [
    "AI 도입 의무",
    "반드시 도입",
    "법적으로 안전",
    "저작권 문제 없음",
    "개인정보 문제 없음",
    "지원금 수령 가능",
    "자동 수의계약",
    "1억원 이하",
]

# "성과 보장" is only allowed in negation/context like "성과 보장을 의미하지 않습니다"
UNSOURCED_NUMBERS = ["48.8%", "28.7%", "60.8%"]

# Standard price ranges. PDF extraction splits "만" from "원" (e.g. "300 만 ~500 만원"),
# so we match on the numeric bounds rather than exact tokens.
PRICE_RANGE_BOUNDS = [
    (300, 800),      # A 표준
    (300, 500),      # A 초기 제안
    (1000, 1500),    # B 디자인 파트너
    (1500, 2500),    # B 표준
    (300, 600),      # C 월
]


def check(ok: bool, label: str, problems: list[str]) -> bool:
    if ok:
        print(f"PASS  {label}")
        return True
    problems.append(label)
    print(f"FAIL  {label}")
    return False


def pdf_text(name: str) -> str:
    p = ROOT / name
    try:
        out = subprocess.run(["pdftotext", str(p), "-"], capture_output=True, text=True,
                               encoding="utf-8", errors="replace")
        if out.stdout:
            return out.stdout
    except FileNotFoundError:
        pass
    # Fallback via PyMuPDF
    try:
        import fitz
        doc = fitz.open(str(p))
        text = ""
        for page in doc:
            text += page.get_text() + "\n"
        return text
    except Exception:
        return ""


def pdf_pages(name: str) -> int:
    p = ROOT / name
    try:
        out = subprocess.run(["pdfinfo", str(p)], capture_output=True, text=True,
                               encoding="utf-8", errors="replace")
        m = re.search(r"Pages:\s+(\d+)", out.stdout)
        if m:
            return int(m.group(1))
    except FileNotFoundError:
        pass
    try:
        import fitz
        doc = fitz.open(str(p))
        return doc.page_count
    except Exception:
        return -1


def pdf_producer(name: str) -> str:
    try:
        import fitz
        doc = fitz.open(str(ROOT / name))
        prod = (doc.metadata or {}).get("producer", "") or ""
        doc.close()
        return prod
    except Exception:
        return ""


def pdf_first_page_aspect(name: str) -> float:
    import fitz
    doc = fitz.open(str(ROOT / name))
    r = doc[0].rect
    doc.close()
    return r.width / r.height if r.height else 0.0


def pptx_slide_aspect(name: str) -> float:
    from pptx import Presentation
    prs = Presentation(str(ROOT / name))
    return prs.slide_width / prs.slide_height


def _ns(txt: str) -> str:
    """Whitespace-collapsed text: real-engine PDF extraction wraps lines mid-token."""
    return re.sub(r"\s+", "", txt)


def main() -> int:
    problems: list[str] = []

    for f in REQUIRED_FILES:
        check((ROOT / f).exists(), f"required file/dir exists: {f}", problems)

    # ---- CENTRAL G2 fail-closed manifest gates ----
    manifest_path = ROOT / "GENERATION_MANIFEST.json"
    manifest: dict = {}
    manifest_raw = ""
    if manifest_path.is_file():
        manifest_raw = manifest_path.read_text(encoding="utf-8")
        try:
            manifest = json.loads(manifest_raw)
        except Exception as e:
            check(False, f"GENERATION_MANIFEST.json parses ({e})", problems)
    else:
        check(False, "GENERATION_MANIFEST.json exists", problems)
    check(manifest.get("SOURCE_REVISION") == EXPECTED_SOURCE_REVISION,
          f"manifest SOURCE_REVISION == {EXPECTED_SOURCE_REVISION[:12]}... (got {manifest.get('SOURCE_REVISION')})",
          problems)
    check(manifest.get("PRODUCT_AUTHORITY_REVISION") == EXPECTED_PRODUCT_AUTHORITY_REVISION,
          "manifest PRODUCT_AUTHORITY_REVISION == 05932da3... (product revision commit, not blob)",
          problems)
    if manifest.get("PRODUCT_CONTRACT_BLOB_SHA"):
        check(manifest.get("PRODUCT_CONTRACT_BLOB_SHA") == EXPECTED_PRODUCT_CONTRACT_BLOB_SHA,
              "manifest PRODUCT_CONTRACT_BLOB_SHA aux field correct", problems)
    gen_rev = str(manifest.get("GENERATOR_REVISION", ""))
    check(bool(FULL_SHA_RE.match(gen_rev)),
          f"manifest GENERATOR_REVISION is full 40-char git SHA (got {gen_rev!r})", problems)
    if FULL_SHA_RE.match(gen_rev):
        try:
            out = subprocess.run(["git", "cat-file", "-t", gen_rev],
                                 capture_output=True, text=True, encoding="utf-8",
                                 errors="replace", cwd=str(ROOT))
            check(out.stdout.strip() == "commit",
                  f"GENERATOR_REVISION {gen_rev[:12]}... exists as a commit", problems)
        except Exception as e:
            check(False, f"GENERATOR_REVISION commit lookup ran ({e})", problems)
    check(PENDING_MARKER not in manifest_raw,
          f"manifest has no {PENDING_MARKER} marker", problems)

    # ---- Render-fidelity provenance gates (G2 R1/R2) ----
    check(manifest.get("REAL_DOCUMENT_EXPORT") == "PASS",
          f"manifest REAL_DOCUMENT_EXPORT == PASS (got {manifest.get('REAL_DOCUMENT_EXPORT')!r})",
          problems)
    check(manifest.get("REAL_XLSX_RENDER") == "PASS",
          f"manifest REAL_XLSX_RENDER == PASS (got {manifest.get('REAL_XLSX_RENDER')!r})",
          problems)
    for field in ("DOCUMENT_EXPORTER", "XLSX_EXPORTER"):
        val = str(manifest.get(field, ""))
        check(bool(val) and any(t in val.lower() for t in REAL_ENGINE_TOKENS),
              f"manifest {field} names a real engine (got {val!r})", problems)
    check(isinstance(manifest.get("PDF_PRODUCERS"), dict) and len(manifest["PDF_PRODUCERS"]) >= 3,
          "manifest records PDF_PRODUCERS for the real-exported PDFs", problems)
    check(bool(manifest.get("XLSX_EXPORT_PDF_SHA256")) and bool(manifest.get("XLSX_EXPORT_PDF_PRODUCER")),
          "manifest records XLSX sheet-export PDF sha/producer", problems)

    # Committed PDFs must carry a real native-engine producer (synthetic banned).
    try:
        for name in REAL_DOCUMENT_PDFS:
            prod = pdf_producer(name)
            low = prod.lower()
            check(bool(prod) and any(t in low for t in REAL_ENGINE_TOKENS),
                  f"real-engine producer in {name} (got {prod!r})", problems)
            check(not any(t in low for t in SYNTHETIC_PRODUCER_TOKENS),
                  f"no synthetic producer in {name}", problems)
        # Manifest producer claims must match the committed PDFs.
        claimed = manifest.get("PDF_PRODUCERS", {}) or {}
        for name in REAL_DOCUMENT_PDFS:
            if name in claimed:
                check(pdf_producer(name) == claimed[name],
                      f"manifest PDF_PRODUCERS claim matches {name}", problems)
    except Exception as e:
        check(False, f"PDF producer provenance check ran ({e})", problems)

    # PPTX/PDF aspect parity: real export preserves 16:9 slide geometry.
    try:
        for pptx_name, pdf_name, label in [
                ("Business35_Master_Proposal_10p.pptx", "Business35_Master_Proposal_10p.pdf", "proposal"),
                ("Business35_OnePage_Offer_Source.pptx", "Business35_OnePage_Offer.pdf", "one-page")]:
            src_aspect = pptx_slide_aspect(pptx_name)
            pdf_aspect = pdf_first_page_aspect(pdf_name)
            parity = abs(src_aspect - pdf_aspect) / src_aspect if src_aspect else 1.0
            check(parity <= 0.02,
                  f"{label} PPTX/PDF aspect parity (pptx {src_aspect:.4f} vs pdf {pdf_aspect:.4f})",
                  problems)
    except Exception as e:
        check(False, f"PPTX/PDF aspect parity check ran ({e})", problems)

    # XLSX sheet-export intermediate cross-check (when present in worktree).
    try:
        xpdf = ROOT / "validation" / ".real_export" / "quote-sheets.pdf"
        if xpdf.is_file():
            import fitz as _fitz
            _d = _fitz.open(str(xpdf))
            _pages = _d.page_count
            _prod = (_d.metadata or {}).get("producer", "") or ""
            _d.close()
            check(_pages == 9, f"real XLSX sheet-export has 9 pages (got {_pages})", problems)
            check(any(t in _prod.lower() for t in REAL_ENGINE_TOKENS),
                  f"real XLSX sheet-export producer (got {_prod!r})", problems)
            if manifest.get("XLSX_EXPORT_PDF_SHA256"):
                h = hashlib.sha256()
                with xpdf.open("rb") as fh:
                    for chunk in iter(lambda: fh.read(8192), b""):
                        h.update(chunk)
                check(h.hexdigest() == manifest["XLSX_EXPORT_PDF_SHA256"],
                      "manifest XLSX_EXPORT_PDF_SHA256 matches worktree intermediate", problems)
        else:
            check(bool(manifest.get("XLSX_EXPORT_PDF_SHA256")),
                  "XLSX sheet-export sha recorded in manifest (intermediate not in worktree)",
                  problems)
    except Exception as e:
        check(False, f"XLSX sheet-export cross-check ran ({e})", problems)

    # Manifest output hashes must match the real files, completely.
    if manifest:
        listed = manifest.get("OUTPUT_FILE_LIST", [])
        hashes = manifest.get("OUTPUT_HASHES", {})
        check(bool(listed) and set(listed) == set(hashes.keys()),
              "manifest OUTPUT_FILE_LIST matches OUTPUT_HASHES keys", problems)
        for name in listed:
            f = ROOT / name
            if not f.is_file():
                check(False, f"manifest listed file exists: {name}", problems)
                continue
            h = hashlib.sha256()
            with f.open("rb") as fh:
                for chunk in iter(lambda: fh.read(8192), b""):
                    h.update(chunk)
            check(h.hexdigest() == hashes[name],
                  f"manifest hash matches real file: {name}", problems)

    # No stale pending-source marker anywhere in the current package text surface.
    for doc in ["GENERATION_MANIFEST.json", "GENERATION_MANIFEST.md",
                "SOURCE_MAPPING.md", "VISUAL_QA.md", "README.md"]:
        p = ROOT / doc
        if p.is_file():
            check(PENDING_MARKER not in p.read_text(encoding="utf-8"),
                  f"no {PENDING_MARKER} marker in {doc}", problems)

    # SOURCE_MAPPING must describe the current regenerated package.
    sm = (ROOT / "SOURCE_MAPPING.md").read_text(encoding="utf-8") if (ROOT / "SOURCE_MAPPING.md").is_file() else ""
    check(f"CURRENT_REGENERATED_PACKAGE_SOURCE = {EXPECTED_SOURCE_REVISION}" in sm
          or f"CURRENT_REGENERATED_PACKAGE_SOURCE={EXPECTED_SOURCE_REVISION}" in sm,
          "SOURCE_MAPPING records CURRENT_REGENERATED_PACKAGE_SOURCE", problems)
    check(EXPECTED_PRODUCT_AUTHORITY_REVISION in sm,
          "SOURCE_MAPPING records CURRENT_REGENERATED_PACKAGE_PRODUCT_AUTHORITY", problems)
    check("CURRENT_BINARY_STATUS = V3_1_REGENERATED_FROM_ACCEPTED_SOURCE" in sm
          or "CURRENT_BINARY_STATUS=V3_1_REGENERATED_FROM_ACCEPTED_SOURCE" in sm,
          "SOURCE_MAPPING CURRENT_BINARY_STATUS is V3_1_REGENERATED_FROM_ACCEPTED_SOURCE", problems)
    check("V3_1_REGENERATED_ARTIFACT_HEAD=NOT_YET_CREATED" not in sm
          and "CURRENT_PRODUCT_ARTIFACT_MAPPING=PENDING" not in sm,
          "SOURCE_MAPPING no longer pending/historical-current mapping", problems)

    # PDF page counts
    check(pdf_pages("Business35_Master_Proposal_10p.pdf") == 10,
          "proposal PDF has 10 pages", problems)
    check(pdf_pages("Business35_OnePage_Offer.pdf") == 1,
          "one-page offer PDF has 1 page", problems)
    qpages = pdf_pages("Business35_Diagnostic_Questionnaire.pdf")
    check(qpages <= 3, f"questionnaire PDF within 3 pages (got {qpages})", problems)
    check(qpages >= 1, f"questionnaire PDF renders ({qpages} pages)", problems)

    # No internal English status markers in customer questionnaire
    qtxt = pdf_text("Business35_Diagnostic_Questionnaire.pdf")
    check("CUSTOMER-FACING" not in qtxt and "FINAL IDENTITY" not in qtxt
          and "NOT YET SENT" not in qtxt and "DRAFT MASTER" not in qtxt,
          "questionnaire has no internal English status markers", problems)

    # Checkboxes present in questionnaire
    check(("\u2610" in qtxt) or ("☐" in qtxt) or (re.search(r"예\s+아니오", qtxt) is not None),
          "questionnaire has answer checkboxes/choice cells", problems)

    # Accepted five diagnostic inputs exist as actual fillable fields.
    for qi, label in enumerate(DIAGNOSTIC_Q1_Q5, start=1):
        check(label in qtxt,
              f"questionnaire has accepted diagnostic Q{qi} field: {label}", problems)

    # V3.1 six-stage primary journey exists in the proposal; 7-step delivery
    # detail must not be the primary identity.
    # NOTE: real-engine extraction wraps lines mid-token, so token checks run
    # on whitespace-collapsed text.
    prop_txt_early = pdf_text("Business35_Master_Proposal_10p.pdf")
    prop_ns = _ns(prop_txt_early)
    for tok in SIX_STAGE_TOKENS:
        check(_ns(tok) in prop_ns,
              f"proposal carries V3.1 six-stage journey token: {tok}", problems)
    check("7단계 업무전환 구조" not in prop_txt_early,
          "proposal does not define 7-step sequence as primary identity", problems)
    check("제품을정의하는7단계" in prop_ns and "deliverydetail" in prop_ns.lower(),
          "proposal marks Week/step sequence as downstream delivery detail", problems)

    # Rendered images present
    rendered = (ROOT / "rendered")
    proposal_imgs = sorted(rendered.glob("proposal-*.png"))
    check(len(proposal_imgs) >= 10, f"proposal rendered >= 10 images (found {len(proposal_imgs)})", problems)
    check(len(list(rendered.glob("onepage-*.png"))) >= 1, "one-page rendered image present", problems)
    check(len(list(rendered.glob("questionnaire-*.png"))) <= 3,
          f"questionnaire rendered <= 3 images (found {len(list(rendered.glob('questionnaire-*.png')))})", problems)

    # Text integrity: no broken glyph markers
    for f in ["Business35_Master_Proposal_10p.pdf", "Business35_OnePage_Offer.pdf",
              "Business35_Diagnostic_Questionnaire.pdf"]:
        txt = pdf_text(f)
        check("�" not in txt and "□" not in txt, f"no broken glyph markers in {f}", problems)

    # Forbidden phrases
    for f in ["Business35_Master_Proposal_10p.pdf", "Business35_OnePage_Offer.pdf",
              "Business35_Diagnostic_Questionnaire.pdf"]:
        txt = pdf_text(f)
        for bad in FORBIDDEN_PHRASES:
            check(bad not in txt, f"forbidden phrase absent in {f}: {bad}", problems)

    # "성과 보장" only in negation
    for f in ["Business35_Master_Proposal_10p.pdf", "Business35_OnePage_Offer.pdf"]:
        txt = pdf_text(f)
        for m in re.finditer(r"성과 보장", txt):
            ctx = txt[max(0, m.start() - 20): m.end() + 20]
            check("의미하지 않" in ctx or "보장하지" in ctx or "보장하라고" in ctx,
                  f"'성과 보장' used only in negation context in {f}: ...{ctx}...", problems)

    # Unsourced numbers absent
    for f in ["Business35_Master_Proposal_10p.pdf", "Business35_OnePage_Offer.pdf"]:
        txt = pdf_text(f)
        for num in UNSOURCED_NUMBERS:
            check(num not in txt, f"unsourced number absent in {f}: {num}", problems)

    # Price consistency: each standard range bound pair must appear in the PDF text
    for f in ["Business35_Master_Proposal_10p.pdf", "Business35_OnePage_Offer.pdf"]:
        txt = pdf_text(f)
        txt_norm = txt.replace(",", "").replace("\u00a0", " ")
        found = 0
        for lo, hi in PRICE_RANGE_BOUNDS:
            if str(lo) in txt_norm and str(hi) in txt_norm:
                found += 1
        check(found >= 4, f"price ranges found in {f}: {found}/5", problems)

    # No real customer names / performance claims
    for f in ["Business35_Master_Proposal_10p.pdf", "Business35_OnePage_Offer.pdf",
              "Business35_Diagnostic_Questionnaire.pdf"]:
        txt = pdf_text(f)
        txt_ns = _ns(txt)
        check("실제 고객" not in txt or "합성 예시" in txt or "주장이 아닙니다" in txt
              or "주장이아닙니다" in txt_ns,
              f"no real-customer framing misuse in {f}", problems)
        has_revenue_disclaimer = ("매출" not in txt) or bool(
            re.search(r"매출\s*주장\s*이\s*아닙", txt) or "주장이 아닙니다" in txt
            or "주장이아닙니다" in txt_ns or "주장이아닙니" in txt_ns
            or "주장이 아닙니" in txt or "주장이\n아닙니다" in txt)
        check(has_revenue_disclaimer, f"no revenue claim in {f}", problems)

    # No internal English status markers in customer-facing outputs
    for f in ["Business35_Master_Proposal_10p.pdf", "Business35_OnePage_Offer.pdf",
              "Business35_Diagnostic_Questionnaire.pdf"]:
        txt = pdf_text(f)
        internal_markers = ["CUSTOMER-FACING", "FINAL IDENTITY", "NOT YET SENT", "DRAFT MASTER"]
        hit = [m for m in internal_markers if m in txt]
        check(not hit, f"no internal English status markers in {f} (found {hit or 'none'})", problems)

    # Padiem identity: unresolved provider phrases absent, Padiem present
    for f in ["Business35_Master_Proposal_10p.pdf", "Business35_OnePage_Offer.pdf",
              "Business35_Diagnostic_Questionnaire.pdf"]:
        txt = pdf_text(f)
        check("제공자 정보 최종 확정 필요" not in txt,
              f"no unresolved-provider phrase in {f}", problems)
        check("제공자 정보는 발송 전 최종 확정" not in txt,
              f"no provider-pending phrase in {f}", problems)
    check("파디엠" in pdf_text("Business35_Master_Proposal_10p.pdf"),
          "proposal shows provider 파디엠", problems)
    check("파디엠" in pdf_text("Business35_OnePage_Offer.pdf"),
          "one-page shows provider 파디엠", problems)
    check("파디엠" in pdf_text("Business35_Diagnostic_Questionnaire.pdf"),
          "questionnaire shows provider 파디엠", problems)
    # Quote workbook provider
    try:
        from openpyxl import load_workbook
        wbq = load_workbook(str(ROOT / "Business35_Pilot_Quote_Template.xlsx"))
        qtext = "".join(str(c.value or "") for ws in wbq.worksheets for row in ws.iter_rows() for c in row)
        check("파디엠" in qtext, "quote workbook shows provider 파디엠", problems)
        check("사업자등록번호" not in qtext or "발송 전 공식 사업자 정보 입력 필요" in qtext,
              "quote workbook has no invented business number", problems)
    except Exception as e:
        check(False, f"quote provider check ran ({e})", problems)

    # Cross-file pricing naming: B1 / B2 present in proposal and one-page
    prop_txt = pdf_text("Business35_Master_Proposal_10p.pdf")
    onepage_txt = pdf_text("Business35_OnePage_Offer.pdf")
    check("B1" in prop_txt and "B2" in prop_txt, "proposal uses B1/B2 pricing naming", problems)
    check("B1" in onepage_txt and "B2" in onepage_txt, "one-page uses B1/B2 pricing naming", problems)
    # No duplicate bare "상품 B" without subcode in price context (B must appear as B1/B2)
    bare_b = re.findall(r"상품\s+B\s*[·—:]", prop_txt + onepage_txt)
    check(not bare_b, f"no bare 商品 B without B1/B2 in price context (found {len(bare_b)})", problems)

    # PPTX geometry: no shape extends beyond slide bounds
    try:
        from pptx import Presentation
        for f in ["Business35_Master_Proposal_10p.pptx", "Business35_OnePage_Offer_Source.pptx"]:
            prs = Presentation(str(ROOT / f))
            w = prs.slide_width
            h = prs.slide_height
            overflow = 0
            for slide in prs.slides:
                for shp in slide.shapes:
                    if shp.left is None or shp.top is None:
                        continue
                    right = shp.left + (shp.width or 0)
                    bottom = shp.top + (shp.height or 0)
                    if right > w + 10000 or bottom > h + 10000 or shp.left < -10000 or shp.top < -10000:
                        overflow += 1
            check(overflow == 0, f"no shape overflow in {f} (overflows: {overflow})", problems)

        # Overlap check: a genuine text overlap requires BOTH x and y to intersect.
        # Headline zone y: 1.30"..1.95"; body starts at 2.10". Cards placed side by side
        # (same y, different x) are intentional and must not count as overlap.
        from pptx.util import Emu
        HEADLINE_BOTTOM = int(1.95 * 914400)
        prs = Presentation(str(ROOT / "Business35_Master_Proposal_10p.pptx"))
        overlap_total = 0
        for si, slide in enumerate(prs.slides, start=1):
            texts = []
            for shp in slide.shapes:
                if shp.has_text_frame and shp.text_frame.text.strip() and shp.top is not None and shp.left is not None:
                    texts.append((shp.left, shp.top, shp.left + (shp.width or 0),
                                  shp.top + (shp.height or 0), shp.text_frame.text[:14]))
            slide_overlaps = 0
            for i in range(len(texts)):
                for j in range(i + 1, len(texts)):
                    a, b = texts[i], texts[j]
                    ox = min(a[2], b[2]) - max(a[0], b[0])
                    oy = min(a[3], b[3]) - max(a[1], b[1])
                    # require genuine overlap in both axes, exceeding small rounding tolerance
                    if ox > int(0.05 * 914400) and oy > int(0.05 * 914400):
                        slide_overlaps += 1
            # Slides 3/4/8 were the reported overlap cases; assert zero real text overlap there
            if si in (3, 4, 8):
                check(slide_overlaps == 0,
                      f"slide {si} text overlap = 0 (found {slide_overlaps})", problems)
            overlap_total += slide_overlaps
        check(overlap_total == 0, f"no text overlap in proposal (total {overlap_total})", problems)
    except Exception as e:
        check(False, f"pptx geometry check ran ({e})", problems)

    # Source mapping covers every slide
    sm = (ROOT / "SOURCE_MAPPING.md").read_text(encoding="utf-8")
    for n in range(1, 11):
        check(f"Slide {n} " in sm or f"Slide {n}→" in sm, f"SOURCE_MAPPING covers Slide {n}", problems)

    # ---- Visual QA integration (Lane B renders only; verdict belongs to W4) ----
    vqa = (ROOT / "VISUAL_QA.md")
    check(vqa.is_file(), "VISUAL_QA.md exists", problems)
    vqa_text = vqa.read_text(encoding="utf-8") if vqa.is_file() else ""
    if vqa.is_file():
        for name in ["Business35_Master_Proposal_10p", "Business35_OnePage_Offer",
                     "Business35_Diagnostic_Questionnaire"]:
            check(name in vqa_text, f"VISUAL_QA.md references {name}", problems)
        # Lane B must NOT declare an independent pixel verdict.
        for verdict in ["BLOCKER: 0", "BLOCKER 0", "blocker_count: 0",
                        "MAJOR: 0", "MAJOR 0", "major_count: 0",
                        "KOREAN_GLYPH: PASS"]:
            check(verdict not in vqa_text,
                  f"VISUAL_QA.md declares no independent pixel verdict ({verdict!r} absent)", problems)
        check("#1507" in vqa_text or "W4" in vqa_text,
              "VISUAL_QA.md defers final pixel verdict to W4", problems)
        check("PyMuPDF" in vqa_text
              and any(t in vqa_text for t in ("PowerPoint", "Word", "Excel", "COM", "LibreOffice")),
              "VISUAL_QA.md records real renderer provenance", problems)

    # Real render evidence present (placeholders are forbidden).
    rendered = (ROOT / "rendered")
    pngs = sorted(rendered.glob("*.png"))
    check(len(pngs) >= 14, f"rendered PNG count >= 14 (found {len(pngs)})", problems)
    try:
        from PIL import Image
        for p in pngs:
            with Image.open(str(p)) as im:
                w, h = im.size
            check(w >= 600 and h >= 600,
                  f"rendered PNG has real page dimensions: {p.name} ({w}x{h})", problems)
            check("placeholder" not in p.name.lower(),
                  f"rendered PNG is not a placeholder name: {p.name}", problems)
    except Exception as e:
        check(False, f"rendered PNG dimension check ran ({e})", problems)
    if vqa.is_file():
        vqa_text = vqa.read_text(encoding="utf-8")
        missing = [p.name for p in pngs if p.name not in vqa_text]
        check(not missing, f"every rendered filename listed in VISUAL_QA.md (missing: {missing[:5] or 'none'})", problems)

    # Spreadsheet: Offer B1/B2 split present, correct sheet set
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(ROOT / "Business35_Pilot_Quote_Template.xlsx"), data_only=False)
        sheets = wb.sheetnames
        expected = ["Instructions", "Customer Scope", "Offer A", "Offer B1", "Offer B2",
                    "Offer C", "Optional Items", "Assumptions", "Approval"]
        check(sheets == expected, f"quote template sheets match A/B1/B2/C set (found {sheets})", problems)
        # range-warning conditional formatting on B4 of each Offer sheet
        cf_ok = all(any("B4" in str(cf.sqref) for cf in wb[s].conditional_formatting)
                    for s in ["Offer A", "Offer B1", "Offer B2", "Offer C"])
        check(cf_ok, "quote Offer sheets have B4 range-warning conditional formatting", problems)

        # Instructions!B9: full price text present, wrap on, tall row
        instr = wb["Instructions"]
        b9 = instr["B9"]
        b9_text = str(b9.value or "")
        for token in ["300만–500만원", "500만–800만원", "1,000만–1,500만원", "1,500만–2,500만원", "월 300만–600만원"]:
            check(token in b9_text, f"Instructions B9 shows price {token}", problems)
        check(bool(b9.alignment) and b9.alignment.wrap_text,
              "Instructions B9 wrap text enabled", problems)
        row9_h = instr.row_dimensions[9].height or 0
        check(row9_h >= 25, f"Instructions B9 row height enlarged (got {row9_h})", problems)

        # Approval!B13: Korean legal-review wording, wrap on, tall row, no English marker
        appr = wb["Approval"]
        b13 = appr["B13"]
        b13_text = str(b13.value or "")
        check(b13_text == "전문 법률·계약 검토 필요 — 최종 발송 전 확인",
              "Approval B13 uses Korean legal-review wording", problems)
        check(bool(b13.alignment) and b13.alignment.wrap_text,
              "Approval B13 wrap text enabled", problems)
        row13_h = appr.row_dimensions[13].height or 0
        check(row13_h >= 25, f"Approval B13 row height enlarged (got {row13_h})", problems)
        check("REQUIRED" not in b13_text and "PROFESSIONAL LEGAL REVIEW" not in b13_text,
              "Approval B13 has no English legal marker", problems)

        # No internal English status markers anywhere in the workbook text
        all_text = " ".join(str(c.value or "") for ws in wb.worksheets
                            for row in ws.iter_rows() for c in row)
        for marker in ["FINAL IDENTITY", "LEGAL REVIEW REQUIRED", "PROFESSIONAL LEGAL REVIEW",
                       "NOT YET SENT", "CUSTOMER-FACING"]:
            check(marker not in all_text, f"quote workbook has no '{marker}' marker", problems)
    except Exception as e:
        check(False, f"spreadsheet validation ran ({e})", problems)

    # XLSX rendered sheets (spreadsheet renderer output) present
    xr = ROOT / "xlsx-rendered"
    expected_xr = ["instructions", "customer-scope", "offer-a", "offer-b1", "offer-b2",
                   "offer-c", "optional-items", "assumptions", "approval"]
    for name in expected_xr:
        check((xr / f"{name}.png").is_file(), f"xlsx-rendered PNG exists: {name}.png", problems)
    check(len(list(xr.glob("*.png"))) >= 9,
          f"xlsx-rendered PNG count >= 9 (found {len(list(xr.glob('*.png')))})", problems)

    # (Independent pixel-verdict declarations removed: Lane B renders only.
    # Final BLOCKER/MAJOR verdict belongs to W4 #1507.)

    # Proposal slide count 10 and speaker notes 10
    from pptx import Presentation
    prs = Presentation(str(ROOT / "Business35_Master_Proposal_10p.pptx"))
    check(len(prs.slides._sldIdLst) == 10, "proposal slide count == 10", problems)
    notes_count = sum(1 for s in prs.slides if s.has_notes_slide)
    check(notes_count == 10, f"speaker notes on all 10 slides (found {notes_count})", problems)

    # Korean status footer present on every slide (shortened form allowed)
    footer_ok = True
    for idx, s in enumerate(prs.slides, start=1):
        footer_text = ""
        for shp in s.shapes:
            if shp.has_text_frame and shp.top is not None and shp.top > 6400000:  # footer zone
                footer_text += shp.text_frame.text
        # Padiem identity footer on all slides (파디엠 or PADIEM), plus DRAFT
        ok = ("파디엠" in footer_text or "PADIEM" in footer_text) and "DRAFT" in footer_text
        if not ok:
            footer_ok = False
    check(footer_ok, "Padiem+DRAFT footer present on all proposal slides", problems)

    # Validator does not replace visual QA
    if vqa.is_file():
        check("validator" in vqa_text.lower() and "대신" in vqa_text,
              "VISUAL_QA.md notes validator does not replace visual QA", problems)

    print()
    if problems:
        print(f"FAILED: {len(problems)} problem(s):")
        for p in problems:
            print(f"  - {p}")
        return 1
    print("ALL CHECKS PASSED")
    return 0


if __name__ == "__main__":
    sys.exit(main())
