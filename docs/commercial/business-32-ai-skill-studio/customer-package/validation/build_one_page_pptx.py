"""Build Business32_OnePage_Offer_Source.pptx (exactly 1 slide)."""
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from _b32theme import *  # noqa: F403
from pptx import Presentation
from pptx.util import Inches

OUT = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
    "Business32_OnePage_Offer_Source.pptx",
)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

s = blank_slide(prs)
rect(s, Inches(0.0), Inches(0.0), Inches(13.333), Inches(0.14), ORANGE)
rect(s, Inches(0.0), Inches(7.36), Inches(13.333), Inches(0.14), BLUE)

t = s.shapes.add_textbox(Inches(0.5), Inches(0.28), Inches(12.4), Inches(0.5))
set_run(t.text_frame.paragraphs[0].add_run(),
        "Business 32 · AI 업무 스킬 전환 프로그램 — 10초 요약", size=16, bold=True, color=BLUE)
cap = s.shapes.add_textbox(Inches(0.5), Inches(0.62), Inches(12.4), Inches(0.3))
set_run(cap.text_frame.paragraphs[0].add_run(),
        "제안검토용 · 합성 샘플 기반 · 사람 검토와 승인 필수", size=11, color=GRAY)

rect(s, Inches(0.5), Inches(0.95), Inches(12.35), Inches(0.7), PAPER, line_color=BLUE,
     text="누구를 위한 것: 반복업무를 수행하는 조직과 팀", size=13, text_color=INK)
rect(s, Inches(0.5), Inches(1.75), Inches(12.35), Inches(0.7), PAPER, line_color=RED,
     text="문제: AI를 개인별·일회성으로 사용해 결과물 품질이 사람마다 다름", size=13, text_color=INK)
rect(s, Inches(0.5), Inches(2.55), Inches(12.35), Inches(0.7), PAPER, line_color=GREEN,
     text="업무 스킬: 입력자료·단계·증거·검토·예외·승인 기준이 붙은 재사용 가능한 절차", size=13, text_color=INK)

offers = [
    ("Offer A — 설계 워크숍", "1~2일 · 업무 1개 · 300만~500만원"),
    ("Offer B — 전환 스프린트", "2~3주 · 업무 1개 · 500만~800만원"),
    ("Offer C — 팀 라이브러리 파일럿", "4~6주 · 업무 3개 · 팀 1개 · 1,200만~2,000만원"),
]
for i, (name, desc) in enumerate(offers):
    left = Inches(0.5 + i * 4.2)
    rect(s, left, Inches(3.4), Inches(4.0), Inches(1.3), BLUE if i == 1 else GRAY,
         text="", size=12)
    nb = s.shapes.add_textbox(left + Inches(0.2), Inches(3.55), Inches(3.6), Inches(0.4))
    set_run(nb.text_frame.paragraphs[0].add_run(), name, size=14, bold=True, color=WHITE)
    db = s.shapes.add_textbox(left + Inches(0.2), Inches(4.0), Inches(3.6), Inches(0.6))
    set_run(db.text_frame.paragraphs[0].add_run(), desc, size=12, color=WHITE)

rect(s, Inches(0.5), Inches(4.95), Inches(12.35), Inches(0.7), PAPER, line_color=ORANGE,
     text="첫 행동: 반복업무 1개를 Skill Discovery Worksheet으로 선정 (실제 파일 불필요)",
     size=13, text_color=INK)
rect(s, Inches(0.5), Inches(5.75), Inches(12.35), Inches(0.7), PAPER, line_color=YELLOW,
     text="가격은 가설입니다. 확정 가격이 아닙니다.", size=13, text_color=INK)
rect(s, Inches(0.5), Inches(6.55), Inches(12.35), Inches(0.7), PAPER, line_color=GREEN,
     text="사람 검토와 승인이 필수입니다. AI가 자동 승인하지 않습니다.", size=13, text_color=INK)

add_footer(s, "cover", prs)

assert len(prs.slides._sldIdLst) == 1, "one-page offer must be exactly 1 slide"
prs.save(OUT)
print("saved", OUT)
