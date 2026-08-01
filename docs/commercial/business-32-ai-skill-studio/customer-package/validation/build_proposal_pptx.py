"""Build Business32_Master_Proposal_10p.pptx (exactly 10 slides)."""
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from _b32theme import *  # noqa: F403
from pptx import Presentation
from pptx.util import Inches

OUT = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
    "Business32_Master_Proposal_10p.pptx",
)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def cover():
    s = blank_slide(prs)
    rect(s, Inches(0.0), Inches(0.0), Inches(13.333), Inches(0.18), ORANGE)
    rect(s, Inches(0.0), Inches(7.32), Inches(13.333), Inches(0.18), BLUE)
    t = s.shapes.add_textbox(Inches(1.2), Inches(1.9), Inches(10.9), Inches(2.2))
    tf = t.text_frame
    tf.word_wrap = True
    set_run(tf.paragraphs[0].add_run(), "조직의 반복업무를", size=46, bold=True, color=INK)
    p = tf.add_paragraph()
    set_run(p.add_run(), "검증된 AI 업무 스킬로", size=46, bold=True, color=BLUE)
    sub = s.shapes.add_textbox(Inches(1.2), Inches(4.2), Inches(10.9), Inches(1.0))
    set_run(sub.text_frame.paragraphs[0].add_run(),
            "입력자료·단계·증거·검토·예외·승인 기준이 포함된 재사용 가능한 AI 업무 스킬 전환 프로그램",
            size=18, color=GRAY)
    idbox = s.shapes.add_textbox(Inches(1.2), Inches(5.4), Inches(10.9), Inches(0.6))
    set_run(idbox.text_frame.paragraphs[0].add_run(),
            "Business 32 · AI Skill Studio — AI 업무 스킬 전환 프로그램", size=14, bold=True, color=ORANGE)
    note = s.shapes.add_textbox(Inches(1.2), Inches(6.0), Inches(10.9), Inches(0.5))
    set_run(note.text_frame.paragraphs[0].add_run(),
            "제안검토용 DRAFT · 합성 샘플 기반 · 가격은 가설", size=12, color=GRAY)
    add_footer(s, "cover", prs)


def problem():
    s = blank_slide(prs)
    add_header(s, 2, "현재 문제 — 개인별 프롬프트와 일회성 AI 사용", prs)
    flow = [
        ("개인별 AI 사용", "팀 자산으로 축적되지 않음"),
        ("품질 편차", "결과물이 사람마다 다름"),
        ("근거 부재", "만든 근거를 확인할 수 없음"),
        ("재작업·재시작", "같은 업무를 반복할 때마다 다시 시작"),
    ]
    for i, (title, sub) in enumerate(flow):
        left = Inches(0.6 + i * 3.13)
        rect(s, left, Inches(1.45), Inches(2.85), Inches(1.45), PAPER, line_color=BLUE)
        tb = s.shapes.add_textbox(left + Inches(0.15), Inches(1.6), Inches(2.55), Inches(0.4))
        set_run(tb.text_frame.paragraphs[0].add_run(), title, size=14, bold=True, color=INK)
        sb = s.shapes.add_textbox(left + Inches(0.15), Inches(2.02), Inches(2.55), Inches(0.8))
        set_run(sb.text_frame.paragraphs[0].add_run(), sub, size=11, color=GRAY)
        if i < 3:
            ab = s.shapes.add_textbox(left + Inches(2.9), Inches(1.95), Inches(0.25), Inches(0.5))
            set_run(ab.text_frame.paragraphs[0].add_run(), "→", size=20, bold=True, color=BLUE)
    bullets(s, Inches(0.6), Inches(3.15), Inches(12.0), Inches(2.4), [
        "프롬프트가 개인 기억에만 남아 팀 자산으로 축적되지 않습니다.",
        "검토·승인 기준이 명확하지 않아 결과물을 다시 고쳐야 합니다.",
        "실수를 발견해도 다음에 같은 실수가 반복됩니다.",
        "AI가 무엇을 근거로 만들었는지 확인할 방법이 없습니다.",
    ], size=14)
    rect(s, Inches(0.6), Inches(5.8), Inches(12.1), Inches(1.0), PAPER, line_color=RED,
         text="해결 방향: 반복업무 하나를, 증거와 검토·승인 기준이 붙은 재사용 가능한 절차로 전환합니다.",
         text_color=INK)


def why_generic():
    s = blank_slide(prs)
    add_header(s, 3, "왜 일반 AI 교육만으로는 부족한가", prs)
    generic = [
        "'할 수 있는 일'을 보여주지만 특정 업무를 고정하지 않습니다.",
        "교육 후 개인별 사용으로 돌아가 조직 절차로 남지 않습니다.",
        "업무별 입력자료·증거·예외·승인 기준은 조직마다 달라 규격화되지 않습니다.",
        "사람 검토·승인 흐름을 업무 절차로 묶지 않습니다.",
    ]
    b32 = [
        "반복업무 1개를 고객과 함께 선정하고 고정합니다.",
        "입력자료·단계·증거·검토·예외·승인 기준을 스킬에 포함합니다.",
        "검토·승인 흐름을 업무 절차로 통합합니다.",
        "증거와 함께 재사용 가능한 스킬 카드로 납품합니다.",
    ]
    rect(s, Inches(0.6), Inches(1.45), Inches(5.9), Inches(3.9), PAPER, line_color=RED)
    hb = s.shapes.add_textbox(Inches(0.85), Inches(1.6), Inches(5.4), Inches(0.4))
    set_run(hb.text_frame.paragraphs[0].add_run(), "일반 AI 교육", size=15, bold=True, color=RED)
    bullets(s, Inches(0.85), Inches(2.1), Inches(5.4), Inches(3.0), generic, size=12)
    rect(s, Inches(6.85), Inches(1.45), Inches(5.9), Inches(3.9), PAPER, line_color=GREEN)
    hb = s.shapes.add_textbox(Inches(7.1), Inches(1.6), Inches(5.4), Inches(0.4))
    set_run(hb.text_frame.paragraphs[0].add_run(), "Business 32 전환 프로그램", size=15, bold=True, color=GREEN)
    bullets(s, Inches(7.1), Inches(2.1), Inches(5.4), Inches(3.0), b32, size=12)
    rect(s, Inches(0.6), Inches(5.6), Inches(12.15), Inches(1.0), BLUE,
         text="핵심 차이: '할 수 있는 일'이 아니라, 조직 업무 1개를 증거와 검토·승인 기준이 붙은 "
              "재사용 가능한 스킬로 전환합니다.",
         size=13)


def method():
    s = blank_slide(prs)
    add_header(s, 4, "Business 32의 전환 방식", prs)
    steps = [
        ("01", "반복업무 선정", "업무 1개를 고객과 함께 선정"),
        ("02", "입력·증거 정의", "필수 입력자료와 증거 요구사항"),
        ("03", "단계 설계", "실행·AI 보조·사람 판단 단계"),
        ("04", "검토·승인 정의", "검토 기준·승인 기준·예외"),
        ("05", "합성 실행", "검증된 프론트엔드로 합성 실행"),
        ("06", "스킬 카드 납품", "검증된 조직 AI 업무 스킬 패키지"),
    ]
    for i, (no, title, sub) in enumerate(steps):
        col = i % 3
        row = i // 3
        left = Inches(0.7 + col * 4.15)
        top = Inches(1.5 + row * 2.4)
        rect(s, left, top, Inches(3.9), Inches(2.0), PAPER, line_color=BLUE)
        nb = s.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), Inches(1.4), Inches(0.6))
        set_run(nb.text_frame.paragraphs[0].add_run(), no, size=28, bold=True, color=BLUE)
        tb = s.shapes.add_textbox(left + Inches(0.2), top + Inches(0.9), Inches(3.5), Inches(0.5))
        set_run(tb.text_frame.paragraphs[0].add_run(), title, size=16, bold=True, color=INK)
        sb = s.shapes.add_textbox(left + Inches(0.2), top + Inches(1.35), Inches(3.5), Inches(0.5))
        set_run(sb.text_frame.paragraphs[0].add_run(), sub, size=12, color=GRAY)


def components():
    s = blank_slide(prs)
    add_header(s, 5, "검증된 업무 스킬의 구성 요소", prs)
    items = [
        ("업무 목적", "이 스킬이 해결하는 업무"), ("사용자와 검토자", "실행자·검토자 역할"),
        ("입력자료", "필수 입력과 확인"), ("실행 단계", "업무 수행 순서"),
        ("AI 보조 단계", "AI가 돕는 단계"), ("사람 판단 단계", "사람이 결정하는 단계"),
        ("필수 증거", "근거 연결"), ("검토·승인 기준", "검토 체크와 승인"),
        ("누락·충돌 처리", "자동 추정·자동 판정 금지"), ("예외·금지 사용", "예외 경고와 금지 기준"),
        ("재실행 조건", "수정·재실행"), ("버전·다음 검토일", "분기 검토·버전 갱신"),
    ]
    for i, (name, desc) in enumerate(items):
        col = i % 3
        row = i // 3
        left = Inches(0.7 + col * 4.15)
        top = Inches(1.4 + row * 1.35)
        rect(s, left, top, Inches(3.9), Inches(1.2), PAPER, line_color=GREEN)
        tb = s.shapes.add_textbox(left + Inches(0.2), top + Inches(0.12), Inches(3.5), Inches(0.4))
        set_run(tb.text_frame.paragraphs[0].add_run(), name, size=13, bold=True, color=GREEN)
        sb = s.shapes.add_textbox(left + Inches(0.2), top + Inches(0.52), Inches(3.5), Inches(0.55))
        set_run(sb.text_frame.paragraphs[0].add_run(), desc, size=11, color=GRAY)


def offer_block(page_no, title, period, scope, price, deliverables, prs, last=False):
    s = blank_slide(prs)
    add_header(s, page_no, title, prs, kind="last" if last else "internal")
    rect(s, Inches(0.7), Inches(1.4), Inches(3.6), Inches(0.7), BLUE,
         text="기간: " + period, size=13)
    rect(s, Inches(4.45), Inches(1.4), Inches(3.6), Inches(0.7), BLUE,
         text="범위: " + scope, size=13)
    rect(s, Inches(8.2), Inches(1.4), Inches(4.4), Inches(0.7), ORANGE,
         text=price + " — 가격 가설", size=13)
    rect(s, Inches(0.7), Inches(2.35), Inches(5.9), Inches(3.6), PAPER, line_color=BLUE)
    dl = s.shapes.add_textbox(Inches(0.95), Inches(2.5), Inches(5.4), Inches(0.4))
    set_run(dl.text_frame.paragraphs[0].add_run(), "산출물", size=15, bold=True, color=INK)
    if len(deliverables) > 6:
        for i, item in enumerate(deliverables):
            col = i // 5
            row = i % 5
            left = Inches(0.95 + col * 2.9)
            top = Inches(3.0 + row * 0.55)
            ib = s.shapes.add_textbox(left, top, Inches(2.7), Inches(0.5))
            set_run(ib.text_frame.paragraphs[0].add_run(), "• " + item, size=12, color=INK)
    else:
        for i, item in enumerate(deliverables):
            top = Inches(3.0 + i * 0.46)
            ib = s.shapes.add_textbox(Inches(0.95), top, Inches(5.4), Inches(0.4))
            set_run(ib.text_frame.paragraphs[0].add_run(), "• " + item, size=12, color=INK)
    rect(s, Inches(6.85), Inches(2.35), Inches(5.8), Inches(3.6), PAPER, line_color=GREEN)
    fb = s.shapes.add_textbox(Inches(7.1), Inches(2.5), Inches(5.3), Inches(0.4))
    set_run(fb.text_frame.paragraphs[0].add_run(), "실행 흐름", size=15, bold=True, color=INK)
    flow = [
        "반복업무 선정", "분석·설계", "합성 실행", "사람 검토·승인", "스킬 카드 납품", "분기 검토·버전 갱신",
    ]
    for i, step in enumerate(flow):
        top = Inches(3.0 + i * 0.46)
        nb = s.shapes.add_textbox(Inches(7.1), top, Inches(0.55), Inches(0.4))
        set_run(nb.text_frame.paragraphs[0].add_run(), "%02d" % (i + 1), size=12, bold=True, color=GREEN)
        sb = s.shapes.add_textbox(Inches(7.7), top, Inches(4.7), Inches(0.4))
        set_run(sb.text_frame.paragraphs[0].add_run(), step, size=12, color=INK)
    note = s.shapes.add_textbox(Inches(0.7), Inches(6.15), Inches(12.0), Inches(0.7))
    set_run(note.text_frame.paragraphs[0].add_run(),
            "모든 실행은 합성 데이터 기반이며, 사람 검토와 승인을 포함합니다. 가격은 시장 검증 전 가설입니다.",
            size=12, color=GRAY)


def pricing_risk():
    s = blank_slide(prs)
    add_header(s, 9, "가격 가설 · 적합 고객 · 위험 경계", prs)
    zones = [
        (0.6, 1.45, RED, "가격 가설 (시장 검증 전)", [
            "Offer A (1~2일): 300만~500만원",
            "Offer B (2~3주): 500만~800만원",
            "Offer C (4~6주): 1,200만~2,000만원",
            "구독·분기 검토: 가설 범위·산정 기준만 안내",
        ]),
        (6.8, 1.45, BLUE, "적합 고객", [
            "반복업무가 주 1회 이상 반복되는 조직",
            "검토·승인자가 존재하는 조직",
            "업무 1개 파일럿부터 시작 가능한 조직",
        ]),
        (0.6, 4.15, ORANGE, "비적합 고객", [
            "전사 자동화를 한 번에 요구하는 경우",
            "사람 검토를 배제하려는 경우",
            "성과(수익) 보장을 요구하는 경우",
        ]),
        (6.8, 4.15, YELLOW, "위험 경계", [
            "정확성 보장·오류 없음·직원 대체 표현 금지",
            "사람 검토와 승인을 항상 포함",
            "파일럿 결과에 따라 확대 여부 결정",
        ]),
    ]
    for left, top, color, title, items in zones:
        rect(s, Inches(left), Inches(top), Inches(5.9), Inches(2.5), PAPER, line_color=color)
        hb = s.shapes.add_textbox(Inches(left + 0.25), Inches(top + 0.15), Inches(5.4), Inches(0.4))
        set_run(hb.text_frame.paragraphs[0].add_run(), title, size=14, bold=True, color=color)
        bullets(s, Inches(left + 0.25), Inches(top + 0.6), Inches(5.4), Inches(1.8), items, size=12)


def next_steps():
    s = blank_slide(prs)
    add_header(s, 10, "다음 단계 — 반복업무 1개 선정", prs, kind="last")
    bullets(s, Inches(0.7), Inches(1.5), Inches(11.5), Inches(3.5), [
        "1단계: Skill Discovery Worksheet으로 반복업무 1개 선정",
        "2단계: 업무 흐름·소요시간·검토 구조만 확인 (실제 파일 불필요)",
        "3단계: Offer A(워크숍) 또는 B(스프린트) 파일럿 진행 여부 검토",
        "4단계: 검증된 스킬 카드 납품 후 분기별 검토·버전 갱신",
    ], size=15)
    rect(s, Inches(0.7), Inches(5.2), Inches(11.9), Inches(1.3), PAPER, line_color=GREEN,
         text="첫 상담에서는 고객의 실제 비공개 파일을 요구하지 않습니다. "
              "업무명·단계·소요시간·검토 구조만 확인합니다.",
         text_color=INK)


cover()
problem()
why_generic()
method()
components()
offer_block(6, "Offer A — AI 업무 스킬 설계 워크숍", "1~2일", "업무 1개", "300만~500만원",
            ["업무 범위", "필요 입력자료", "실행 단계", "AI 보조 단계", "사람 판단 단계",
             "검토 기준", "금지 업무", "예외 목록", "스킬 초안"], prs)
offer_block(7, "Offer B — AI 업무 스킬 전환 스프린트", "2~3주", "업무 1개", "500만~800만원",
            ["업무 분석", "합성 실행", "증거 요구사항", "누락·충돌 처리", "실행자·검토자 역할",
             "수정·재실행", "최종 승인 기준", "검증된 스킬 카드", "운영 가이드"], prs)
offer_block(8, "Offer C — 팀 AI 스킬 라이브러리 파일럿", "4~6주", "업무 3개 · 팀 1개", "1,200만~2,000만원",
            ["검증된 스킬 3개", "역할·승인 기준", "버전 관리 기준", "금지 사용 기준",
             "분기 검토 일정", "팀 운영 플레이북"], prs)
pricing_risk()
next_steps()

assert len(prs.slides._sldIdLst) == 10, "proposal must be exactly 10 slides"
prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides._sldIdLst))
