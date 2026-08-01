"""Build Business32_Verified_Skill_Card_Sample.pptx (2-3 slides)."""
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from _b32theme import *  # noqa: F403
from pptx import Presentation
from pptx.util import Inches

OUT = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
    "Business32_Verified_Skill_Card_Sample.pptx",
)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SAMPLES = [
    "SAMPLE", "SYNTHETIC", "HUMAN REVIEW REQUIRED", "NOT AUTOMATICALLY APPROVED",
]


def sample_ribbon(s):
    rect(s, Inches(0.0), Inches(7.1), Inches(13.333), Inches(0.4), RED)
    rb = s.shapes.add_textbox(Inches(0.3), Inches(7.12), Inches(12.7), Inches(0.35))
    set_run(rb.text_frame.paragraphs[0].add_run(),
            "SAMPLE · SYNTHETIC · HUMAN REVIEW REQUIRED · NOT AUTOMATICALLY APPROVED",
            size=11, bold=True, color=WHITE)


def card_fields(s, title, fields, top=1.35, cols=2):
    tb = s.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.1), Inches(0.4))
    set_run(tb.text_frame.paragraphs[0].add_run(), title, size=16, bold=True, color=GREEN)
    for i, (name, value) in enumerate(fields):
        col = i % cols
        row = i // cols
        left = Inches(0.6 + col * 6.25)
        box_top = Inches(top + 0.45 + row * 0.6)
        nb = s.shapes.add_textbox(left, box_top, Inches(2.2), Inches(0.4))
        set_run(nb.text_frame.paragraphs[0].add_run(), name, size=11, bold=True, color=INK)
        vb = s.shapes.add_textbox(left + Inches(2.35), box_top, Inches(3.8), Inches(0.5))
        set_run(vb.text_frame.paragraphs[0].add_run(), value, size=11, color=BODY_STRONG)


def footer_above_ribbon(s, kind):
    box = s.shapes.add_textbox(Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.3))
    text = {"internal": FOOTER_INTERNAL, "last": FOOTER_LAST}[kind]
    set_run(box.text_frame.paragraphs[0].add_run(), text, size=10, color=GRAY)


def page1():
    s = blank_slide(prs)
    rect(s, Inches(0.0), Inches(0.0), Inches(13.333), Inches(0.14), ORANGE)
    h = s.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.35))
    set_run(h.text_frame.paragraphs[0].add_run(),
            "Business 32 · Verified Skill Card Sample (합성)", size=11, bold=True, color=BLUE)
    t = s.shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(12.3), Inches(0.7))
    set_run(t.text_frame.paragraphs[0].add_run(),
            "교육 프로그램 안내문 작성 및 검토 — VERIFIED ORGANIZATIONAL AI SKILL PACKAGE",
            size=24, bold=True, color=INK)
    sub = s.shapes.add_textbox(Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.3))
    set_run(sub.text_frame.paragraphs[0].add_run(), "업무 목적", size=12, bold=True, color=GREEN)
    card_fields(s, "기본 정보", [
        ("소유자(owner)", "가상 기관명 A · 교육기획팀"),
        ("실행자(active operator)", "홍보 담당"),
        ("검토자(reviewer)", "팀장 · 담당 책임자"),
        ("허용 사용(allowed use)", "합성 콘텐츠·프로그램 안내문 작성"),
        ("금지 사용(prohibited use)", "실제 개인정보 입력·실제 배포 확정 금지"),
        ("입력자료(required inputs)", "가상 프로그램명·일정·대상자·승인자"),
    ], top=1.35)
    card_fields(s, "실행 구조", [
        ("실행 단계(execution steps)", "입력 확인 → 초안 → 증거 연결 → 검토 → 승인 → 배포안"),
        ("AI 보조 단계(AI-assisted steps)", "안내문 초안 작성 · 문구 정리"),
        ("사람 판단 단계(human actions)", "사실 확인 · 일정·대상자 확정 · 검토 · 승인"),
        ("필수 증거(evidence requirements)", "가상 일정·대상자 근거 연결"),
        ("검토 기준(review checks)", "사실 일치 · 대상자 적합 · 금지 사용 없음"),
        ("승인 기준(approval record)", "담당 책임자 최종 승인"),
    ], top=4.2)
    sample_ribbon(s)
    footer_above_ribbon(s, "internal")


def page2():
    s = blank_slide(prs)
    h = s.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.3), Inches(0.35))
    set_run(h.text_frame.paragraphs[0].add_run(),
            "Business 32 · Verified Skill Card Sample (합성) — 2/2", size=11, bold=True, color=BLUE)
    card_fields(s, "처리·경계", [
        ("누락 처리(missing-evidence)", "보완 요청 또는 중단 · 자동 추정 금지"),
        ("충돌 처리(conflicting-evidence)", "사람 판단 · 자동 최선 금지"),
        ("예외(known exceptions)", "긴급 변경 시 규칙 재검토"),
        ("재실행 조건(rollback condition)", "승인 전 수정 시 절차 재실행 · 이전 버전 복귀 가능"),
        ("출력 형식", "안내문 원고(문서·이메일·SNS 초안)"),
        ("버전(version)", "1.0 (합성)"),
    ], top=1.35)
    rb = s.shapes.add_textbox(Inches(0.6), Inches(3.7), Inches(12.1), Inches(0.4))
    set_run(rb.text_frame.paragraphs[0].add_run(), "다음 검토일(next review date): 2026-11-01 (합성)", size=14, bold=True, color=ORANGE)
    box = s.shapes.add_textbox(Inches(0.6), Inches(4.4), Inches(12.1), Inches(1.7))
    tf = box.text_frame
    tf.word_wrap = True
    set_run(tf.paragraphs[0].add_run(),
            "이 카드는 SAMPLE·SYNTHETIC이며, 실제 조직 스킬이 아닙니다. "
            "사람 검토와 승인 없이는 어떤 결과도 확정되지 않습니다. "
            "AI가 자동 승인하지 않습니다.", size=13, color=INK)
    sample_ribbon(s)
    footer_above_ribbon(s, "last")


page1()
page2()

assert 2 <= len(prs.slides._sldIdLst) <= 3, "skill card must be 2-3 slides"
prs.save(OUT)
print("saved", OUT, "slides:", len(prs.slides._sldIdLst))
