"""Shared blueprint-workshop theme for Business 32 customer-package generators."""
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

PAPER = RGBColor(0xF5, 0xF1, 0xE7)
INK = RGBColor(0x1F, 0x29, 0x33)
BLUE = RGBColor(0x2B, 0x6C, 0xB0)
ORANGE = RGBColor(0xC0, 0x56, 0x21)
GREEN = RGBColor(0x2F, 0x85, 0x5A)
RED = RGBColor(0xC5, 0x30, 0x30)
YELLOW = RGBColor(0xB7, 0x79, 0x1F)
GRAY = RGBColor(0x6B, 0x72, 0x80)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "NanumSquareRound"

FOOTER_COVER = "DRAFT · 제공자 정보 최종 확정 필요"
FOOTER_INTERNAL = "DRAFT"
FOOTER_LAST = "가격 가설 · 제공자 정보 최종 확정 필요"


def set_run(run, text, size=14, bold=False, color=INK, font=FONT):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def blank_slide(prs, width=13.333, height=7.5):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PAPER
    bg.line.fill.background()
    bg.shadow.inherit = False
    return slide


def add_footer(slide, kind, prs):
    box = slide.shapes.add_textbox(Inches(0.4), Inches(7.08), Inches(12.5), Inches(0.35))
    tf = box.text_frame
    tf.word_wrap = True
    text = {"cover": FOOTER_COVER, "internal": FOOTER_INTERNAL, "last": FOOTER_LAST}[kind]
    set_run(tf.paragraphs[0].add_run(), text, size=10, color=GRAY)


def add_header(slide, page_no, title, prs, kind="internal"):
    box = slide.shapes.add_textbox(Inches(0.4), Inches(0.18), Inches(11.5), Inches(0.3))
    set_run(box.text_frame.paragraphs[0].add_run(), "Business 32 · AI Skill Studio · AI 업무 스킬 전환 프로그램", size=10, color=GRAY)
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.45), Inches(12.5), Inches(0.75))
    set_run(title_box.text_frame.paragraphs[0].add_run(), title, size=30, bold=True, color=INK)
    page_box = slide.shapes.add_textbox(Inches(12.4), Inches(0.2), Inches(0.7), Inches(0.3))
    p = page_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    set_run(p.add_run(), str(page_no), size=10, color=GRAY)
    add_footer(slide, kind, prs)


def rect(slide, left, top, width, height, color, line_color=None, text="", size=12, bold=True, text_color=WHITE):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.5)
    shape.shadow.inherit = False
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.12)
        tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.06)
        tf.margin_bottom = Inches(0.06)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        set_run(tf.paragraphs[0].add_run(), text, size=size, bold=bold, color=text_color)
    return shape


def tag(slide, left, top, text, color=BLUE, size=10):
    w = Inches(0.15 + 0.11 * len(text))
    box = slide.shapes.add_textbox(left, top, w, Inches(0.3))
    p = box.text_frame.paragraphs[0]
    set_run(p.add_run(), text, size=size, bold=True, color=color)
    return box


def bullets(slide, left, top, width, height, items, size=14, color=INK):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        set_run(p.add_run(), item, size=size, color=color)
    return box
