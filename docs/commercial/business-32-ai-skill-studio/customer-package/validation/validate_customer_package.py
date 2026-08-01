#!/usr/bin/env python3
"""Validate the Business 32 customer-facing pilot package.

Checks required files, page counts (proposal 10, one-page 1, worksheet <=3,
skill card 2-3), rendered PNG parity, external runtime 0, no real customer/org
data, no backend/SaaS/auto-approval claims, price-hypothesis presence, human
review wording, and that only the customer-package scope changed.
"""
import glob
import math
import os
import re
import subprocess
import sys

from pptx import Presentation
from docx import Document
from openpyxl import load_workbook
from pypdf import PdfReader

VALIDATED_UX_HEAD = "73ec4718d0835248ab20d56bc68f3956536112b4"
VALIDATED_HANDOFF_HEAD = "29068281998b7f1a59d76a95174807ffbf20cb38"
COMMERCIAL_HEAD = "30565f4ddcf99296751109df3a0973d7ba79eaa8"

PACKAGE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))  # customer-package/
REPO_ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.dirname(PACKAGE))))
PRODUCT_WORKSPACE = os.path.join(REPO_ROOT, "reference", "business-32-ai-skill-studio-ux")

REQUIRED_FILES = [
    "README.md",
    "SOURCE_MAPPING.md",
    "CUSTOMIZATION_CHECKLIST.md",
    "VISUAL_QA.md",
    "Business32_Master_Proposal_10p.pptx",
    "Business32_Master_Proposal_10p.pdf",
    "Business32_OnePage_Offer_Source.pptx",
    "Business32_OnePage_Offer.pdf",
    "Business32_Skill_Discovery_Worksheet.docx",
    "Business32_Skill_Discovery_Worksheet.pdf",
    "Business32_Verified_Skill_Card_Sample.pptx",
    "Business32_Verified_Skill_Card_Sample.pdf",
    "Business32_Pilot_Quote_Template.xlsx",
    "Business32_Customer_Meeting_Script.md",
    "Business32_Followup_Email_Templates.md",
    "rendered/manifest.md",
    "validation/_b32theme.py",
    "validation/build_proposal_pptx.py",
    "validation/build_one_page_pptx.py",
    "validation/build_discovery_worksheet.py",
    "validation/build_skill_card_pptx.py",
    "validation/build_quote_xlsx.py",
    "validation/validate_customer_package.py",
]

PDFS = {
    "Business32_Master_Proposal_10p.pdf": 10,
    "Business32_OnePage_Offer.pdf": 1,
    "Business32_Skill_Discovery_Worksheet.pdf": 3,
    "Business32_Verified_Skill_Card_Sample.pdf": 3,
}

PPTX_SLIDES = {
    "Business32_Master_Proposal_10p.pptx": 10,
    "Business32_OnePage_Offer_Source.pptx": 1,
    "Business32_Verified_Skill_Card_Sample.pptx": 3,
}

XLSX_SHEETS = [
    "안내", "고객·업무", "Offer 선택", "작업 범위", "산출물", "일정", "견적", "가정·제외사항",
]

PRICES = ["300만~500만원", "500만~800만원", "1,200만~2,000만원"]

failures = []


def path(rel):
    return os.path.join(PACKAGE, rel)


def read_md(rel):
    with open(path(rel), "r", encoding="utf-8") as fh:
        return fh.read()


def check(name, fn):
    try:
        fn()
        print("PASS " + name)
    except AssertionError as error:
        failures.append(name)
        print("FAIL " + name + ": " + str(error))


def pdf_text(rel):
    reader = PdfReader(path(rel))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def normalize(text):
    return re.sub(r"\s+", "", text)


def main():
    for rel in REQUIRED_FILES:
        check("required file exists: " + rel, lambda r=rel: assert_exists(r))

    for rel, max_pages in PDFS.items():
        check("pdf page count within bound: " + rel, lambda r=rel, m=max_pages: assert_pdf_pages(r, m))

    check("proposal pdf is exactly 10 pages", lambda: assert_pdf_exact("Business32_Master_Proposal_10p.pdf", 10))
    check("one-page pdf is exactly 1 page", lambda: assert_pdf_exact("Business32_OnePage_Offer.pdf", 1))
    check("worksheet pdf is 3 pages or fewer", lambda: assert_pdf_exact("Business32_Skill_Discovery_Worksheet.pdf", None, max_pages=3))
    check("skill card pdf is 2-3 pages", lambda: assert_pdf_exact("Business32_Verified_Skill_Card_Sample.pdf", None, min_pages=2, max_pages=3))

    check("proposal pptx is exactly 10 slides", lambda: assert_slides("Business32_Master_Proposal_10p.pptx", 10))
    check("one-page pptx is exactly 1 slide", lambda: assert_slides("Business32_OnePage_Offer_Source.pptx", 1))
    check("skill card pptx is 2-3 slides", lambda: assert_slides("Business32_Verified_Skill_Card_Sample.pptx", None, min_pages=2, max_pages=3))

    check("discovery worksheet has 13 questions", lambda: assert_worksheet_questions())
    check("worksheet checkbox only on check-type question", lambda: assert_worksheet_checkboxes())
    check("worksheet is exactly 2 pages with repeated header", lambda: assert_worksheet_pages())
    check("worksheet splits Q1-7/Q8-13 across pages", lambda: assert_worksheet_split())

    check("pptx shapes stay within page bounds", lambda: assert_pptx_bounds())
    check("footer does not overlap content", lambda: assert_footer_overlap())
    check("pptx text fits inside every shape", lambda: assert_text_fit())
    check("offer deliverable text stays inside panel", lambda: assert_offer_deliverables_fit())
    check("offer disclaimer does not overlap panel", lambda: assert_offer_disclaimer_clear())

    check("quote workbook has all 8 sheets", lambda: assert_xlsx_sheets())

    check("rendered PNGs match every pdf page", lambda: assert_png_parity())

    check("external runtime dependency 0 in docs", lambda: assert_no_external_runtime())

    md_all = "\n".join(read_md(rel) for rel in REQUIRED_FILES if rel.endswith(".md"))
    check("no real customer/org data in docs", lambda: (
        assert_no_email(md_all),
        assert_no_phone(md_all),
        assert_no_business_number(md_all),
    ))

    for rel, _max in PDFS.items():
        text = pdf_text(rel)
        check("no PII in pdf: " + rel, lambda t=text: (
            assert_no_email(t),
            assert_no_phone(t),
            assert_no_business_number(t),
        ))
        check("no backend/SaaS/auto-approval claims in pdf: " + rel, lambda t=text: (
            assert_not_contain(t, "정확성을 보장"),
            assert_not_contain(t, "보장합니다"),
            assert_not_contain(t, "직원을 대체"),
            assert_not_contain(t, "대체합니다"),
            assert_not_contain(t, "자동 승인합니다"),
        ))
        check("human review wording present in pdf: " + rel, lambda t=text: (
            assert_contain(t, "사람 검토"),
            assert_contain(t, "합성"),
        ))

    proposal = pdf_text("Business32_Master_Proposal_10p.pdf")
    check("all offer prices shown as hypotheses in proposal", lambda: (
        tuple(assert_contain(normalize(proposal), normalize(price)) for price in PRICES),
        assert_contain(proposal, "가격 가설"),
        assert_contain(proposal, "가설"),
    ))

    one = pdf_text("Business32_OnePage_Offer.pdf")
    check("one-page shows price hypothesis and human review", lambda: (
        assert_contain(one, "가설"),
        assert_contain(one, "사람 검토"),
    ))

    check("validated heads referenced", lambda: (
        assert_contain(read_md("README.md"), VALIDATED_UX_HEAD),
        assert_contain(read_md("README.md"), VALIDATED_HANDOFF_HEAD),
        assert_contain(read_md("README.md"), COMMERCIAL_HEAD),
    ))

    check("pixel visual QA not declared", lambda: (
        assert_contain(read_md("VISUAL_QA.md"), "PIXEL_VISUAL_QA_PASS: NOT DECLARED"),
        assert_contain(read_md("VISUAL_QA.md"), "CUSTOMER_SEND_READY:  NOT DECLARED"),
    ))

    check("no backend implementation files in workspace", lambda: assert_no_backend_files())

    check("validated product workspace unchanged", lambda: assert_product_unchanged())

    check("only customer-package scope changed", lambda: assert_scope())

    print()
    if failures:
        print("%d validation failure(s)" % len(failures))
        return 1
    print("customer package validation ok")
    return 0


def assert_pdf_pages(rel, max_pages):
    n = len(PdfReader(path(rel)).pages)
    assert n <= max_pages, "pdf pages %d > %d: %s" % (n, max_pages, rel)


def assert_exists(rel):
    assert os.path.isfile(path(rel)), "missing file: " + rel


def assert_pdf_exact(rel, exact=None, min_pages=None, max_pages=None):
    n = len(PdfReader(path(rel)).pages)
    if exact is not None:
        assert n == exact, "pdf pages %d != %d: %s" % (n, exact, rel)
    if min_pages is not None:
        assert n >= min_pages, "pdf pages %d < %d: %s" % (n, min_pages, rel)
    if max_pages is not None:
        assert n <= max_pages, "pdf pages %d > %d: %s" % (n, max_pages, rel)


def assert_slides(rel, exact=None, min_pages=None, max_pages=None):
    n = len(Presentation(path(rel)).slides._sldIdLst)
    if exact is not None:
        assert n == exact, "slides %d != %d: %s" % (n, exact, rel)
    if min_pages is not None:
        assert n >= min_pages, "slides %d < %d: %s" % (n, min_pages, rel)
    if max_pages is not None:
        assert n <= max_pages, "slides %d > %d: %s" % (n, max_pages, rel)


def assert_worksheet_questions():
    doc = Document(path("Business32_Skill_Discovery_Worksheet.docx"))
    count = 0
    for p in doc.paragraphs:
        if re.match(r"^(☐ )?Q\d+\.", p.text):
            count += 1
    assert count == 13, "expected 13 questions, got %d" % count


def assert_xlsx_sheets():
    wb = load_workbook(path("Business32_Pilot_Quote_Template.xlsx"), data_only=False)
    assert wb.sheetnames == XLSX_SHEETS, "unexpected sheets: %s" % wb.sheetnames


def assert_png_parity():
    for rel, _max in PDFS.items():
        base = rel[:-4]
        n = len(PdfReader(path(rel)).pages)
        pngs = glob.glob(os.path.join(PACKAGE, "rendered", base + "_p-*.png"))
        assert len(pngs) == n, "png count %d != pdf pages %d for %s" % (len(pngs), n, base)


def assert_no_external_runtime():
    for rel in REQUIRED_FILES:
        if not rel.endswith(".md"):
            continue
        content = read_md(rel)
        stripped = content.replace('xmlns="http://www.w3.org/2000/svg"', "")
        assert "http://" not in stripped and "https://" not in stripped, "external URL in " + rel


def assert_no_email(text):
    assert not re.search(r"[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}", text), "email-like pattern found"


def assert_no_phone(text):
    assert not re.search(r"010[- ]?\d{3,4}[- ]?\d{4}", text), "phone pattern found"


def assert_no_business_number(text):
    assert not re.search(r"\d{3}-\d{2}-\d{5}", text), "business registration pattern found"


def assert_contain(text, needle):
    assert needle in text, "missing: " + needle


def assert_not_contain(text, needle):
    assert needle not in text, "forbidden claim present: " + needle


def assert_no_backend_files():
    backend_markers = [".sql", "schema.py", "drizzle", "migration", "api/", "worker.py", "server.py"]
    for root, _dirs, files in os.walk(PACKAGE):
        for name in files:
            full = os.path.join(root, name)
            rel = os.path.relpath(full, PACKAGE)
            if rel.startswith("validation") and name in ("validate_customer_package.py", "build_proposal_pptx.py"):
                continue
            for marker in backend_markers:
                assert marker not in rel.lower(), "backend file marker in package: " + rel


def assert_product_unchanged():
    out = subprocess.run(
        ["git", "status", "--porcelain", "--untracked-files=all", "--", PRODUCT_WORKSPACE],
        capture_output=True, text=True, cwd=REPO_ROOT,
    ).stdout.strip()
    assert not out, "validated product workspace has changes:\n" + out


def assert_scope():
    base = "origin/docs/business-32-commercial-package"
    out = subprocess.run(
        ["git", "diff", "--name-only", base + "...HEAD"],
        capture_output=True, text=True, cwd=REPO_ROOT,
    ).stdout.strip()
    allowed = "docs/commercial/business-32-ai-skill-studio/customer-package/"
    if not out:
        return
    for line in out.splitlines():
        assert line.startswith(allowed), "out-of-scope path: " + line


def assert_worksheet_checkboxes():
    doc = Document(path("Business32_Skill_Discovery_Worksheet.docx"))
    descriptive = 0
    checkable = 0
    for p in doc.paragraphs:
        m = re.match(r"^(☐ )?Q(\d+)\.", p.text)
        if not m:
            continue
        if m.group(1):
            checkable += 1
            assert m.group(2) == "12", "checkbox on unexpected question Q%s" % m.group(2)
        else:
            descriptive += 1
    assert checkable == 1, "expected 1 checkbox question, got %d" % checkable
    assert descriptive == 12, "expected 12 descriptive questions, got %d" % descriptive


def assert_worksheet_pages():
    reader = PdfReader(path("Business32_Skill_Discovery_Worksheet.pdf"))
    assert len(reader.pages) == 2, "worksheet pdf pages %d != 2" % len(reader.pages)
    p1 = reader.pages[0].extract_text() or ""
    p2 = reader.pages[1].extract_text() or ""
    n1, n2 = normalize(p1), normalize(p2)
    assert "Business32·AISkillStudio" in n1, "worksheet page 1 missing header"
    assert "Business32·AISkillStudio" in n2, "worksheet page 2 missing repeated header"
    assert "SkillDiscoveryWorksheet" in n2, "worksheet page 2 missing repeated title"
    assert "페이지1/2" in n1, "worksheet page 1 missing page number"
    assert "페이지2/2" in n2, "worksheet page 2 missing page number"
    assert len(n2.strip()) > 300, "worksheet page 2 is excessively blank"


def assert_worksheet_split():
    reader = PdfReader(path("Business32_Skill_Discovery_Worksheet.pdf"))
    p1 = reader.pages[0].extract_text() or ""
    p2 = reader.pages[1].extract_text() or ""
    for i in range(1, 8):
        assert "Q%d." % i in p1, "Q%d missing on page 1" % i
        assert "Q%d." % i not in p2, "Q%d unexpectedly on page 2" % i
    for i in range(8, 14):
        assert "Q%d." % i in p2, "Q%d missing on page 2" % i
        assert "Q%d." % i not in p1, "Q%d unexpectedly on page 1" % i


EMU_PER_INCH = 914400
FOOTER_KEYWORDS = ("DRAFT", "제공자 정보")
FOOTER_REGION_IN = 6.5
PPTX_FILES = [
    "Business32_Master_Proposal_10p.pptx",
    "Business32_OnePage_Offer_Source.pptx",
    "Business32_Verified_Skill_Card_Sample.pptx",
]


def _shape_text(shape):
    try:
        return shape.text_frame.text or ""
    except Exception:
        return ""


def _is_footer(shape):
    if shape.top is None or shape.height is None:
        return False
    top_in = shape.top / EMU_PER_INCH
    bottom_in = (shape.top + shape.height) / EMU_PER_INCH
    if not (top_in >= FOOTER_REGION_IN or bottom_in >= 7.0):
        return False
    text = _shape_text(shape)
    return any(k in text for k in FOOTER_KEYWORDS)


def _is_edge_bar(shape, sw, sh):
    if None in (shape.left, shape.top, shape.width, shape.height):
        return False
    l_in = shape.left / EMU_PER_INCH
    w_in = shape.width / EMU_PER_INCH
    t_in = shape.top / EMU_PER_INCH
    h_in = shape.height / EMU_PER_INCH
    sw_in = sw / EMU_PER_INCH
    sh_in = sh / EMU_PER_INCH
    if w_in < sw_in - 0.05:
        return False
    return t_in <= 0.05 or (t_in + h_in) >= sh_in - 0.05


def _shape_box(shape):
    return (shape.left, shape.top, shape.left + shape.width, shape.top + shape.height)


def _rects_overlap(a, b):
    return a[0] < b[2] and b[0] < a[2] and a[1] < b[3] and b[1] < a[3]


def assert_pptx_bounds():
    tol = 1000
    for rel in PPTX_FILES:
        prs = Presentation(path(rel))
        sw, sh = prs.slide_width, prs.slide_height
        for idx, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                vals = (shape.left, shape.top, shape.width, shape.height)
                if any(v is None for v in vals):
                    continue
                left, top, width, height = vals
                assert left >= -tol, "%s slide %d shape left %d < 0" % (rel, idx, left)
                assert top >= -tol, "%s slide %d shape top %d < 0" % (rel, idx, top)
                assert left + width <= sw + tol, (
                    "%s slide %d shape right overflow %d > %d" % (rel, idx, left + width, sw))
                assert top + height <= sh + tol, (
                    "%s slide %d shape bottom overflow %d > %d" % (rel, idx, top + height, sh))


def assert_footer_overlap():
    for rel in PPTX_FILES:
        prs = Presentation(path(rel))
        sw, sh = prs.slide_width, prs.slide_height
        for idx, slide in enumerate(prs.slides, start=1):
            shapes = list(slide.shapes)
            for f in [s for s in shapes if _is_footer(s)]:
                fb = _shape_box(f)
                for other in shapes:
                    if other is f or _is_footer(other) or _is_edge_bar(other, sw, sh):
                        continue
                    if other.width == sw and other.height == sh:
                        continue
                    if _rects_overlap(fb, _shape_box(other)):
                        raise AssertionError(
                            "%s slide %d footer overlaps shape: %s" % (rel, idx, _shape_text(other)[:40]))


def _char_width(ch):
    if ord(ch) >= 0x2E80:
        return 1.0
    if ch in "‘’'\"()[]{}·":
        return 0.3
    return 0.55


def _estimate_text_height(text, usable_width_in, font_pt):
    per_line = max(usable_width_in * 72.0 / font_pt, 0.1)
    lines = 0
    for para in text.split("\n"):
        chars = sum(_char_width(c) for c in para)
        lines += max(1, math.ceil(chars / per_line))
    return lines * font_pt / 72.0 * 1.3


def _text_fits(shape):
    tf = shape.text_frame
    text = tf.text or ""
    if not text.strip() or shape.width is None or shape.height is None:
        return True
    width_in = shape.width / EMU_PER_INCH
    height_in = shape.height / EMU_PER_INCH
    usable_w = width_in - tf.margin_left / EMU_PER_INCH - tf.margin_right / EMU_PER_INCH
    usable_h = height_in - tf.margin_top / EMU_PER_INCH - tf.margin_bottom / EMU_PER_INCH
    if usable_h <= 0:
        return True
    est = 0.0
    for p in tf.paragraphs:
        if not p.runs:
            continue
        font_pt = max((r.font.size.pt if r.font.size else 12) for r in p.runs)
        if tf.word_wrap and usable_w > 0.05:
            est += _estimate_text_height(p.text, usable_w, font_pt)
        else:
            est += font_pt / 72.0 * 1.3
    return est <= usable_h + 0.03


def assert_text_fit():
    for rel in PPTX_FILES:
        prs = Presentation(path(rel))
        for idx, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                assert _text_fits(shape), "%s slide %d text overflows shape: %s" % (
                    rel, idx, _shape_text(shape)[:40])


def _offer_panel_bounds():
    return tuple(v * EMU_PER_INCH for v in (0.7, 2.35, 6.6, 5.95))


def assert_offer_deliverables_fit():
    prs = Presentation(path("Business32_Master_Proposal_10p.pptx"))
    p_left, p_top, p_right, p_bottom = _offer_panel_bounds()
    tol = 0.05 * EMU_PER_INCH
    for slide_idx in (6, 7):
        slide = list(prs.slides)[slide_idx - 1]
        items = []
        for shape in slide.shapes:
            if shape.top is None or shape.left is None or shape.width is None:
                continue
            if not (p_left - tol <= shape.left < p_right
                    and p_top - tol <= shape.top < p_bottom):
                continue
            if not _shape_text(shape).startswith("• "):
                continue
            items.append((shape, _shape_text(shape)))
        assert len(items) == 9, "slide %d deliverable items %d != 9" % (slide_idx, len(items))
        for shape, text in items:
            assert shape.top + shape.height <= p_bottom + 0.02 * EMU_PER_INCH, (
                "slide %d deliverable '%s' bottom %.3f exceeds panel" % (
                    slide_idx, text, (shape.top + shape.height) / EMU_PER_INCH))
            assert _text_fits(shape), "slide %d deliverable text overflows shape: %s" % (slide_idx, text)


def assert_offer_disclaimer_clear():
    prs = Presentation(path("Business32_Master_Proposal_10p.pptx"))
    _p_left, _p_top, _p_right, p_bottom = _offer_panel_bounds()
    for slide_idx in (6, 7, 8):
        slide = list(prs.slides)[slide_idx - 1]
        for shape in slide.shapes:
            if "모든 실행은 합성 데이터 기반" not in _shape_text(shape):
                continue
            assert shape.top >= p_bottom + 0.05 * EMU_PER_INCH, (
                "slide %d disclaimer top %.3f overlaps panel bottom %.3f" % (
                    slide_idx, shape.top / EMU_PER_INCH, p_bottom / EMU_PER_INCH))


if __name__ == "__main__":
    sys.exit(main())
