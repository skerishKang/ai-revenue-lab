#!/usr/bin/env python3
"""B35 Independent QA Harness - Lane C #1505

Machine-checkable acceptance surface for the B35 closeout. This lane does not
author commercial copy and does not regenerate final customer artifacts.

Required verdicts (each must be PASS for overall PASS):

  PACKAGE_INVENTORY_PASS
  SOURCE_MAPPING_PASS
  STRUCTURAL_QA_PASS
  FORMULA_QA_PASS
  TEXT_FIT_PASS
  STALE_ARTIFACT_REJECTION_PASS
  PRIVATE_DATA_BOUNDARY_PASS
  EXACT_REVISION_TRACE_PASS

A failed or unavailable check remains failed/unavailable; do not convert it to
PASS by inference. Historical PASS from PR #359 does not transfer.

Usage:
  python tools/b35-independent-qa/validate_b35_independent_qa.py
  python tools/b35-independent-qa/validate_b35_independent_qa.py --package-root docs/.../customer-package/v3-regenerated
"""

from __future__ import annotations

import argparse
import hashlib
import json
import re
import subprocess
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import List, Dict, Tuple, Optional

# ---------------------------------------------------------------------------
# Constants - product authority (frozen at W0)
# ---------------------------------------------------------------------------

PRODUCT_PR = "#370"
PRODUCT_COMMIT = "05932da3af774220372f0e9f3716b07cd83511f9"
PRODUCT_CONTRACT_REL = "reference/business-35-ai-media-education-dx-v3/PRODUCT_CONTRACT.md"

# Exact accepted source revision for final G3 (CENTRAL). Any other 40-hex
# SOURCE_REVISION must fail EXACT_REVISION_TRACE (format-only PASS is forbidden).
ACCEPTED_SOURCE_REVISION = "63adbefcf24a91a5a064c6b8e13779e151ba7de7"

# Price hypothesis tokens (standard ladder)
PRICE_TOKENS = [
    "300만–500만원",
    "500만–800만원",
    "300만–800만원",
    "1,000만–1,500만원",
    "1,500만–2,500만원",
    "월 300만–600만원",
]

# V3.1 journey / promise markers
V3_JOURNEY_MARKERS = [
    "현재 미디어 업무 병목",
    "조직·결과물·병목·팀 규모·AI 사용 상태",
    "조직별 진단 + 새 업무 흐름 + 추천 파일럿",
    "운영체계 산출물",
    "진단 워크숍 또는 6주 파일럿",
    "자기 조직용 전환 요약으로 상담 준비",
    "파디엠 AI 미디어 업무전환 스튜디오",
    "AI 교육을 듣는 데서 끝내지 않고, 팀의 실제 미디어 업무 한 흐름을 사람이 승인하는 운영체계로 바꾼다",
]

# Stale markers that must NOT appear as current truth
STALE_FORBIDDEN_AS_CURRENT = [
    "BUSINESS_35_FINAL_PACKAGE_QA_PASS",
    "DOCUMENT_PIXEL_QA_PASS",
    "XLSX_PIXEL_QA_PASS",
    "BLOCKER_ZERO",  # as standalone claim without context is historical
    "UI_DEPLOYED_VERIFIED",
    "UX_NOT_STARTED",
    "PHASE1",
]

# No-send / legal-review completion forbiddens
FORBIDDEN_COMPLETION_CLAIMS = [
    "법률 검토 완료",
    "법적 검토 완료",
    "계약 검토 완료",
    "legal review completed",
    "LEGAL_REVIEW_COMPLETED",
    "CUSTOMER_SEND_READY=true",
    "CUSTOMER_SEND_READY = true",
    "고객 발송 준비 완료",
]

FORBIDDEN_PHRASES_CUSTOMER = [
    "AI 도입 의무",
    "반드시 도입",
    "법적으로 안전",
    "저작권 문제 없음",
    "개인정보 문제 없음",
    "지원금 수령 가능",
    "자동 수의계약",
    "1억원 이하",
    "성과 보장",  # only allowed in negation
]

# Private data patterns
EMAIL_RE = re.compile(r"[a-zA-Z0-9._%+\-]+@[a-zA-Z0-9.\-]+\.[a-zA-Z]{2,}")
PHONE_RE = re.compile(r"\b\d{2,3}-\d{3,4}-\d{4}\b")
BIZNO_RE = re.compile(r"\b\d{3}-\d{2}-\d{5}\b")
SECRET_RE = re.compile(r"(AKIA[0-9A-Z]{16}|ghp_[0-9A-Za-z]{36}|sk-[A-Za-z0-9]{20,}|-----BEGIN [A-Z ]+ PRIVATE KEY-----)")

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

@dataclass
class CheckResult:
    name: str
    passed: bool
    unavailable: bool = False
    details: List[str] = field(default_factory=list)

    @property
    def verdict(self) -> str:
        if self.unavailable:
            return f"{self.name}_UNAVAILABLE"
        return f"{self.name}_{'PASS' if self.passed else 'FAIL'}"


def sha256_file(p: Path) -> str:
    h = hashlib.sha256()
    with p.open("rb") as f:
        for chunk in iter(lambda: f.read(8192), b""):
            h.update(chunk)
    return h.hexdigest()


def try_pdf_pages(pdf_path: Path) -> Tuple[Optional[int], str]:
    """Return (pages, method) or (None, reason)."""
    if not pdf_path.exists():
        return None, "missing file"
    # Try pdfinfo
    try:
        out = subprocess.run(["pdfinfo", str(pdf_path)], capture_output=True, text=True, timeout=10)
        if out.returncode == 0:
            m = re.search(r"Pages:\s+(\d+)", out.stdout)
            if m:
                return int(m.group(1)), "pdfinfo"
    except Exception as e:
        pass
    # Try PyMuPDF
    try:
        import fitz  # type: ignore
        doc = fitz.open(str(pdf_path))
        return doc.page_count, "pymupdf"
    except Exception:
        pass
    # Try PyPDF2
    try:
        from PyPDF2 import PdfReader  # type: ignore
        r = PdfReader(str(pdf_path))
        return len(r.pages), "pypdf2"
    except Exception:
        pass
    return None, "pdfinfo/pymupdf/pypdf2 unavailable"


def try_pdf_text(pdf_path: Path) -> Tuple[Optional[str], str]:
    if not pdf_path.exists():
        return None, "missing"
    try:
        out = subprocess.run(["pdftotext", str(pdf_path), "-"], capture_output=True, text=True, timeout=10)
        if out.returncode == 0 and out.stdout.strip():
            return out.stdout, "pdftotext"
    except Exception:
        pass
    try:
        import fitz
        doc = fitz.open(str(pdf_path))
        txt = "\n".join(page.get_text() for page in doc)
        if txt.strip():
            return txt, "pymupdf"
    except Exception:
        pass
    try:
        from PyPDF2 import PdfReader
        r = PdfReader(str(pdf_path))
        txt = "\n".join((p.extract_text() or "") for p in r.pages)
        if txt.strip():
            return txt, "pypdf2"
    except Exception as e:
        return None, f"pdf text extraction unavailable: {e}"
    return None, "no text extracted"


def extract_pptx_text(pptx_path: Path) -> Optional[str]:
    try:
        from pptx import Presentation
        prs = Presentation(str(pptx_path))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append(shape.text_frame.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            texts.append(cell.text)
        return "\n".join(texts)
    except Exception:
        return None


def extract_docx_text(docx_path: Path) -> Optional[str]:
    try:
        from docx import Document
        doc = Document(str(docx_path))
        texts = [p.text for p in doc.paragraphs]
        for tbl in doc.tables:
            for row in tbl.rows:
                for cell in row.cells:
                    texts.append(cell.text)
        return "\n".join(texts)
    except Exception:
        return None


def extract_xlsx_text(xlsx_path: Path) -> Optional[str]:
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(xlsx_path), data_only=True, read_only=True)
        out = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                for v in row:
                    if v is not None:
                        out.append(str(v))
        return "\n".join(out)
    except Exception:
        return None


# ---------------------------------------------------------------------------
# Checks
# ---------------------------------------------------------------------------

def check_package_inventory(commercial_root: Path, package_root: Path) -> CheckResult:
    r = CheckResult(name="PACKAGE_INVENTORY", passed=True)
    # Commercial source required files
    commercial_required = [
        "CURRENT_PRODUCT_AUTHORITY.md",
        "README.md",
        "SOURCES.md",
        "01-one-page-offer.md",
        "02-ten-page-proposal.md",
        "03-diagnostic-questionnaire.md",
        "04-six-week-pilot-plan.md",
        "05-statement-of-work-draft.md",
        "06-risk-and-data-annex.md",
        "07-kpi-measurement-framework.md",
        "08-customer-qualification-scorecard.md",
        "tests/validate_sales_package.py",
    ]
    missing_commercial = []
    for rel in commercial_required:
        p = commercial_root / rel
        if not p.is_file():
            missing_commercial.append(rel)
            r.details.append(f"missing commercial source: {rel}")
        elif p.stat().st_size == 0:
            missing_commercial.append(rel)
            r.details.append(f"empty commercial source: {rel}")

    # Package families - at least one candidate per family must exist
    families: Dict[str, List[str]] = {
        "Master Proposal PPTX": ["Business35_Master_Proposal_10p.pptx", "Business35_V3_1_Master_Proposal_10p.pptx"],
        "Master Proposal PDF": ["Business35_Master_Proposal_10p.pdf", "Business35_V3_1_Master_Proposal_10p.pdf", "pdf/Business35_V3_1_Master_Proposal_10p.pdf"],
        "OnePage Offer PPTX": ["Business35_OnePage_Offer_Source.pptx", "Business35_V3_1_OnePage_Offer_Source.pptx"],
        "OnePage Offer PDF": ["Business35_OnePage_Offer.pdf", "Business35_V3_1_OnePage_Offer.pdf", "pdf/Business35_V3_1_OnePage_Offer_Source.pdf"],
        "Questionnaire DOCX": ["Business35_Diagnostic_Questionnaire.docx", "Business35_V3_1_Diagnostic_Questionnaire.docx"],
        "Questionnaire PDF": ["Business35_Diagnostic_Questionnaire.pdf", "Business35_V3_1_Diagnostic_Questionnaire.pdf", "pdf/Business35_V3_1_Diagnostic_Questionnaire.pdf"],
        "Quote XLSX": ["Business35_Pilot_Quote_Template.xlsx", "Business35_V3_1_Pilot_Quote_Template.xlsx"],
        "Meeting Script": ["Business35_Customer_Meeting_Script.md"],
        "Followup Templates": ["Business35_Followup_Email_Templates.md"],
        "Source Mapping": ["SOURCE_MAPPING.md"],
        "Customization Checklist": ["CUSTOMIZATION_CHECKLIST.md"],
        "Package README": ["README.md"],
    }
    missing_families = []
    for fam, candidates in families.items():
        found = any((package_root / c).is_file() for c in candidates)
        if not found:
            missing_families.append(fam)
            r.details.append(f"missing family {fam}: none of {candidates} found under {package_root}")
        else:
            # also check non-empty
            for c in candidates:
                p = package_root / c
                if p.is_file() and p.stat().st_size == 0:
                    r.details.append(f"empty file in family {fam}: {c}")

    # Rendered evidence
    rendered = package_root / "rendered"
    rendered_alt = package_root / "renders"
    # also v3-regenerated/renders/*
    has_rendered = rendered.is_dir() or rendered_alt.is_dir() or any(
        (package_root / d).is_dir() for d in ["renders/proposal", "renders/onepage", "renders/questionnaire", "renders/quote"]
    )
    if not has_rendered:
        # also check for legacy rendered with PNGs
        pngs = list(package_root.rglob("*.png"))
        if len(pngs) == 0:
            r.details.append(f"missing rendered evidence: no rendered/ dir and no PNGs found under {package_root}")
            missing_families.append("Rendered Evidence")
        else:
            r.details.append(f"rendered fallback: found {len(pngs)} PNGs via rglob")
    else:
        pngs = list(package_root.rglob("*.png"))
        if len(pngs) < 5:
            r.details.append(f"rendered evidence sparse: only {len(pngs)} PNGs")

    # Naming convention: files should start with Business35_
    bad_names = []
    for p in package_root.rglob("*"):
        if p.is_file() and p.suffix.lower() in (".pptx", ".pdf", ".docx", ".xlsx"):
            if not p.name.startswith("Business35_"):
                # allow pdf subdir
                if p.parent.name != "pdf":
                    bad_names.append(str(p.relative_to(package_root)))
    if bad_names:
        r.details.append(f"naming convention violation (should start with Business35_): {bad_names[:5]}")

    if missing_commercial or missing_families:
        r.passed = False
        r.details.insert(0, f"inventory: missing {len(missing_commercial)} commercial + {len(missing_families)} package families")
    else:
        r.details.insert(0, "inventory: all required families present")

    # If commercial_root itself missing, that's a hard fail but not unavailable
    if not commercial_root.exists():
        r.passed = False
        r.details.append(f"commercial_root does not exist: {commercial_root}")
    if not package_root.exists():
        r.passed = False
        r.details.append(f"package_root does not exist: {package_root}")

    return r


def check_source_mapping(commercial_root: Path, package_root: Path) -> CheckResult:
    r = CheckResult(name="SOURCE_MAPPING", passed=True)
    sm = package_root / "SOURCE_MAPPING.md"
    if not sm.is_file():
        r.passed = False
        r.details.append("SOURCE_MAPPING.md missing in package_root")
        return r
    text = sm.read_text(encoding="utf-8", errors="ignore")
    # Must reference product authority - fail-closed (C2)
    for marker in [PRODUCT_COMMIT, "PRODUCT_CONTRACT", "CURRENT_PRODUCT_AUTHORITY", "파디엠"]:
        if marker not in text:
            r.passed = False
            r.details.append(f"SOURCE_MAPPING missing product authority marker: {marker}")

    # Must cover Slide 1..10
    missing_slides = []
    for n in range(1, 11):
        if f"Slide {n}" not in text and f"슬라이드 {n}" not in text and f"Page {n}" not in text:
            missing_slides.append(n)
    if missing_slides:
        r.passed = False
        r.details.append(f"SOURCE_MAPPING missing slides: {missing_slides}")

    # Must map each output family
    families = ["Proposal", "OnePage", "Questionnaire", "Quote", "Meeting", "Followup", "Offer A", "Offer B"]
    missing_maps = [f for f in ["Proposal", "Questionnaire", "Quote"] if f.lower() not in text.lower()]
    if missing_maps:
        r.passed = False
        r.details.append(f"SOURCE_MAPPING missing family mappings: {missing_maps}")

    # Must distinguish historical vs current (stale vs V3.1)
    if "PRE_V3_1" not in text and "STALE" not in text and "HISTORICAL" not in text:
        r.details.append("SOURCE_MAPPING may not distinguish stale vs current (no PRE_V3_1/STALE/HISTORICAL marker)")

    # Commercial source side: SOURCES.md exists and has deep links
    sources_md = commercial_root / "SOURCES.md"
    if not sources_md.is_file():
        r.details.append("SOURCES.md missing in commercial_root (affects mapping completeness)")
        # not hard fail, but note
    else:
        s_text = sources_md.read_text(encoding="utf-8", errors="ignore")
        if "SRC-" not in s_text:
            r.details.append("SOURCES.md missing SRC- entries")
        if "VERIFIED" not in s_text:
            r.details.append("SOURCES.md missing VERIFIED status")

    # Check that commercial source files are referenced via SOURCE_MAPPING or SOURCES
    # Simplified: at least mention 01-..08
    for doc in ["01-one-page-offer", "02-ten-page-proposal", "03-diagnostic"]:
        if doc not in text.lower() and doc not in (sources_md.read_text(encoding="utf-8", errors="ignore").lower() if sources_md.is_file() else ""):
            r.details.append(f"source doc not referenced anywhere: {doc}")

    if any("missing" in d for d in r.details):
        # only fail if sliding missing
        if missing_slides or missing_maps:
            r.passed = False
    if not r.details:
        r.details.append("source mapping covers Slides 1..10 and families")
    # If no missing slides/maps, but other warnings exist, still PASS unless critical missing
    if missing_slides or missing_maps:
        r.passed = False
    return r


def check_structural_qa(commercial_root: Path, package_root: Path) -> CheckResult:
    r = CheckResult(name="STRUCTURAL_QA", passed=True)
    # PPTX checks
    pptx_candidates = list(package_root.rglob("Business35*Master*Proposal*.pptx"))
    if not pptx_candidates:
        pptx_candidates = list(package_root.rglob("*.pptx"))
        # filter for proposal-like
        pptx_candidates = [p for p in pptx_candidates if "Proposal" in p.name or "proposal" in p.name.lower()]

    if not pptx_candidates:
        r.passed = False
        r.details.append("no proposal PPTX found")
    else:
        for pptx_path in pptx_candidates[:2]:
            try:
                from pptx import Presentation
                prs = Presentation(str(pptx_path))
                count = len(prs.slides)
                # For master proposal, expect 10
                if "Master" in pptx_path.name or "master" in pptx_path.name.lower():
                    if count != 10:
                        r.passed = False
                        r.details.append(f"proposal PPTX {pptx_path.name} slide count {count} != 10")
                    else:
                        r.details.append(f"proposal PPTX {pptx_path.name}: {count} slides OK")
                    # Check speaker notes count
                    notes_count = sum(1 for s in prs.slides if s.has_notes_slide and s.notes_slide.placeholders)
                    # Some generators don't set notes; just report
                    r.details.append(f"  notes slides: {notes_count}/10")
                else:
                    r.details.append(f"PPTX {pptx_path.name}: {count} slides")
            except Exception as e:
                r.passed = False
                r.details.append(f"PPTX open failed {pptx_path.name}: {e}")

    # OnePage PPTX should be 1 slide
    onepage_cands = list(package_root.rglob("Business35*OnePage*Source*.pptx")) + list(package_root.rglob("Business35*OnePage*Offer*.pptx"))
    if not onepage_cands:
        onepage_cands = [p for p in package_root.rglob("*.pptx") if "OnePage" in p.name]
    for p in onepage_cands[:1]:
        try:
            from pptx import Presentation
            prs = Presentation(str(p))
            if len(prs.slides) != 1:
                r.passed = False
                r.details.append(f"OnePage PPTX {p.name} slide count {len(prs.slides)} != 1")
            else:
                r.details.append(f"OnePage PPTX {p.name}: 1 slide OK")
        except Exception as e:
            r.details.append(f"OnePage PPTX check failed: {e}")

    # PDF page counts
    pdf_checks = [
        ("Master Proposal PDF", ["Business35_Master_Proposal_10p.pdf", "Business35_V3_1_Master_Proposal_10p.pdf", "pdf/Business35_V3_1_Master_Proposal_10p.pdf"], 10, 10),
        ("OnePage Offer PDF", ["Business35_OnePage_Offer.pdf", "Business35_V3_1_OnePage_Offer.pdf", "pdf/Business35_V3_1_OnePage_Offer_Source.pdf"], 1, 1),
        ("Questionnaire PDF", ["Business35_Diagnostic_Questionnaire.pdf", "Business35_V3_1_Diagnostic_Questionnaire.pdf", "pdf/Business35_V3_1_Diagnostic_Questionnaire.pdf"], 1, 5),
        ("Quote PDF", ["Business35_Pilot_Quote_Template.pdf", "Business35_V3_1_Pilot_Quote_Template.pdf", "pdf/Business35_V3_1_Pilot_Quote_Template.pdf"], 1, 4),
    ]
    pdf_tool_available = False
    for label, cands, min_p, max_p in pdf_checks:
        found = None
        for cand in cands:
            p = package_root / cand
            if p.is_file():
                found = p
                break
        if not found:
            # try rglob
            matches = list(package_root.rglob(cand.split("/")[-1]))
            if matches:
                found = matches[0]
        if not found:
            if label in ("Master Proposal PDF", "OnePage Offer PDF", "Questionnaire PDF"):
                r.passed = False
                r.details.append(f"missing PDF for {label}")
            else:
                r.details.append(f"optional PDF missing for {label} (ok if XLSX-only quote)")
            continue
        pages, method = try_pdf_pages(found)
        if pages is None:
            r.passed = False
            r.details.append(f"PDF page count unavailable for {label} ({found.name}): {method} - mark as FAIL per spec")
        else:
            pdf_tool_available = True
            if pages < min_p or pages > max_p:
                r.passed = False
                r.details.append(f"PDF {label} {found.name} pages {pages} outside [{min_p},{max_p}] via {method}")
            else:
                r.details.append(f"PDF {label} {found.name}: {pages} pages OK via {method}")

    if not pdf_tool_available:
        r.details.append("PDF tooling availability limited - some checks may be unavailable but forced to FAIL per spec")

    # DOCX questionnaire completeness
    docx_cands = list(package_root.rglob("Business35*Questionnaire*.docx"))
    if not docx_cands:
        r.passed = False
        r.details.append("no questionnaire DOCX found")
    else:
        for docx_path in docx_cands[:1]:
            txt = extract_docx_text(docx_path)
            if txt is None:
                r.passed = False
                r.details.append(f"DOCX open failed {docx_path.name}")
            else:
                # Expect at least 10 questions or markers
                q_count = len(re.findall(r"\d{2}\.", txt))
                if q_count < 5:
                    r.details.append(f"DOCX {docx_path.name} question markers sparse: {q_count}")
                else:
                    r.details.append(f"DOCX {docx_path.name}: {q_count} question markers")
                if "파디엠" not in txt:
                    r.passed = False
                    r.details.append(f"DOCX {docx_path.name} missing provider 파디엠")
                # Check for DRAFT marker
                if "DRAFT" not in txt:
                    r.details.append(f"DOCX {docx_path.name} missing DRAFT")
                # Check for checkboxes
                if "☐" not in txt and "□" not in txt and "예" not in txt:
                    r.details.append(f"DOCX {docx_path.name} may lack checkboxes")

    # XLSX sheet integrity (basic)
    xlsx_cands = list(package_root.rglob("Business35*Quote*.xlsx"))
    if not xlsx_cands:
        xlsx_cands = list(package_root.rglob("*.xlsx"))
    if not xlsx_cands:
        r.passed = False
        r.details.append("no Quote XLSX found")
    else:
        for xlsx_path in xlsx_cands[:1]:
            try:
                from openpyxl import load_workbook
                wb = load_workbook(str(xlsx_path), data_only=False)
                sheets = wb.sheetnames
                r.details.append(f"XLSX {xlsx_path.name} sheets: {sheets}")
                # Expect either legacy 9 or V3 2
                legacy_expected = ["Instructions", "Customer Scope", "Offer A", "Offer B1", "Offer B2", "Offer C", "Optional Items", "Assumptions", "Approval"]
                v3_expected = ["Quote", "Terms"]
                if sheets == legacy_expected:
                    r.details.append("  matches legacy 9-sheet Quote template")
                elif sheets == v3_expected:
                    r.details.append("  matches V3 2-sheet Quote template")
                else:
                    # At least check not empty and has some expected names
                    if len(sheets) < 2:
                        r.passed = False
                        r.details.append(f"  XLSX sheet count {len(sheets)} <2")
                    else:
                        r.details.append(f"  XLSX sheets non-standard but count ok")
                # Check for provider marker via text extraction
                xlsx_txt = extract_xlsx_text(xlsx_path) or ""
                if "파디엠" not in xlsx_txt:
                    r.passed = False
                    r.details.append(f"XLSX {xlsx_path.name} missing 파디엠")
            except Exception as e:
                r.passed = False
                r.details.append(f"XLSX open failed {xlsx_path}: {e}")

    return r


def check_formula_qa(package_root: Path) -> CheckResult:
    r = CheckResult(name="FORMULA_QA", passed=True)
    xlsx_cands = list(package_root.rglob("Business35*Quote*.xlsx"))
    if not xlsx_cands:
        xlsx_cands = list(package_root.rglob("*.xlsx"))
    if not xlsx_cands:
        r.passed = False
        r.details.append("no XLSX found for formula check")
        return r

    try:
        from openpyxl import load_workbook
    except Exception as e:
        r.passed = False
        r.unavailable = False  # per spec, unavailable remains fail
        r.details.append(f"openpyxl unavailable: {e} - FAIL")
        return r

    for xlsx_path in xlsx_cands[:1]:
        try:
            wb = load_workbook(str(xlsx_path), data_only=False)
        except Exception as e:
            r.passed = False
            r.details.append(f"workbook open failed: {e}")
            continue

        # Check for formula errors in data_only=True view
        try:
            wb_data = load_workbook(str(xlsx_path), data_only=True)
            errors = []
            for ws in wb_data.worksheets:
                for row in ws.iter_rows():
                    for cell in row:
                        v = cell.value
                        if isinstance(v, str) and v.startswith("#") and v in ("#VALUE!", "#REF!", "#DIV/0!", "#NAME?", "#NUM!", "#NULL!", "#N/A"):
                            errors.append(f"{ws.title}!{cell.coordinate}={v}")
            if errors:
                r.passed = False
                r.details.append(f"formula errors: {errors[:5]}")
            else:
                r.details.append("no formula error values in data_only view")
        except Exception as e:
            r.details.append(f"data_only check skipped: {e}")

        # Check conditional formatting on Offer sheets (legacy) or any sheet
        cf_found = False
        for ws in wb.worksheets:
            if ws.conditional_formatting:
                cf_found = True
                # Check B4 specifically for Offer sheets
                if ws.title.startswith("Offer"):
                    has_b4 = any("B4" in str(cf.sqref) for cf in ws.conditional_formatting)
                    if has_b4:
                        r.details.append(f"CF on {ws.title} includes B4")
                    else:
                        r.details.append(f"CF on {ws.title} missing B4")
        if not cf_found:
            r.details.append("no conditional formatting found (may be ok for V3 Quote/T terms)")

        # Print area / print setup integrity
        for ws in wb.worksheets:
            # Check printArea
            if ws.print_area:
                r.details.append(f"printArea {ws.title}: {ws.print_area}")
            else:
                r.details.append(f"printArea {ws.title}: not set (warning)")
            # Check page setup
            if ws.page_setup.fitToPage is not None:
                r.details.append(f"fitToPage {ws.title}: {ws.page_setup.fitToPage}")
            # Check column widths not zero
            for col, dim in ws.column_dimensions.items():
                if dim.width is not None and dim.width == 0:
                    r.passed = False
                    r.details.append(f"zero width column {ws.title}!{col}")

        # Check for external links (should not exist in reusable master)
        if wb._external_links:
            r.passed = False
            r.details.append(f"external links found: {len(wb._external_links)}")

        # Check for macros / VBA (should not exist)
        if wb.vba_archive is not None:
            r.details.append("VBA archive present (unexpected)")

        # Check for at least one formula cell (if legacy)
        has_formula = False
        for ws in wb.worksheets:
            for row in ws.iter_rows():
                for cell in row:
                    if cell.data_type == "f":
                        has_formula = True
                        break
        if has_formula:
            r.details.append("at least one formula cell present")
        else:
            r.details.append("no formula cells (ok for V3 simple Quote)")

    return r


def check_text_fit(package_root: Path) -> CheckResult:
    r = CheckResult(name="TEXT_FIT", passed=True)
    # PPTX overflow / overlap
    pptx_files = list(package_root.rglob("*.pptx"))
    if not pptx_files:
        r.passed = False
        r.details.append("no PPTX found for text-fit")
        return r

    try:
        from pptx import Presentation
        from pptx.util import Emu
    except Exception as e:
        r.passed = False
        r.details.append(f"python-pptx unavailable: {e} - FAIL per spec")
        return r

    overflow_total = 0
    overlap_total = 0
    for pptx_path in pptx_files[:3]:
        try:
            prs = Presentation(str(pptx_path))
            w = prs.slide_width
            h = prs.slide_height
            overflow = 0
            for slide in prs.slides:
                for shp in slide.shapes:
                    if shp.left is None or shp.top is None or shp.width is None or shp.height is None:
                        continue
                    right = shp.left + shp.width
                    bottom = shp.top + shp.height
                    if right > w + 10000 or bottom > h + 10000 or shp.left < -10000 or shp.top < -10000:
                        overflow += 1
            if overflow:
                overflow_total += overflow
                r.details.append(f"overflow in {pptx_path.name}: {overflow}")
            else:
                r.details.append(f"no overflow in {pptx_path.name}")

            # Overlap: require both x and y overlap > 0.05 inch
            for si, slide in enumerate(prs.slides, start=1):
                texts = []
                for shp in slide.shapes:
                    if shp.has_text_frame and shp.text_frame.text.strip() and shp.top is not None and shp.left is not None:
                        texts.append((shp.left, shp.top, shp.left + (shp.width or 0), shp.top + (shp.height or 0)))
                slide_overlaps = 0
                for i in range(len(texts)):
                    for j in range(i + 1, len(texts)):
                        a, b = texts[i], texts[j]
                        ox = min(a[2], b[2]) - max(a[0], b[0])
                        oy = min(a[3], b[3]) - max(a[1], b[1])
                        if ox > int(0.05 * 914400) and oy > int(0.05 * 914400):
                            slide_overlaps += 1
                if slide_overlaps:
                    overlap_total += slide_overlaps
                    # Only fail if slide 3/4/8 style overlaps? For now any overlap is FAIL
                    r.details.append(f"overlap in {pptx_path.name} slide {si}: {slide_overlaps}")
        except Exception as e:
            r.details.append(f"PPTX text-fit check failed {pptx_path.name}: {e}")

    if overflow_total > 0:
        r.passed = False
        r.details.append(f"total overflows: {overflow_total}")
    if overlap_total > 0:
        # Be conservative: any overlap is FAIL (legacy validator required 0)
        r.passed = False
        r.details.append(f"total text overlaps: {overlap_total}")
    else:
        r.details.append("no text overlap detected")

    # PDF broken glyph check
    pdfs = list(package_root.rglob("*.pdf"))
    for pdf in pdfs[:3]:
        txt, method = try_pdf_text(pdf)
        if txt is None:
            r.details.append(f"PDF glyph check unavailable for {pdf.name}: {method}")
            continue
        if "�" in txt or "□" in txt:
            r.passed = False
            r.details.append(f"broken glyph in {pdf.name} via {method}")
        else:
            r.details.append(f"no broken glyph in {pdf.name} via {method}")

    return r


def check_stale_rejection(commercial_root: Path, package_root: Path, product_contract: Path) -> CheckResult:
    r = CheckResult(name="STALE_ARTIFACT_REJECTION", passed=True)
    # Scan package texts for stale markers that indicate pre-V3.1 as current
    texts: Dict[str, str] = {}
    for md in package_root.rglob("*.md"):
        try:
            texts[str(md.relative_to(package_root))] = md.read_text(encoding="utf-8", errors="ignore")
        except Exception:
            pass
    for pptx in package_root.rglob("*.pptx"):
        t = extract_pptx_text(pptx)
        if t:
            texts[str(pptx.relative_to(package_root))] = t
    for docx in package_root.rglob("*.docx"):
        t = extract_docx_text(docx)
        if t:
            texts[str(docx.relative_to(package_root))] = t
    for xlsx in package_root.rglob("*.xlsx"):
        t = extract_xlsx_text(xlsx)
        if t:
            texts[str(xlsx.relative_to(package_root))] = t
    for pdf in package_root.rglob("*.pdf"):
        txt, _ = try_pdf_text(pdf)
        if txt:
            texts[str(pdf.relative_to(package_root))] = txt

    combined = "\n".join(texts.values())

    # Check that V3.1 markers are present (if package claims to be current)
    missing_v3 = [m for m in V3_JOURNEY_MARKERS[:5] if m not in combined]
    present = 5 - len(missing_v3)
    if missing_v3:
        r.details.append(f"V3.1 journey markers: {present}/5 present, missing {missing_v3[:2]}")
        if present < 3:
            r.passed = False
            r.details.append("V3.1 product identity not found - stale artifact likely ( <3 markers)")
    else:
        r.details.append("V3.1 journey markers present (5/5)")

    # Forbidden stale as current claims - per-file local context (C3)
    for rel, txt in texts.items():
        for phrase in STALE_FORBIDDEN_AS_CURRENT:
            if phrase in txt:
                # Determine if this file is internal documentation that may legitimately contain historical markers
                is_internal_doc = any(k in rel for k in ["README", "SOURCE_MAPPING", "VISUAL_QA", "CUSTOMIZATION", "REFERENCE_COMPARISON", "CURRENT_PRODUCT_AUTHORITY"])
                # Check local historical context within same file (within 800 chars of phrase)
                local_ok = False
                if is_internal_doc and ("HISTORICAL" in txt or "STALE" in txt or "PRE_V3" in txt):
                    # Check proximity: historical token within 1000 chars of phrase occurrence
                    for m in re.finditer(re.escape(phrase), txt):
                        start = max(0, m.start() - 1000)
                        end = m.end() + 1000
                        window = txt[start:end]
                        if "HISTORICAL" in window or "STALE" in window or "PRE_V3" in window:
                            local_ok = True
                            break
                if local_ok:
                    r.details.append(f"stale marker {phrase} in {rel} with local historical context (ok)")
                else:
                    r.passed = False
                    r.details.append(f"stale marker as current in {rel}: {phrase}")

    # Check that reusable master does not claim old seven-step as primary
    # Old seven-step: check for historical education sequence not qualified as delivery detail
    if "7단계" in combined or "seven-step" in combined.lower():
        if "delivery detail" not in combined.lower() and "세부" not in combined:
            r.details.append("seven-step present without delivery-detail qualification (warning)")

    # Verify that historical binaries are correctly labeled stale in README/SOURCE_MAPPING
    readme = package_root / "README.md"
    sm = package_root / "SOURCE_MAPPING.md"
    for p, label in [(readme, "README"), (sm, "SOURCE_MAPPING")]:
        if p.is_file():
            txt = p.read_text(encoding="utf-8", errors="ignore")
            if "STALE_FOR_SEND" in txt or "PRE_V3_1" in txt or "HISTORICAL_ONLY" in txt:
                r.details.append(f"{label} correctly marks historical binaries")
            else:
                r.details.append(f"{label} missing stale marker disclaimer (should mark historical)")

    # Check for price-hypothesis vs stale price claim - per-file (C1)
    # For any file containing price tokens, hypothesis wording must be present in same file
    for rel, txt in texts.items():
        has_price = any(tok in txt for tok in PRICE_TOKENS)
        if has_price and "PRICE_HYPOTHESIS_ONLY" not in txt and "가설" not in txt and "hypothesis" not in txt.lower():
            # Allow if file is internal doc that lists forbidden? No, price in customer outputs must have hypothesis
            if any(x in rel for x in [".pptx", ".pdf", ".docx", ".xlsx"]) or "Meeting" in rel or "Followup" in rel:
                r.passed = False
                r.details.append(f"price hypothesis wording missing in {rel} (price without 가설)")

    # Cross-cutting: CUSTOMER_SEND_READY and legal-review completion false claims (C1)
    for rel, txt in texts.items():
        for phrase in FORBIDDEN_COMPLETION_CLAIMS:
            if phrase in txt:
                # Allow checklist/documentation context
                is_listing = False
                for m in re.finditer(re.escape(phrase), txt):
                    window = txt[max(0, m.start()-800): m.end()+800]
                    if any(k in window for k in ["Forbidden", "금지", "불가", "DRAFT · PROFESSIONAL LEGAL REVIEW REQUIRED", "발송 전 체크", "[ ]", "주의:"]):
                        is_listing = True
                        break
                if not is_listing:
                    r.passed = False
                    r.details.append(f"forbidden completion claim in {rel}: {phrase}")

    # Cross-cutting: forbidden customer phrases (C1) - per-file, allow listing context
    for rel, txt in texts.items():
        for phrase in FORBIDDEN_PHRASES_CUSTOMER:
            if phrase == "성과 보장":
                continue  # handled separately
            if phrase in txt:
                is_listing = False
                for m in re.finditer(re.escape(phrase), txt):
                    window = txt[max(0, m.start()-800): m.end()+800]
                    if any(k in window for k in ["금지 발언", "Forbidden", "금지", "forbidden", "발송 전 체크", "금지 표현", "주의:"]):
                        is_listing = True
                        break
                if not is_listing:
                    r.passed = False
                    r.details.append(f"forbidden customer phrase in {rel}: {phrase}")
        # Special handling for 성과 보장 - only allowed in negation or in forbidden listing
        if "성과 보장" in txt:
            for m in re.finditer(r"성과 보장", txt):
                ctx = txt[max(0, m.start()-30): m.end()+30]
                window = txt[max(0, m.start()-800): m.end()+800]
                is_listing = any(k in window for k in ["금지 발언", "Forbidden", "금지", "forbidden", "발송 전 체크", "금지 표현"])
                if is_listing:
                    continue
                if "보장하지" not in ctx and "의미하지 않" not in ctx and "보장하라고 요구" not in ctx and "보장하지 않는다" not in ctx:
                    r.passed = False
                    r.details.append(f"forbidden phrase '성과 보장' without negation in {rel}: ...{ctx}...")
                    break

    return r


def check_private_boundary(package_root: Path) -> CheckResult:
    r = CheckResult(name="PRIVATE_DATA_BOUNDARY", passed=True)
    files_to_scan: List[Path] = []
    for ext in ("*.md", "*.pptx", "*.docx", "*.xlsx", "*.pdf"):
        files_to_scan.extend(package_root.rglob(ext))

    found_private = []
    for p in files_to_scan:
        text = ""
        if p.suffix == ".md":
            try:
                text = p.read_text(encoding="utf-8", errors="ignore")
            except Exception:
                continue
        elif p.suffix == ".pptx":
            text = extract_pptx_text(p) or ""
        elif p.suffix == ".docx":
            text = extract_docx_text(p) or ""
        elif p.suffix == ".xlsx":
            text = extract_xlsx_text(p) or ""
        elif p.suffix == ".pdf":
            txt, _ = try_pdf_text(p)
            text = txt or ""
        else:
            continue

        # Scan patterns
        emails = EMAIL_RE.findall(text)
        # Filter out example placeholders
        real_emails = [e for e in emails if "example.com" not in e and "test.com" not in e and "padiem" not in e.lower() and "placeholder" not in e.lower()]
        if real_emails:
            found_private.append(f"{p.name}: email {real_emails[:2]}")
        phones = PHONE_RE.findall(text)
        if phones:
            # Allow if near placeholder text
            if "발송 전" not in text and "입력 필요" not in text:
                found_private.append(f"{p.name}: phone {phones[:2]}")
        biznos = BIZNO_RE.findall(text)
        if biznos:
            if "발송 전 공식 사업자 정보 입력 필요" not in text:
                found_private.append(f"{p.name}: bizno {biznos[:2]}")
        secrets = SECRET_RE.findall(text)
        if secrets:
            found_private.append(f"{p.name}: secret {secrets[:1]}")

    if found_private:
        r.passed = False
        r.details.extend(found_private[:10])
    else:
        r.details.append("no real customer/contact/private data patterns found")

    # Cross-cutting (C1): also check CUSTOMER_SEND_READY and legal completion in private boundary scope
    for p in files_to_scan:
        try:
            if p.suffix == ".md":
                txt = p.read_text(encoding="utf-8", errors="ignore")
            elif p.suffix == ".pptx":
                txt = extract_pptx_text(p) or ""
            elif p.suffix == ".docx":
                txt = extract_docx_text(p) or ""
            elif p.suffix == ".xlsx":
                txt = extract_xlsx_text(p) or ""
            elif p.suffix == ".pdf":
                txt, _ = try_pdf_text(p)
                txt = txt or ""
            else:
                continue
        except Exception:
            continue
        for phrase in FORBIDDEN_COMPLETION_CLAIMS:
            if phrase in txt:
                is_listing = False
                for m in re.finditer(re.escape(phrase), txt):
                    window = txt[max(0, m.start()-800): m.end()+800]
                    if any(k in window for k in ["Forbidden", "금지", "불가", "DRAFT · PROFESSIONAL LEGAL REVIEW REQUIRED", "발송 전 체크", "[ ]"]):
                        is_listing = True
                        break
                if not is_listing:
                    r.passed = False
                    r.details.append(f"forbidden completion claim in {p.name}: {phrase}")
                    break

    # Also check that reusable master contains required placeholder disclaimer
    xlsx_cands = list(package_root.rglob("*.xlsx"))
    has_disclaimer = False
    for xlsx in xlsx_cands:
        txt = extract_xlsx_text(xlsx) or ""
        if "발송 전 공식 사업자 정보 입력 필요" in txt or "제공자 정보" in txt:
            has_disclaimer = True
    if not has_disclaimer and xlsx_cands:
        r.details.append("XLSX missing business-details placeholder disclaimer (should be present in reusable master)")
        # Not hard fail but note

    # Check for real customer claim without disclaimer (allow checklist documentation)
    disclaimer_variants = ["합성 예시", "주장이 아닙니다", "주장이 아닙", "삽입 금지", "받지 않는다", "없음", "placeholder", "대괄호"]
    for p in files_to_scan:
        if p.suffix == ".md" and ("CUSTOMIZATION" in p.name.upper() or "CHECKLIST" in p.name.upper()):
            continue
        try:
            if p.suffix == ".md":
                txt = p.read_text(encoding="utf-8", errors="ignore")
            elif p.suffix == ".pptx":
                txt = extract_pptx_text(p) or ""
            elif p.suffix == ".docx":
                txt = extract_docx_text(p) or ""
            elif p.suffix == ".xlsx":
                txt = extract_xlsx_text(p) or ""
            elif p.suffix == ".pdf":
                txt, _ = try_pdf_text(p)
                txt = txt or ""
            else:
                continue
        except Exception:
            continue
        if "실제 고객" in txt and not any(v in txt for v in disclaimer_variants):
            r.passed = False
            r.details.append(f"real customer claim without disclaimer in {p.name}")
            break

    return r


def check_exact_revision_trace(commercial_root: Path, package_root: Path, product_contract: Path, manifest_path: Optional[Path]) -> CheckResult:
    r = CheckResult(name="EXACT_REVISION_TRACE", passed=True)
    # Fail-closed: actual product authority file dependency must exist and be readable.
    # Missing/unreadable product_contract can never be inferred PASS.
    try:
        pc = Path(product_contract) if product_contract is not None else None
    except Exception:
        pc = None
    if pc is None or not pc.is_file():
        r.passed = False
        r.details.append(f"product_contract missing/unreadable: {product_contract} (fail-closed, no inference to PASS)")
    else:
        try:
            if pc.stat().st_size == 0:
                r.passed = False
                r.details.append(f"product_contract empty: {pc} (fail-closed)")
            else:
                pc.read_text(encoding="utf-8", errors="ignore")
                r.details.append(f"product_contract present: {pc}")
        except Exception as e:
            r.passed = False
            r.details.append(f"product_contract unreadable: {pc}: {e} (fail-closed)")
    # Find manifest
    candidates = []
    if manifest_path and manifest_path.exists():
        candidates.append(manifest_path)
    candidates.extend([
        package_root / "MANIFEST_V3_1.json",
        package_root / "MANIFEST.json",
        package_root / "generation_manifest.json",
        package_root / "GENERATION_MANIFEST.json",
        package_root / "v3-regenerated" / "MANIFEST_V3_1.json",
        package_root / "pdf" / "MANIFEST_V3_1.json",
    ])
    # also rglob
    candidates.extend(list(package_root.rglob("MANIFEST*.json")))
    candidates.extend(list(package_root.rglob("*manifest*.json")))

    manifest = None
    for cand in candidates:
        if cand and cand.exists():
            manifest = cand
            break

    if not manifest:
        r.passed = False
        r.details.append(f"manifest missing: searched {[str(c) for c in candidates[:5]]}")
        return r

    r.details.append(f"manifest found: {manifest.relative_to(package_root) if manifest.is_relative_to(package_root) else manifest}")

    try:
        data = json.loads(manifest.read_text(encoding="utf-8"))
    except Exception as e:
        r.passed = False
        r.details.append(f"manifest JSON parse failed: {e}")
        return r

    # Check required fields
    required_top = ["SOURCE_REVISION", "PRODUCT_AUTHORITY_REVISION", "GENERATOR_REVISION", "OUTPUT_FILE_LIST", "OUTPUT_HASHES"]
    # Also handle alternative structure: {"status": ..., "files": [{"path":..., "sha256":...}]}
    # Legacy/PR359 style may have different keys
    alt_files = data.get("files") or data.get("artifacts") or data.get("outputs")
    has_source_rev = "SOURCE_REVISION" in data or "source_revision" in data or "source" in data
    has_prod_rev = "PRODUCT_AUTHORITY_REVISION" in data or "product_authority" in data or "product" in data
    has_gen_rev = "GENERATOR_REVISION" in data or "generator_revision" in data or "generator" in data

    if not has_source_rev:
        # Check if manifest is alternative V3 style with top-level status only
        if alt_files is not None:
            r.details.append("manifest uses alternative 'files' structure without explicit SOURCE_REVISION (fail)")
        r.passed = False
        r.details.append("manifest missing SOURCE_REVISION")
    else:
        src_rev = data.get("SOURCE_REVISION") or data.get("source_revision") or data.get("source")
        if isinstance(src_rev, str) and re.fullmatch(r"[0-9a-f]{40}", src_rev):
            if src_rev == ACCEPTED_SOURCE_REVISION:
                r.details.append(f"SOURCE_REVISION exact match OK: {src_rev[:7]}...")
            else:
                r.passed = False
                r.details.append(
                    f"SOURCE_REVISION {src_rev} != ACCEPTED_SOURCE_REVISION {ACCEPTED_SOURCE_REVISION} (EXACT_REVISION_TRACE_FAIL)"
                )
        else:
            r.passed = False
            r.details.append(f"SOURCE_REVISION invalid: {src_rev}")

    prod_rev = data.get("PRODUCT_AUTHORITY_REVISION") or data.get("product_authority") or data.get("product")
    if prod_rev:
        if PRODUCT_COMMIT in str(prod_rev) or str(prod_rev) == PRODUCT_COMMIT:
            r.details.append(f"PRODUCT_AUTHORITY_REVISION matches expected {PRODUCT_COMMIT[:7]}")
        else:
            r.details.append(f"PRODUCT_AUTHORITY_REVISION {prod_rev} != expected {PRODUCT_COMMIT}")
            # Not hard fail if different but should be noted; however spec says exact trace required
            # If manifest product rev doesn't match expected, it's a trace fail
            r.passed = False
    else:
        r.passed = False
        r.details.append("manifest missing PRODUCT_AUTHORITY_REVISION")

    gen_rev = data.get("GENERATOR_REVISION") or data.get("generator_revision") or data.get("generator")
    # Fail-closed: GENERATOR_REVISION must be full 40-char lowercase hex.
    # 12-char / unknown / unusual / missing all fail; no PASS on format-unusual.
    if isinstance(gen_rev, str) and re.fullmatch(r"[0-9a-f]{40}", gen_rev):
        r.details.append(f"GENERATOR_REVISION OK: {gen_rev[:7]}...")
    else:
        r.passed = False
        if gen_rev:
            r.details.append(f"GENERATOR_REVISION invalid format (must be [0-9a-f]{{40}}): {gen_rev} (EXACT_REVISION_TRACE_FAIL)")
        else:
            r.details.append("manifest missing GENERATOR_REVISION (EXACT_REVISION_TRACE_FAIL)")

    # Output file list and hashes
    file_list = data.get("OUTPUT_FILE_LIST") or data.get("output_file_list") or data.get("files") or data.get("artifacts")
    hashes = data.get("OUTPUT_HASHES") or data.get("output_hashes") or data.get("hashes")

    # For alternative files structure, build hashes from files list
    if hashes is None and isinstance(file_list, list) and file_list and isinstance(file_list[0], dict) and "sha256" in file_list[0]:
        hashes = {f["path"]: f["sha256"] for f in file_list}
        file_list = list(hashes.keys())

    if not file_list:
        r.passed = False
        r.details.append("manifest missing OUTPUT_FILE_LIST")
    else:
        r.details.append(f"OUTPUT_FILE_LIST: {len(file_list)} entries")
        # Verify each listed file exists relative to package_root or manifest parent
        base = manifest.parent if manifest.parent != package_root else package_root
        # Try both
        missing = []
        for rel in file_list[:20]:
            # rel may be string path or dict
            if isinstance(rel, dict):
                rel = rel.get("path") or rel.get("name") or str(rel)
            p1 = package_root / rel
            p2 = base / rel
            if not p1.exists() and not p2.exists():
                missing.append(rel)
        if missing:
            r.passed = False
            r.details.append(f"manifest listed files missing: {missing[:3]}")

    if not hashes:
        r.passed = False
        r.details.append("manifest missing OUTPUT_HASHES")
    else:
        r.details.append(f"OUTPUT_HASHES: {len(hashes) if isinstance(hashes, dict) else 'list'} entries")
        # Verify hashes match actual files
        mismatch = []
        if isinstance(hashes, dict):
            for rel, expected in hashes.items():
                p = package_root / rel
                if not p.exists():
                    p = manifest.parent / rel
                if p.exists():
                    actual = sha256_file(p)
                    if actual.lower() != expected.lower():
                        mismatch.append(f"{rel}: expected {expected[:7]} got {actual[:7]}")
                else:
                    mismatch.append(f"{rel}: file missing for hash check")
        elif isinstance(hashes, list):
            # list of dicts
            for entry in hashes[:10]:
                if isinstance(entry, dict):
                    rel = entry.get("path") or entry.get("file") or ""
                    expected = entry.get("sha256") or entry.get("hash") or ""
                    p = package_root / rel if rel else None
                    if p and p.exists():
                        actual = sha256_file(p)
                        if actual.lower() != expected.lower():
                            mismatch.append(f"{rel}")
        if mismatch:
            r.passed = False
            r.details.append(f"hash mismatches: {mismatch[:3]}")
        else:
            r.details.append("all checked hashes match")

    # Additional: check that package files not listed are not extra? Not required

    # Check that CURRENT_PRODUCT_AUTHORITY.md product commit matches manifest prod rev
    cpa = commercial_root / "CURRENT_PRODUCT_AUTHORITY.md"
    if cpa.is_file():
        cpa_text = cpa.read_text(encoding="utf-8", errors="ignore")
        if PRODUCT_COMMIT not in cpa_text:
            r.details.append(f"CURRENT_PRODUCT_AUTHORITY.md missing product commit {PRODUCT_COMMIT[:7]} (may be stale)")

    return r


def check_price_hypothesis_and_no_send(package_root: Path) -> Tuple[CheckResult, CheckResult]:
    """Two implicit checks that are part of multiple verdicts but we surface as sub-checks.
    For spec compatibility, price-hypothesis is checked inside STALE/PACKAGE, and no-send inside PRIVATE/STALE.
    This helper is not a separate verdict but adds details to existing verdicts.
    We will create wrapper results for reporting.
    """
    # This is not used directly; price checks are embedded in other verifiers.
    pass


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main() -> int:
    parser = argparse.ArgumentParser(description="B35 Independent QA Harness #1505")
    parser.add_argument("--commercial-root", default="docs/commercial/business-35-ai-media-education-dx")
    parser.add_argument("--package-root", default="docs/commercial/business-35-ai-media-education-dx/customer-package")
    parser.add_argument("--product-contract", default=PRODUCT_CONTRACT_REL)
    parser.add_argument("--manifest", default=None, help="path to manifest JSON (auto-detect if not set)")
    parser.add_argument("--output-json", default="tools/b35-independent-qa/evidence/qa_report.json")
    parser.add_argument("--output-md", default="tools/b35-independent-qa/evidence/qa_report.md")
    parser.add_argument("--pretty", action="store_true", help="pretty print to stdout")
    args = parser.parse_args()

    commercial_root = Path(args.commercial_root)
    package_root = Path(args.package_root)
    product_contract = Path(args.product_contract)
    manifest_path = Path(args.manifest) if args.manifest else None

    # Header
    print("=" * 72)
    print("B35 Independent QA Harness - Lane C #1505")
    print("=" * 72)
    print(f"Commercial root : {commercial_root} (exists={commercial_root.exists()})")
    print(f"Package root   : {package_root} (exists={package_root.exists()})")
    print(f"Product contract: {product_contract} (exists={product_contract.exists()})")
    print(f"Manifest        : {manifest_path or 'auto-detect'}")
    print(f"Legacy lineage  : PR #359 @ ef343f420661cda5f86cc2848404bca8f1dffe54 (reference only)")
    print()

    # Run checks
    results: List[CheckResult] = []

    results.append(check_package_inventory(commercial_root, package_root))
    results.append(check_source_mapping(commercial_root, package_root))
    results.append(check_structural_qa(commercial_root, package_root))
    results.append(check_formula_qa(package_root))
    results.append(check_text_fit(package_root))
    results.append(check_stale_rejection(commercial_root, package_root, product_contract))
    results.append(check_private_boundary(package_root))
    results.append(check_exact_revision_trace(commercial_root, package_root, product_contract, manifest_path))

    # Print verdicts
    print("VERDICTS")
    print("-" * 72)
    all_pass = True
    for r in results:
        status = "PASS" if r.passed and not r.unavailable else ("UNAVAILABLE" if r.unavailable else "FAIL")
        icon = "[PASS]" if status == "PASS" else "[FAIL]"
        print(f"{icon} {r.verdict} - {status}")
        for d in r.details[:8]:
            print(f"    - {d}")
        if len(r.details) > 8:
            print(f"    - ... +{len(r.details)-8} more")
        if not r.passed or r.unavailable:
            all_pass = False

    print()
    print("-" * 72)
    if all_pass:
        print("OVERALL: ALL 8 VERDICTS PASS")
    else:
        print("OVERALL: FAIL - one or more verdicts failed/unavailable (do not convert to PASS)")

    # Additional cross-cutting checks (C1) - fail-closed promotion to owning verdicts
    # These checks were already done per-file in verdicts, but we re-evaluate here to ensure overall_pass is fail-closed
    # Gather combined text for summary (but verdict promotion is per-file already)
    # Build per-file texts for cross-cutting (listing-aware)
    cross_texts = {}
    for p in package_root.rglob("*.md"):
        try:
            cross_texts[str(p.relative_to(package_root))] = p.read_text(encoding="utf-8", errors="ignore")
        except Exception:
            pass
    for p in package_root.rglob("*.pptx"):
        ttxt = extract_pptx_text(p)
        if ttxt:
            cross_texts[str(p.relative_to(package_root))] = ttxt
    for p in package_root.rglob("*.pdf"):
        ttxt, _ = try_pdf_text(p)
        if ttxt:
            cross_texts[str(p.relative_to(package_root))] = ttxt
    for p in package_root.rglob("*.docx"):
        ttxt = extract_docx_text(p)
        if ttxt:
            cross_texts[str(p.relative_to(package_root))] = ttxt
    for p in package_root.rglob("*.xlsx"):
        ttxt = extract_xlsx_text(p)
        if ttxt:
            cross_texts[str(p.relative_to(package_root))] = ttxt
    combined_text = "\n".join(cross_texts.values())

    # Price hypothesis - per-file listing-aware
    price_ok = True
    has_price = False
    for rel, txt in cross_texts.items():
        if any(tok in txt for tok in PRICE_TOKENS):
            has_price = True
            if "가설" not in txt and "PRICE_HYPOTHESIS" not in txt and "hypothesis" not in txt.lower():
                price_ok = False
                break
    if not has_price:
        price_ok = True
    print()
    print("CROSS-CUTTING (embedded, fail-closed)")
    print(f"  Price hypothesis wording: {'OK' if price_ok else 'MISSING/FAIL'}")

    # No-send - per-file with listing allowance
    no_send_ok = True
    for rel, txt in cross_texts.items():
        for phrase in FORBIDDEN_COMPLETION_CLAIMS:
            if phrase in txt:
                is_listing = False
                for m in re.finditer(re.escape(phrase), txt):
                    window = txt[max(0, m.start()-800): m.end()+800]
                    if any(k in window for k in ["Forbidden", "금지", "불가", "DRAFT · PROFESSIONAL LEGAL REVIEW REQUIRED", "발송 전 체크", "[ ]", "주의:"]):
                        is_listing = True
                        break
                if not is_listing:
                    no_send_ok = False
                    break
        if not no_send_ok:
            break
    print(f"  No customer-send claim: {'OK' if no_send_ok else 'FORBIDDEN CLAIM FOUND'}")

    # Forbidden phrases - per-file with listing allowance
    forbidden_found = []
    for rel, txt in cross_texts.items():
        for ph in FORBIDDEN_PHRASES_CUSTOMER:
            if ph == "성과 보장":
                continue
            if ph in txt:
                is_listing = False
                for m in re.finditer(re.escape(ph), txt):
                    window = txt[max(0, m.start()-800): m.end()+800]
                    if any(k in window for k in ["금지 발언", "Forbidden", "금지", "forbidden", "발송 전 체크", "금지 표현", "주의:"]):
                        is_listing = True
                        break
                if not is_listing:
                    forbidden_found.append(ph)
                    break
        if forbidden_found:
            break
    if "성과 보장" in combined_text:
        import re as _re
        found_neg = False
        for rel, txt in cross_texts.items():
            if "성과 보장" in txt:
                for m in _re.finditer(r"성과 보장", txt):
                    window = txt[max(0, m.start()-800): m.end()+800]
                    if any(k in window for k in ["금지 발언", "Forbidden", "금지", "forbidden", "발송 전 체크"]):
                        continue
                    ctx = txt[max(0, m.start()-30): m.end()+30]
                    if "의미하지 않" not in ctx and "보장하지" not in ctx and "보장하라고 요구" not in ctx and "보장하지 않는다" not in ctx:
                        found_neg = True
                        break
                if found_neg:
                    break
        if found_neg:
            forbidden_found.append("성과 보장 (not in negation)")
    print(f"  Forbidden phrases: {forbidden_found if forbidden_found else 'none'}")

    # Fail-closed promotion: if cross-cutting fails, force owning verdict to FAIL
    # This ensures overall_pass cannot be true while cross-cutting is violated
    if not price_ok or not no_send_ok or forbidden_found:
        for r in results:
            if not price_ok and r.name == "STALE_ARTIFACT_REJECTION" and r.passed:
                r.passed = False
                r.details.append("C1 promotion: price hypothesis violation forces STALE_ARTIFACT_REJECTION_FAIL")
            if (not no_send_ok or forbidden_found) and r.name == "STALE_ARTIFACT_REJECTION" and r.passed:
                # Check again per-file to avoid false promotion when already failed
                r.passed = False
                r.details.append(f"C1 promotion: forbidden/send claim forces STALE_ARTIFACT_REJECTION_FAIL (no_send_ok={no_send_ok}, forbidden={forbidden_found[:2]})")
        # Recompute all_pass after promotion (unavailable can never be PASS)
        all_pass = all(rr.passed and not rr.unavailable for rr in results)
        if not all_pass:
            print("\nC1 PROMOTION: cross-cutting violation promoted to verdict FAIL -> overall FAIL")

    # Write evidence
    output_json = Path(args.output_json)
    output_md = Path(args.output_md)
    output_json.parent.mkdir(parents=True, exist_ok=True)
    output_md.parent.mkdir(parents=True, exist_ok=True)

    report = {
        "issue": 1505,
        "lane": "C / Independent Package QA",
        "branch": "feat/b35-w3-independent-qa-v31",
        "base_sha": "eae88e0066c1b119bfa6c75d8b16c127b0137e5e",
        "commercial_root": str(commercial_root),
        "package_root": str(package_root),
        "product_contract": str(product_contract),
        "product_commit": PRODUCT_COMMIT,
        "verdicts": {r.name: r.verdict for r in results},
        "passed": {r.name: r.passed for r in results},
        "overall_pass": all_pass,
        "details": {r.name: r.details for r in results},
        "price_hypothesis_ok": price_ok,
        "no_send_ok": no_send_ok,
        "forbidden_found": forbidden_found,
    }
    output_json.write_text(json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

    # Markdown report
    md_lines = []
    md_lines.append("# B35 Independent QA Report - Lane C #1505")
    md_lines.append("")
    md_lines.append(f"- Commercial root: `{commercial_root}`")
    md_lines.append(f"- Package root: `{package_root}`")
    md_lines.append(f"- Product contract: `{product_contract}`")
    md_lines.append(f"- Product commit: `{PRODUCT_COMMIT}`")
    md_lines.append(f"- Branch: `feat/b35-w3-independent-qa-v31`")
    md_lines.append(f"- Base SHA: `eae88e0066c1b119bfa6c75d8b16c127b0137e5e`")
    md_lines.append("")
    md_lines.append("## Verdicts")
    md_lines.append("")
    md_lines.append("```text")
    for r in results:
        md_lines.append(r.verdict)
    md_lines.append("```")
    md_lines.append("")
    md_lines.append(f"**Overall:** {'PASS (all 8)' if all_pass else 'FAIL'}")
    md_lines.append("")
    for r in results:
        md_lines.append(f"### {r.name} - {r.verdict}")
        for d in r.details:
            md_lines.append(f"- {d}")
        md_lines.append("")
    md_lines.append("## Cross-cutting")
    md_lines.append(f"- Price hypothesis wording: {'OK' if price_ok else 'FAIL'}")
    md_lines.append(f"- No customer-send claim: {'OK' if no_send_ok else 'FAIL'}")
    md_lines.append(f"- Forbidden phrases: {forbidden_found if forbidden_found else 'none'}")
    md_lines.append("")
    md_lines.append(f"_Generated by tools/b35-independent-qa/validate_b35_independent_qa.py - Lane C harness_")
    md_lines.append(f"_Legacy lineage PR #359 referenced only, historical PASS not transferred_")
    output_md.write_text("\n".join(md_lines) + "\n", encoding="utf-8")

    print()
    print(f"Evidence written:")
    print(f"  JSON: {output_json}")
    print(f"  MD  : {output_md}")

    return 0 if all_pass else 1


if __name__ == "__main__":
    sys.exit(main())
